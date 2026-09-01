# -*- coding: utf-8 -*-
"""
MY 설교 AI 스튜디오 Pro  (v3.0 - 원고 기반 정밀 분석 엔진)

[v3.0 핵심 개선]
 1) 원고 기반 근거(Grounding) 엔진 도입
    - 어떤 본문/제목을 넣어도 '그 원고에서 실제로 뽑아낸' 요약/카드뉴스가 나옵니다.
    - AI 프롬프트에 '원고에 없는 내용 창작 금지 + 근거 문장 인용 의무' 규칙을 강제합니다.
 2) AI 연결 실패를 숨기지 않음
    - 예전에는 API가 죽으면 몰래 '템플릿 문구'가 나와서 엉망인 결과처럼 보였습니다.
    - 이제는 상단에 경고를 띄우고, 대신 원고에서 문장을 직접 추출한 실제 요약을 보여줍니다.
 3) 최신 Gemini 모델 목록 + 자동 탐색
    - 구형 모델명(gemini-1.5-*)만 호출하다 전부 실패하던 문제 해결.
 4) 설교 전환 시 세션/캐시 완전 격리
 5) PPTX 배경 이미지가 검은 사각형에 완전히 가려지던 버그 수정(투명도 적용)
 6) 다운로드 버튼이 매 rerun마다 문서를 새로 만들던 성능 폭탄 제거(캐시)
"""

import PIL.Image
import PIL.ImageDraw
import PIL.ImageFont
import PIL.ImageFilter

if not hasattr(PIL.Image, 'ANTIALIAS'):
    PIL.Image.ANTIALIAS = getattr(PIL.Image, 'Resampling', PIL.Image).LANCZOS

import streamlit as st
import google.generativeai as genai
import json
import os
import io
import re
import hashlib
import asyncio
import zipfile
import urllib.parse
import urllib.request
from collections import Counter
from datetime import datetime, timedelta

from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGB
from docx.oxml.ns import qn as docx_qn
from docx.oxml import OxmlElement
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import ParagraphStyle
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# 선택적 의존성 (없어도 앱 전체가 죽지 않도록 보호)
try:
    import edge_tts
    HAS_TTS = True
except Exception:
    HAS_TTS = False

try:
    import yt_dlp
    HAS_YTDLP = True
except Exception:
    HAS_YTDLP = False

# moviepy 1.x 는 최신 파이썬에서 설치가 실패하므로 2.x 를 우선 사용하고,
# 1.x 스타일 메서드 이름(set_position / subclip / resize / crop ...)은 별칭으로 되살린다.
HAS_MOVIEPY = False
try:
    from moviepy import VideoFileClip, ColorClip, CompositeVideoClip   # moviepy 2.x
    HAS_MOVIEPY = True
except Exception:
    try:
        from moviepy.editor import VideoFileClip, ColorClip, CompositeVideoClip   # moviepy 1.x
        HAS_MOVIEPY = True
    except Exception:
        HAS_MOVIEPY = False

if HAS_MOVIEPY:
    for _cls in (VideoFileClip, ColorClip, CompositeVideoClip):
        for _old, _new in (("set_position", "with_position"), ("set_duration", "with_duration"),
                           ("set_start", "with_start"), ("set_opacity", "with_opacity"),
                           ("set_audio", "with_audio"), ("subclip", "subclipped"),
                           ("resize", "resized"), ("crop", "cropped")):
            if not hasattr(_cls, _old) and hasattr(_cls, _new):
                try:
                    setattr(_cls, _old, getattr(_cls, _new))
                except Exception:
                    pass

try:
    from video_engine import create_animated_video, create_pil_text_clip
    HAS_VIDEO_ENGINE = True
except Exception:
    HAS_VIDEO_ENGINE = False

    def create_animated_video(*args, **kwargs):
        raise RuntimeError("video_engine.py 모듈을 찾을 수 없습니다.")

    def create_pil_text_clip(*args, **kwargs):
        raise RuntimeError("video_engine.py 모듈을 찾을 수 없습니다.")


st.set_page_config(
    page_title="MY 설교 AI 스튜디오 Pro",
    page_icon="🕊️",
    layout="wide",
    initial_sidebar_state="expanded"
)

st.markdown("""
<style>
    /* ========== 전역 톤 : 딥 인디고 + 오로라 그라데이션 ========== */
    :root{
        --ink:#070b1c; --ink2:#0e1533; --line:rgba(148,163,255,.18);
        --gold:#ffd766; --sky:#5ec8ff; --violet:#a78bfa; --mint:#5eead4; --rose:#fb7ff0;
        --txt:#eef2ff; --muted:#9aa6d4;
    }
    [data-testid="stAppViewContainer"]{
        background:
          radial-gradient(1100px 620px at 8% -8%, rgba(124,58,237,.30), transparent 60%),
          radial-gradient(900px 560px at 96% 4%, rgba(14,165,233,.24), transparent 60%),
          radial-gradient(800px 700px at 55% 108%, rgba(236,72,153,.16), transparent 60%),
          linear-gradient(180deg,#070b1c 0%,#0b1026 55%,#080c1e 100%);
        background-attachment: fixed;
    }
    [data-testid="stHeader"]{ background: transparent; }
    [data-testid="stSidebar"]{
        background: linear-gradient(180deg, rgba(13,18,45,.96), rgba(8,11,28,.96));
        border-right:1px solid var(--line);
    }
    html, body, [class*="css"]{ color: var(--txt); }

    /* ========== 버튼 ========== */
    .stButton > button, [data-testid="stDownloadButton"] > button{
        border-radius:12px !important;
        border:1px solid var(--line) !important;
        background: linear-gradient(135deg, rgba(99,102,241,.20), rgba(14,165,233,.16)) !important;
        color:#e9edff !important; font-weight:700 !important;
        transition: all .18s ease !important;
    }
    .stButton > button:hover, [data-testid="stDownloadButton"] > button:hover{
        border-color: rgba(167,139,250,.75) !important;
        box-shadow: 0 8px 22px rgba(99,102,241,.35) !important;
        transform: translateY(-1px);
    }
    .stButton > button[kind="primary"]{
        background: linear-gradient(135deg,#7c3aed 0%,#4f46e5 45%,#0ea5e9 100%) !important;
        border:none !important; color:#fff !important;
        box-shadow:0 10px 26px rgba(79,70,229,.45) !important;
    }
    div[data-testid="column"] button{
        width:100% !important; padding:6px 8px !important; font-size:12px !important;
    }

    /* ========== 입력/선택 ========== */
    [data-baseweb="input"] input, [data-baseweb="textarea"] textarea, [data-baseweb="select"] > div{
        background: rgba(11,16,38,.85) !important;
        border-radius:10px !important; border:1px solid var(--line) !important; color:var(--txt) !important;
    }
    [data-testid="stExpander"]{
        border:1px solid var(--line) !important; border-radius:16px !important;
        background: linear-gradient(135deg, rgba(20,26,58,.72), rgba(12,17,40,.72)) !important;
        backdrop-filter: blur(8px);
        margin-bottom:10px;
    }
    [data-testid="stExpander"] summary{ font-weight:700 !important; color:#dbe3ff !important; }

    /* 사역 메뉴 라디오를 알약 버튼처럼 */
    section[data-testid="stSidebar"] [role="radiogroup"] label,
    div[data-testid="stVerticalBlock"] [role="radiogroup"] label{
        border-radius:12px; padding:7px 10px; margin-bottom:5px;
        border:1px solid transparent; transition:all .16s ease;
    }
    div[data-testid="stVerticalBlock"] [role="radiogroup"] label:hover{
        background: rgba(124,58,237,.16); border-color: rgba(167,139,250,.45);
    }

    /* ========== 콘텐츠 카드 ========== */
    .content-box{
        background: linear-gradient(160deg, rgba(21,27,60,.88), rgba(11,15,36,.92));
        border:1px solid var(--line);
        border-radius:20px; padding:26px 28px;
        line-height:1.9; color:var(--txt); font-size:15.2px; margin-top:12px;
        box-shadow:0 18px 44px rgba(2,6,23,.55);
        white-space:normal; word-break:keep-all;
    }
    .p-line{ display:block; margin:3px 0; }
    .content-box h3{ color:var(--gold); font-size:19px; margin-top:18px; font-weight:800; }

    /* 섹션 제목 줄 (📌 / 💡 / 🙏 등으로 시작하는 줄) */
    .sec-head{
        display:block; margin:20px 0 8px 0; padding:9px 14px;
        font-weight:800; font-size:16.5px; color:#fff;
        background: linear-gradient(90deg, rgba(124,58,237,.42), rgba(14,165,233,.14) 70%, transparent);
        border-left:4px solid var(--violet); border-radius:0 12px 12px 0;
    }
    /* 인도자 가이드 — 제목과 내용 전체를 파란 블록으로 */
    .leader-block{
        display:block; margin:8px 0; padding:12px 16px;
        background: linear-gradient(135deg, rgba(56,189,248,.16), rgba(59,130,246,.09));
        border:1px solid rgba(56,189,248,.42);
        border-left:5px solid #38bdf8;
        border-radius:0 14px 14px 0;
        color:#bfe9ff !important; font-weight:600; line-height:1.75;
    }
    .leader-block b{ color:#7dd3fc; }
    .ground-quote{
        display:block; margin:6px 0 6px 10px; padding:8px 14px;
        color:#c7d2fe !important; background:rgba(99,102,241,.14);
        border-left:3px solid #818cf8; border-radius:0 10px 10px 0; font-size:13.8px;
    }
    .num-item{ display:block; padding:2px 0 2px 4px; }
    .num-badge{
        display:inline-block; min-width:22px; height:22px; line-height:22px; text-align:center;
        border-radius:7px; font-size:12px; font-weight:800; margin-right:8px;
        background:linear-gradient(135deg,#7c3aed,#0ea5e9); color:#fff;
    }

    /* ========== 헤더 / 배지 ========== */
    .hero{
        background: linear-gradient(135deg, rgba(124,58,237,.34), rgba(14,165,233,.20));
        border:1px solid var(--line); border-radius:22px;
        padding:20px 26px; margin-bottom:14px;
        box-shadow:0 16px 40px rgba(2,6,23,.5);
    }
    .hero h1{ margin:0; font-size:29px; font-weight:900; color:#fff; letter-spacing:-.5px; }
    .chip{
        display:inline-block; padding:5px 13px; border-radius:999px;
        font-size:12.5px; font-weight:800; margin-right:6px;
    }
    .chip-gold{ background:linear-gradient(135deg,#f59e0b,#fbbf24); color:#1a1200; }
    .chip-sky { background:linear-gradient(135deg,#0ea5e9,#38bdf8); color:#04212e; }
    .chip-vio { background:linear-gradient(135deg,#7c3aed,#a78bfa); color:#fff; }
    .chip-mint{ background:linear-gradient(135deg,#0d9488,#5eead4); color:#012b26; }
    .badge-ok { background:linear-gradient(135deg,#065f46,#10b981); color:#eafff6; padding:4px 11px; border-radius:9px; font-size:12px; font-weight:700;}
    .badge-bad{ background:linear-gradient(135deg,#7f1d1d,#ef4444); color:#fff1f1; padding:4px 11px; border-radius:9px; font-size:12px; font-weight:700;}
    .lib-card{
        background: linear-gradient(150deg, rgba(22,28,62,.9), rgba(11,15,36,.9));
        border:1px solid var(--line); border-radius:18px; padding:20px 22px; margin-bottom:12px;
        box-shadow:0 12px 30px rgba(2,6,23,.42);
    }
    hr{ border-color: var(--line) !important; }

    /* ===== 안내·가이드 문구는 전부 파란 계열로 통일 ===== */
    [data-testid="stCaptionContainer"], [data-testid="stCaptionContainer"] p{
        color:#7dd3fc !important; font-weight:600;
    }
    [data-testid="stCaptionContainer"]{
        border-left:3px solid rgba(56,189,248,.55);
        background:rgba(56,189,248,.07);
        padding:6px 12px; border-radius:0 10px 10px 0; margin:4px 0 8px 0;
    }
    /* st.info = 안내 박스 */
    div[data-testid="stAlertContainer"]{
        border-radius:14px !important; border:1px solid rgba(56,189,248,.35) !important;
    }
    /* 라디오/셀렉트 라벨 */
    label[data-testid="stWidgetLabel"] p{ color:#c3cdf5 !important; font-weight:700; }

    /* 탭 */
    button[data-baseweb="tab"]{ font-weight:700 !important; }
    div[data-baseweb="tab-highlight"]{
        background: linear-gradient(90deg,#7c3aed,#0ea5e9) !important; height:3px !important;
    }
    /* 슬라이더 */
    div[data-testid="stSlider"] [data-baseweb="slider"] div[role="slider"]{
        background: linear-gradient(135deg,#7c3aed,#0ea5e9) !important;
    }
    /* 사이드바 메뉴 */
    section[data-testid="stSidebar"] [role="radiogroup"] label:hover{
        background: rgba(14,165,233,.14);
    }
    /* 카드 미리보기 이미지 */
    [data-testid="stImage"] img{ border-radius:16px; box-shadow:0 14px 34px rgba(2,6,23,.55); }
</style>
""", unsafe_allow_html=True)


# ==============================================================================
# 성경 66권 · 약어 · 분류
# ==============================================================================
OLD_TESTAMENT_BOOKS = [
    "창세기", "출애굽기", "레위기", "민수기", "신명기", "여호수아", "사사기", "룻기",
    "사무엘상", "사무엘하", "열왕기상", "열왕기하", "역대상", "역대하", "에스라", "느헤미야",
    "에스더", "욥기", "시편", "잠언", "전도서", "아가", "이사야", "예레미야", "예레미야애가",
    "에스겔", "다니엘", "호세아", "요엘", "아모스", "오바댜", "요나", "미가", "나훔",
    "하박국", "스바냐", "학개", "스가랴", "말라기"
]
NEW_TESTAMENT_BOOKS = [
    "마태복음", "마가복음", "누가복음", "요한복음", "사도행전", "로마서", "고린도전서", "고린도후서",
    "갈라디아서", "에베소서", "빌립보서", "골로새서", "데살로니가전서", "데살로니가후서", "디모데전서",
    "디모데후서", "디도서", "빌레몬서", "히브리서", "야고보서", "베드로전서", "베드로후서",
    "요한일서", "요한이서", "요한삼서", "유다서", "요한계시록"
]
BIBLE_BOOKS = OLD_TESTAMENT_BOOKS + NEW_TESTAMENT_BOOKS

# 약어 → 정식명 (설교문에 '롬 8:28' 처럼 쓰는 경우 대응)
BIBLE_ABBREV = {
    "창": "창세기", "출": "출애굽기", "레": "레위기", "민": "민수기", "신": "신명기",
    "수": "여호수아", "삿": "사사기", "룻": "룻기", "삼상": "사무엘상", "삼하": "사무엘하",
    "왕상": "열왕기상", "왕하": "열왕기하", "대상": "역대상", "대하": "역대하",
    "스": "에스라", "느": "느헤미야", "에": "에스더", "욥": "욥기", "시": "시편",
    "잠": "잠언", "전": "전도서", "아": "아가", "사": "이사야", "렘": "예레미야",
    "애": "예레미야애가", "겔": "에스겔", "단": "다니엘", "호": "호세아", "욜": "요엘",
    "암": "아모스", "옵": "오바댜", "욘": "요나", "미": "미가", "나": "나훔",
    "합": "하박국", "습": "스바냐", "학": "학개", "슥": "스가랴", "말": "말라기",
    "마": "마태복음", "막": "마가복음", "눅": "누가복음", "요": "요한복음", "행": "사도행전",
    "롬": "로마서", "고전": "고린도전서", "고후": "고린도후서", "갈": "갈라디아서",
    "엡": "에베소서", "빌": "빌립보서", "골": "골로새서", "살전": "데살로니가전서",
    "살후": "데살로니가후서", "딤전": "디모데전서", "딤후": "디모데후서", "딛": "디도서",
    "몬": "빌레몬서", "히": "히브리서", "약": "야고보서", "벧전": "베드로전서",
    "벧후": "베드로후서", "요일": "요한일서", "요이": "요한이서", "요삼": "요한삼서",
    "유": "유다서", "계": "요한계시록",
}


# ==============================================================================
# 설교 작성 : 신학적 관점(렌즈) & 대지 구조
# ==============================================================================
THEOLOGY_LENSES = {
    "개혁주의 (칼빈주의 · 하나님의 주권)": {
        "desc": "하나님의 절대 주권과 전적 은혜를 축으로, 인간의 무능력과 은혜의 우선성을 드러냅니다.",
        "guide": (
            "- 본문에서 '하나님이 주어가 되시는 동사'를 먼저 찾아 강조하십시오.\n"
            "- 인간의 공로나 결단이 아니라 하나님의 작정·섭리·보존이 이야기의 주도권을 쥐고 있음을 밝히십시오.\n"
            "- 웨스트민스터 신앙고백의 어법(작정, 섭리, 유효한 부르심, 견인)을 자연스럽게 녹이십시오.\n"
            "- 적용은 '더 애쓰라'가 아니라 '이미 주신 은혜 위에 서라'는 방향으로 이끄십시오."
        ),
    },
    "장로교 정통 (웨스트민스터 · 언약)": {
        "desc": "웨스트민스터 표준문서의 틀 안에서 교회와 성례, 언약 공동체의 자리를 짚습니다.",
        "guide": (
            "- 본문을 개인 경건에 가두지 말고 언약 공동체(교회)의 자리에서 읽으십시오.\n"
            "- 소요리문답·신앙고백의 관련 조항을 한 번 이상 인용하되 어렵지 않게 풀어 쓰십시오.\n"
            "- 말씀·성례·권징이라는 은혜의 방편과 본문을 연결하십시오.\n"
            "- 세대를 잇는 신앙 전수(언약의 자녀)를 적용 한 가지에 포함하십시오."
        ),
    },
    "복음주의 (십자가 · 은혜 · 복음선포)": {
        "desc": "십자가와 부활, 개인의 회심과 복음 전파를 전면에 세웁니다.",
        "guide": (
            "- 모든 대지가 마지막에 십자가와 부활로 수렴하게 하십시오.\n"
            "- 회심하지 않은 청중을 향한 초청(결단) 문단을 결론에 반드시 두십시오.\n"
            "- 쉬운 언어로, 신학 용어는 즉시 풀어 설명하십시오.\n"
            "- 개인의 삶의 변화와 복음 증거를 적용에 담으십시오."
        ),
    },
    "성경신학 (구속사 · 정경 전체의 흐름)": {
        "desc": "창조–타락–구속–완성의 큰 이야기 속에서 본문의 위치를 밝힙니다.",
        "guide": (
            "- 먼저 이 본문이 구속사(창조–타락–구속–완성)의 어느 지점에 서 있는지 명시하십시오.\n"
            "- 이 본문 이전의 약속과 이후의 성취를 정경 전체 흐름으로 연결하십시오 "
            "(구약이면 그리스도에게로, 신약이면 구약 배경으로).\n"
            "- 인물을 도덕적 모범으로 삼는 '교훈 설교'를 금지합니다. 본문의 주인공은 하나님이십니다.\n"
            "- 모형(type)과 성취(antitype)를 억지로 만들지 말고, 성경이 실제로 연결한 것만 쓰십시오.\n"
            "- 적용은 '이 이야기 안에 있는 우리는 누구인가'에서 도출하십시오."
        ),
    },
    "언약적 관점 (언약의 주와 백성)": {
        "desc": "하나님이 맺으신 언약의 구조(약속·조건·표징·성취)로 본문을 읽습니다.",
        "guide": (
            "- 본문에 작동하는 언약이 무엇인지(아브라함·모세·다윗·새 언약) 먼저 규정하십시오.\n"
            "- 언약의 4요소 — 언약의 주(하나님), 언약 백성, 언약의 약속과 요구, 언약의 표징 — 을 짚으십시오.\n"
            "- 인간의 언약 파기와 하나님의 언약 신실하심(헤세드)을 대비시키십시오.\n"
            "- 그리스도께서 언약의 중보자로 어떻게 성취하셨는지 밝히십시오.\n"
            "- 적용은 '언약 백성답게 사는 삶'으로 귀결시키십시오."
        ),
    },
    "팀 켈러 관점 (복음중심 · 우상 분석 · 도시 변증)": {
        "desc": "마음의 우상을 드러내고, 종교와 비종교를 모두 넘어서는 복음의 제3의 길을 제시합니다.",
        "guide": (
            "- 서론은 오늘 도시인의 실제 질문·불안·욕망에서 출발하십시오 (문화적 변증).\n"
            "- 본문이 겨냥하는 '마음의 우상'(성공, 인정, 통제, 안전, 관계 등)을 이름 붙여 드러내십시오.\n"
            "- '종교적 도덕주의'와 '세속적 방종' 두 길을 먼저 제시하고, 복음이 그 둘을 모두 해체하는 "
            "제3의 길임을 보이십시오.\n"
            "- 문제 → 복음 → 그러므로 그리스도 안에서, 이 흐름을 유지하십시오.\n"
            "- 반드시 그리스도께서 '우리 대신' 하신 일로 착지하십시오. 도덕적 권면으로 끝내지 마십시오.\n"
            "- 회의하는 사람도 끝까지 들을 수 있게, 반론을 먼저 공정하게 말한 뒤 답하십시오."
        ),
    },
    "마틴 로이드 존스 관점 (강해설교 · 교리적 논증)": {
        "desc": "본문을 논리적으로 논증하며 죄와 은혜를 깊이 파고들고, 성령의 능력 있는 선포를 지향합니다.",
        "guide": (
            "- 본문을 한 절 한 절 논리적으로 논증하십시오. '무엇을 말하는가 → 왜 그러한가 → 그러므로 무엇인가' 순서로.\n"
            "- 인간의 죄와 무능력을 회피하지 말고 정면으로 다룬 뒤, 그 위에 은혜를 세우십시오.\n"
            "- 교리(칭의·중생·성화)를 설교 안에서 분명히 가르치되, 청중의 영혼을 향해 말하십시오.\n"
            "- 감정에 호소하기 전에 진리로 설득하고, 마지막에 성령의 역사하심을 간구하는 선포로 맺으십시오.\n"
            "- '논리에 불이 붙은 것'이 설교라는 원칙대로, 뜨겁되 논리가 흐트러지지 않게 하십시오."
        ),
    },
    "존 스토트 관점 (강해 · 이중 경청 · 균형)": {
        "desc": "본문 세계와 오늘의 세계를 잇는 '다리 놓기' 설교로, 균형과 통합을 중시합니다.",
        "guide": (
            "- '이중 경청' — 성경 본문의 세계와 오늘의 세상을 함께 들으십시오. 두 세계를 잇는 다리를 세우는 것이 목표입니다.\n"
            "- 본문이 말하는 것 이상도 이하도 말하지 마십시오. 본문의 지배를 받는 설교여야 합니다.\n"
            "- 복음 전도와 사회적 책임을 대립시키지 말고 함께 다루십시오.\n"
            "- 회의하는 사람도 납득할 수 있도록, 반론을 공정하게 소개한 뒤 성경으로 답하십시오.\n"
            "- 절제된 문장, 명확한 구조, 지적 정직성을 유지하십시오."
        ),
    },
    "매튜 헨리 관점 (본문 축자 주석 · 실천적 적용)": {
        "desc": "본문을 구절 단위로 세밀히 주석하고, 각 구절마다 경건한 적용을 붙입니다.",
        "guide": (
            "- 본문을 구절 단위로 나누어 차례대로 주석하십시오. 건너뛰지 마십시오.\n"
            "- 각 구절마다 '관찰 → 의미 → 우리에게 주는 교훈' 세 걸음을 붙이십시오.\n"
            "- 성경으로 성경을 해석하십시오. 관련 구절을 자주 연결하십시오.\n"
            "- 문장은 따뜻하고 경건하게, 독자의 마음을 살피듯 쓰십시오.\n"
            "- 사변적 논쟁보다 실천적 경건과 위로를 앞세우십시오."
        ),
    },
    "해돈 로빈슨 관점 (빅 아이디어 강해설교)": {
        "desc": "본문의 '큰 개념(Big Idea)' 하나를 찾아, 설교 전체를 그 하나로 통일시킵니다.",
        "guide": (
            "- 먼저 본문의 주해적 개념을 '주제(무엇을 말하는가) + 보충어(무엇이라 말하는가)' 한 문장으로 규정하십시오.\n"
            "- 그 한 문장(빅 아이디어)을 설교 전체가 섬기게 하십시오. 대지는 그 개념을 설명·증명·적용하는 도구입니다.\n"
            "- 설교 목적을 분명히 하십시오: 이 설교를 들은 청중이 무엇을 알고, 느끼고, 행하기를 원하는가.\n"
            "- 빅 아이디어를 설교 안에서 최소 3회 같은 문장으로 반복하십시오.\n"
            "- 곁가지 정보는 아무리 흥미로워도 빅 아이디어를 섬기지 않으면 잘라내십시오."
        ),
    },
    "존 맥아더 관점 (절별 강해 · 원어 주해)": {
        "desc": "본문을 절별로 철저히 주해하며 원어와 문법을 근거로 정확한 의미를 밝힙니다.",
        "guide": (
            "- 본문을 절별로 순서대로 강해하십시오. 각 절의 핵심 단어를 원어로 짚고 문법을 설명하십시오.\n"
            "- 문맥(앞뒤 절, 책 전체)을 근거로 의미를 확정하고, 잘못된 통속적 해석을 분명히 교정하십시오.\n"
            "- 성경의 무오성과 충족성을 전제로, 본문이 말하는 바를 타협 없이 선포하십시오.\n"
            "- 예화는 최소화하고 본문 자체의 설명에 지면을 쓰십시오.\n"
            "- 적용은 순종과 거룩을 향한 분명한 요구로 제시하십시오."
        ),
    },
    "하용조 목사 관점 (성령 · 사도행전적 교회)": {
        "desc": "성령의 인도하심과 치유·회복, 사도행전적 교회를 향한 열정을 담습니다.",
        "guide": (
            "- 본문에서 성령께서 지금 무엇을 하고 계신지를 드러내십시오.\n"
            "- 상처 입은 영혼을 향한 치유와 회복의 메시지를 담되, 값싼 위로가 아니라 복음의 능력으로 이끄십시오.\n"
            "- 사도행전적 교회 — 예배·전도·선교·공동체가 살아 있는 교회상을 그리십시오.\n"
            "- 문장은 따뜻하고 간결하게, 청중과 대화하듯 쓰십시오.\n"
            "- 개인의 결단을 넘어 공동체와 열방을 향한 헌신으로 확장하십시오."
        ),
    },
    "옥한흠 목사 관점 (제자훈련 · 평신도를 깨운다)": {
        "desc": "모든 성도를 그리스도의 제자로 세우는 제자훈련의 시각으로 본문을 읽습니다.",
        "guide": (
            "- 본문을 '이 말씀이 한 사람을 어떻게 제자로 세우는가'라는 질문으로 읽으십시오.\n"
            "- 목회자와 평신도의 이분법을 깨십시오. 모든 성도가 사역자임을 본문에서 끌어내십시오.\n"
            "- 값싼 은혜와 명목상의 신앙을 정직하게 지적하고, 대가를 치르는 제자도로 부르십시오.\n"
            "- 소그룹·삶공부·순종의 훈련이라는 구체적 통로를 적용에 제시하십시오.\n"
            "- 한 영혼을 향한 목회자의 진정성이 문장에 배어나게 쓰십시오."
        ),
    },
    "선교적 교회 관점 (보내심 받은 공동체)": {
        "desc": "교회를 하나님의 선교에 보내심 받은 공동체로 보고, 본문을 그 자리에서 읽습니다.",
        "guide": (
            "- 선교의 주체는 교회가 아니라 하나님(Missio Dei)이십니다. 본문에서 보내시는 하나님을 먼저 드러내십시오.\n"
            "- 교회는 '선교하는 조직'이 아니라 '보내심 받은 백성'입니다. 이 정체성으로 본문을 읽으십시오.\n"
            "- 성도의 일터·이웃·지역이 곧 선교지임을 구체적으로 짚으십시오.\n"
            "- 개인 구원에 머물지 말고 하나님 나라의 회복과 이웃 사랑으로 확장하십시오.\n"
            "- 적용은 '흩어지는 교회'로서 이번 주에 무엇을 할 것인가로 착지시키십시오."
        ),
    },
    "교의학적 관점 (조직신학 주제 중심)": {
        "desc": "본문에서 신론·기독론·구원론 등 교리를 끌어내어 체계적으로 가르칩니다.",
        "guide": (
            "- 본문이 가장 강하게 증언하는 교리 한 가지(신론/기독론/성령론/구원론/교회론/종말론)를 "
            "설교의 중심 교리로 명시하십시오.\n"
            "- 그 교리를 정의 → 성경적 근거 → 흔한 오해 → 바른 이해 순으로 전개하십시오.\n"
            "- 관련 신조·신앙고백을 한 번 인용하되 쉽게 풀어 쓰십시오.\n"
            "- 교리가 차가운 지식이 아니라 예배와 삶으로 이어지도록, 각 대지 끝에 "
            "'그러므로 우리는 이렇게 예배한다'를 붙이십시오.\n"
            "- 교리 용어를 쓸 때마다 즉시 일상 언어로 다시 설명하십시오."
        ),
    },
}

SERMON_LENGTHS = {
    "5분 (짧은 권면)":     {"chars": 1200,  "tokens": 3000,
                        "guide": "새벽기도·짧은 권면 분량입니다. 서론은 3문장 이내, 대지 설명은 핵심만, 예화는 1개만 쓰십시오."},
    "10분 (경건회)":       {"chars": 2400,  "tokens": 4500,
                        "guide": "경건회·주중예배 분량입니다. 예화는 2개 이내, 각 대지는 짧고 선명하게 쓰십시오."},
    "15분 (주중예배)":     {"chars": 3600,  "tokens": 6000,
                        "guide": "본문 주해와 적용을 균형 있게 담되 곁가지를 치지 마십시오."},
    "20분 (간결한 주일설교)": {"chars": 4800,  "tokens": 7000,
                        "guide": "각 대지에 주해 1문단 · 예화 1개 · 적용 1개씩 배치하십시오."},
    "25분 (표준 주일설교)":  {"chars": 6000,  "tokens": 8192,
                        "guide": "표준 주일 강단 분량입니다. 서론·배경·대지·결론·기도를 고루 갖추십시오."},
    "30분 (충실한 주일설교)": {"chars": 7200,  "tokens": 8192,
                        "guide": "본문 배경 설명과 예화를 넉넉히 넣되, 늘어지지 않게 각 대지를 균등하게 배분하십시오."},
    "40분 (강해 심화)":     {"chars": 9600,  "tokens": 12000,
                        "guide": "절별 주해를 상세히 넣고, 원어·배경 설명과 예화를 각 대지마다 충분히 배치하십시오."},
    "60분 (성경공부·특강)":  {"chars": 14000, "tokens": 16000,
                        "guide": "성경공부·특강 분량입니다. 절별 주해, 배경, 신학적 쟁점, 질의응답용 보충 설명까지 담으십시오."},
}

OUTLINE_SHAPES = {
    "원포인트 (One-Point)": {
        "points": ["본론 — 하나의 중심 진리를 세 국면으로 심화"],
        "chars": 5000, "minutes": "22~28분",
        "guide": (
            "- 대지를 나누지 말고, 하나의 중심 진리를 끝까지 밀고 가십시오.\n"
            "- 그 하나의 진리를 (1) 본문에서 발견 (2) 깊이 파고들기 (3) 삶에 적용, 세 국면으로 심화시키십시오.\n"
            "- 중심 문장(한 문장)을 설교 안에서 최소 4번 반복해 청중의 귀에 박히게 하십시오.\n"
            "- 곁가지 주제를 넣지 마십시오. 하나만 남기는 것이 이 형식의 목적입니다."
        ),
    },
    "2대지": {
        "points": ["제1대지", "제2대지"],
        "chars": 5500, "minutes": "24~28분",
        "guide": (
            "- 두 대지는 대비 구조(문제↔해답, 옛것↔새것, 인간↔하나님)로 세우십시오.\n"
            "- 첫 대지에서 긴장을 충분히 쌓고, 둘째 대지에서 해소하십시오.\n"
            "- 각 대지에 예화 1개와 구체적 적용 1개를 반드시 넣으십시오."
        ),
    },
    "3대지": {
        "points": ["제1대지", "제2대지", "제3대지"],
        "chars": 6000, "minutes": "25~30분",
        "guide": (
            "- 세 대지가 본문의 흐름을 따라 순차적으로 전개되게 하십시오.\n"
            "- 세 대지는 병렬이 아니라 점층(약속→위기→성취 등)이 되도록 배열하십시오.\n"
            "- 각 대지에 본문 인용 · 주해 · 예화 · 적용을 모두 넣으십시오."
        ),
    },
    "4대지": {
        "points": ["제1대지", "제2대지", "제3대지", "제4대지"],
        "chars": 7000, "minutes": "30~35분",
        "guide": (
            "- 네 대지는 본문의 절 구분을 따라 나누되, 각 대지를 짧고 선명하게 유지하십시오.\n"
            "- 대지마다 한 문장 요약(소제목)을 먼저 던지고 풀어 가십시오.\n"
            "- 길어지는 만큼 예화는 짧게, 적용은 구체적으로 쓰십시오."
        ),
    },
}


def classify_scripture(scripture_text: str):
    """본문 문자열에서 구약/신약 및 책 이름을 판별한다. (긴 이름 우선 매칭)"""
    if not scripture_text:
        return "기타", "성경전체"
    txt = scripture_text.strip()

    for b in sorted(BIBLE_BOOKS, key=len, reverse=True):
        if b in txt:
            return ("구약" if b in OLD_TESTAMENT_BOOKS else "신약"), b

    for ab in sorted(BIBLE_ABBREV.keys(), key=len, reverse=True):
        if re.search(rf'(^|[^가-힣]){re.escape(ab)}\s*\d', txt):
            full = BIBLE_ABBREV[ab]
            return ("구약" if full in OLD_TESTAMENT_BOOKS else "신약"), full

    return "기타", "성경전체"


# ==============================================================================
# ★★★ 원고 정밀 분석 엔진 (본 앱 품질의 심장) ★★★
#   - 지어내지 않는다. 원고에 실제로 있는 문장/단어만 뽑는다.
# ==============================================================================

KOR_STOPWORDS = set("""
그리고 그러나 그런데 하지만 그래서 그러므로 우리 여러분 오늘 이것 그것 저것 때문 통해 위해 대해
정말 참으로 지금 다시 모든 어떤 무엇 이런 저런 그런 이렇게 그렇게 저렇게 사람 사람들 이야기 경우
때문에 그때 지금은 여기 거기 저기 하나 둘 셋 자신 자기 서로 함께 아주 매우 너무 조금 많이 항상
말씀 하나님 예수님 예수 주님 아멘 성도 성도님 교회 오늘날 우리가 우리는 우리를 저는 제가 여러분이
그분 그것이 있습니다 없습니다 합니다 됩니다 입니다 것입니다 것이다 라고 라는 이라는 하는 되는 있는
없는 같은 위한 통한 한번 다시한번 사실 물론 결국 특히 바로 아마 혹시 만약 그럼 자 이제
안에 밖에 위에 아래 아니라 아니고 말고 자가 나도 너도 여기에 거기에 그때에 오늘은 본문 본문은
읽은 읽는 보면 보니 보십시오 하십시오 합시다 첫째 둘째 셋째 넷째 먼저 다음으로 끝으로 마지막으로
사랑하는 존경하는 축복 은혜 믿음 기도 신앙 삶의 인생 하나님의 우리의 저희 여러분들
""".split())

JOSA_SUFFIX = (
    "으로써", "에게서", "께서는", "에서는", "으로는", "이라는", "에게는", "까지도", "부터는", "라고도",
    "께서", "에게", "에서", "으로", "이라", "라는", "처럼", "까지", "부터", "보다", "조차", "마저",
    "한테", "이나", "이며", "이고", "라도", "든지",
    "은", "는", "이", "가", "을", "를", "에", "의", "도", "와", "과", "로", "만", "랑",
)

VERB_TAIL_PAT = re.compile(
    r'(니다|습니다|십시오|시오|세요|읍시다|했다|한다|하다|하는|하며|하고|해서|해도|'
    r'되어|되고|되는|겠다|것이|것을|것은|것도|이라|라고|면서|는데|지만|어서|아서)$'
)

# 원고 안에서 성경 인용을 찾아내는 정규식
_BOOK_ALT = "|".join(sorted(BIBLE_BOOKS, key=len, reverse=True) +
                     sorted(BIBLE_ABBREV.keys(), key=len, reverse=True))
SCRIPTURE_REF_RE = re.compile(
    rf'((?:{_BOOK_ALT}))\s*(\d{{1,3}})\s*(?:장|[:：])\s*(\d{{1,3}})?\s*(?:[-~–]\s*(\d{{1,3}}))?\s*절?'
)

# 적용/권면 신호 어미 (한국어 설교의 권면 어미를 폭넓게 포착)
APPLY_TAIL_RE = re.compile(
    r'(십시오|ㅂ시다|읍시다|합시다|시기\s*바랍니다|시길\s*바랍니다|기를\s*바랍니다|'
    r'해야\s*합니다|해야만\s*합니다|되시기\s*바랍니다|되기를\s*축복합니다|'
    r'하시기\s*축원합니다|되시기\s*축원합니다|합시다\.|드립시다|삽시다|나아갑시다)'
)

# 대지(포인트) 신호 — 한국 설교의 전형적 구조 표지
POINT_MARKER_RE = re.compile(
    r'^\s*(첫째|둘째|셋째|넷째|다섯째|먼저|다음으로|끝으로|마지막으로|'
    r'제\s*[1-5]\s*대지|[1-5]\s*대지|[1-5][\.\)])[,\s:：]'
)

# 서론·봉독 등 대지 제목으로 부적합한 상투 문장
INTRO_NOISE_RE = re.compile(
    r'(본문은|본문\s*말씀은|함께\s*읽은|봉독|읽어\s*드린|말씀을\s*읽었|오늘\s*우리가\s*함께|'
    r'설교\s*제목은|사랑하는\s*성도\s*여러분|기도하겠습니다)'
)

PRAYER_HINT_RE = re.compile(r'(기도(를)?\s*(드립니다|합니다|드리겠습니다)|예수님의\s*이름으로|아멘)')


_JOSA_SORTED = sorted(JOSA_SUFFIX, key=len, reverse=True)


def _strip_josa(word: str):
    """가장 긴 조사 1개만 제거한 후보를 돌려준다. (없으면 None)"""
    for j in _JOSA_SORTED:
        if len(word) > len(j) + 1 and word.endswith(j):
            return word[:-len(j)]
    return None


def build_stem_map(words):
    """
    조사를 무조건 떼면 '십자가'가 '십자'로 잘려 엉뚱한 키워드가 된다.
    → 어간으로 인정하는 조건:
       (a) 그 어간이 원고에 단독으로도 등장하거나
       (b) 서로 다른 조사 형태가 2개 이상 붙어 나타날 때
    """
    vocab = set(words)
    forms = {}
    for w in vocab:
        cand = _strip_josa(w)
        if cand:
            forms.setdefault(cand, set()).add(w)
    accepted = {s for s, fs in forms.items() if len(fs) >= 2 or s in vocab}
    mapping = {}
    for w in vocab:
        cand = _strip_josa(w)
        if cand and cand in accepted:
            mapping[w] = cand
    return mapping


def _stem(word: str, stem_map=None) -> str:
    if stem_map is not None:
        return stem_map.get(word, word)
    cand = _strip_josa(word)
    return cand if cand else word


def split_sentences(text: str):
    """한국어 문장 분리 (마침표/물음표/느낌표 + '~다.' 종결 고려)"""
    if not text:
        return []
    t = re.sub(r'\s+', ' ', text.replace('\n', ' '))
    raw = re.split(r'(?<=[.!?…])\s+|(?<=다\.)\s*|(?<=요\.)\s*|(?<=까\?)\s*', t)
    out = []
    for s in raw:
        s = s.strip()
        if len(s) >= 8:
            out.append(s)
    return out


def split_paragraphs(text: str):
    if not text:
        return []
    parts = [p.strip() for p in re.split(r'\n\s*\n|\n', text)]
    return [p for p in parts if len(p) >= 15]


def extract_scripture_refs(text: str):
    """원고에 실제로 인용된 성경 구절 목록 (등장 순, 중복 제거)"""
    refs, seen = [], set()
    for m in SCRIPTURE_REF_RE.finditer(text or ""):
        book = BIBLE_ABBREV.get(m.group(1), m.group(1))
        ch = m.group(2)
        v1 = m.group(3)
        v2 = m.group(4)
        if v1 and v2:
            label = f"{book} {ch}:{v1}-{v2}"
        elif v1:
            label = f"{book} {ch}:{v1}"
        else:
            label = f"{book} {ch}장"
        if label not in seen:
            seen.add(label)
            refs.append(label)
    return refs


NOISE_TOKEN_RE = re.compile(r'^(절부터|절까지|장에서|절에서|이라는|말씀은|본문의|장의|절의)$')


def extract_keywords(text: str, top_n: int = 20, stem_map=None):
    """원고에서 '이 설교만의' 특징 키워드 추출"""
    words = re.findall(r'[가-힣]{2,}', text or "")
    if stem_map is None:
        stem_map = build_stem_map(words)
    counter = Counter()
    for w in words:
        s = _stem(w, stem_map)
        if len(s) < 2:
            continue
        if s in KOR_STOPWORDS or w in KOR_STOPWORDS:
            continue
        if VERB_TAIL_PAT.search(s) or NOISE_TOKEN_RE.match(s):
            continue
        if re.match(r'^\d+$', s):
            continue
        counter[s] += 1
    # 1회만 나온 단어는 특징으로 보지 않음 (원고가 아주 짧으면 예외)
    min_cnt = 2 if len(words) > 400 else 1
    return [w for w, c in counter.most_common(top_n * 3) if c >= min_cnt][:top_n]


def _score_sentence(sent: str, kw_rank: dict, stem_map=None) -> float:
    score = 0.0
    for w in re.findall(r'[가-힣]{2,}', sent):
        s = _stem(w, stem_map)
        if s in kw_rank:
            score += kw_rank[s]
    if SCRIPTURE_REF_RE.search(sent):
        score += 3.0
    if '"' in sent or '“' in sent or "'" in sent:
        score += 1.0
    n = len(sent)
    if n < 20:
        score *= 0.5
    elif n > 160:
        score *= 0.8
    return score


@st.cache_data(show_spinner=False, max_entries=64)
def analyze_manuscript(text: str) -> dict:
    """
    설교 원고를 실제로 읽어 구조를 뽑아낸다.
    반환: keywords / refs / paragraphs / sentences / key_sentences / applications /
          prayer / sections(3분할 대표 문단)
    """
    text = (text or "").strip()
    paragraphs = split_paragraphs(text)
    sentences = split_sentences(text)
    stem_map = build_stem_map(re.findall(r'[가-힣]{2,}', text))
    keywords = extract_keywords(text, stem_map=stem_map)
    refs = extract_scripture_refs(text)

    kw_rank = {w: (len(keywords) - i) / len(keywords) * 2.0 for i, w in enumerate(keywords)} if keywords else {}

    def sc(s):
        return _score_sentence(s, kw_rank, stem_map)

    scored = sorted([(s, sc(s)) for s in sentences], key=lambda x: x[1], reverse=True)
    key_sentences = [s for s, _ in scored[:12]]

    # 적용/권면 문장 (원고에 실제로 있는 것만)
    applications = []
    for s in sentences:
        if APPLY_TAIL_RE.search(s) and 15 <= len(s) <= 160:
            applications.append(s)
    applications = sorted(applications, key=sc, reverse=True)[:8]

    # 기도문: 원고 뒤쪽 30% 안에서 기도 신호가 있는 문단
    prayer = ""
    tail_start = int(len(paragraphs) * 0.65)
    for p in reversed(paragraphs[tail_start:] or paragraphs[-3:] if paragraphs else []):
        if PRAYER_HINT_RE.search(p):
            prayer = p
            break

    # ── 대지(sections) 추출 ────────────────────────────────────────────────
    # 1순위: 설교자가 직접 표시한 '첫째/둘째/셋째, 먼저/다음으로/끝으로' 구조를 그대로 사용
    sections = []
    marked = [p for p in paragraphs if POINT_MARKER_RE.match(p)]
    if len(marked) >= 2:
        for p in marked[:4]:
            sents = split_sentences(p)
            head = sents[0] if sents else p
            sections.append({"headline": head[:90], "evidence": p[:420], "chunk": [p]})
    else:
        # 2순위: 원고를 3등분해 각 구간의 대표 문단을 뽑되, 서론·봉독 상투구는 제외
        if paragraphs:
            n = len(paragraphs)
            bounds = [(0, max(1, n // 3)), (max(1, n // 3), max(2, (2 * n) // 3)),
                      (max(2, (2 * n) // 3), n)]
            used = set()
            for a, b in bounds:
                chunk = paragraphs[a:b] or paragraphs[-1:]
                ranked = sorted(chunk, key=sc, reverse=True)
                best = next((p for p in ranked if p not in used and not INTRO_NOISE_RE.search(p[:60])),
                            ranked[0])
                used.add(best)
                cand = [s for s in split_sentences(best) if not INTRO_NOISE_RE.search(s)]
                head = max(cand, key=sc) if cand else best
                sections.append({"headline": head[:90], "evidence": best[:420], "chunk": chunk})

    # 명제로 쓸 문장: 서론·봉독 문장은 제외
    prop_pool = [s for s in key_sentences if not INTRO_NOISE_RE.search(s)]
    key_sentences = prop_pool + [s for s in key_sentences if s not in prop_pool]

    return {
        "keywords": keywords,
        "refs": refs,
        "paragraphs": paragraphs,
        "sentences": sentences,
        "key_sentences": key_sentences,
        "applications": applications,
        "prayer": prayer,
        "sections": sections,
        "char_count": len(text),
    }


def _trim_title(sent: str, limit: int = 42) -> str:
    """문장을 대지 제목처럼 다듬는다 (원문 단어는 유지)"""
    s = re.sub(r'^[0-9]+[\.\)]\s*', '', sent).strip()
    s = re.sub(r'^(첫째|둘째|셋째|넷째|먼저|다음으로|끝으로)[,\s]*', '', s).strip()
    s = re.sub(r'\s+', ' ', s)
    if len(s) <= limit:
        return s.rstrip('.。')
    cut = s[:limit]
    if ' ' in cut:
        cut = cut[:cut.rfind(' ')]
    return cut + "…"


def build_local_summary(title: str, scripture: str, full_text: str) -> str:
    """
    AI 미연결/실패 시에도 '원고에서 실제로 뽑아낸' 요약을 만든다.
    (예전처럼 아무 설교에나 붙는 템플릿 문구를 쓰지 않는다.)
    """
    a = analyze_manuscript(full_text)
    if a["char_count"] < 80:
        return ("⚠️ 설교 원고가 비어 있거나 너무 짧습니다.\n"
                "[📤 새 설교 등록/원고작성] 메뉴에서 설교문 전문을 먼저 등록해 주세요.")

    kw = a["keywords"][:8]
    refs = a["refs"][:8]
    secs = a["sections"]
    apps = a["applications"]

    # 명제는 대지 제목과 중복되지 않는 문장으로 고른다
    sec_heads = {s["headline"] for s in secs}
    prop_src = next((s for s in a["key_sentences"] if s[:90] not in sec_heads),
                    (a["key_sentences"][0] if a["key_sentences"]
                     else (a["paragraphs"][0] if a["paragraphs"] else title)))

    lines = []
    lines.append(f"🧭 [원고 자동 추출 요약]  ·  {title}  ·  {scripture}")
    lines.append(f"※ AI 서버 연결 없이, 업로드하신 설교 원고 본문에서 직접 추출한 결과입니다. (원고 {a['char_count']:,}자 분석)")
    lines.append("")
    lines.append("🎯 설교 핵심 명제")
    lines.append(f"{prop_src}")
    lines.append("")
    lines.append(f"🔑 이 설교의 핵심 키워드: {', '.join(kw) if kw else '(추출 불가)'}")
    if refs:
        lines.append(f"📖 원고에 인용된 성경 구절: {', '.join(refs)}")
    lines.append("")
    lines.append("📌 원고 흐름에 따른 3대지")
    for i, sec in enumerate(secs[:3], start=1):
        lines.append(f"{i}. {_trim_title(sec['headline'])}")
        lines.append(f"   ▸ 원고 근거: \"{sec['evidence'][:220]}{'…' if len(sec['evidence']) > 220 else ''}\"")
        lines.append("")

    lines.append("💡 원고에서 뽑은 실천 적용")
    if apps:
        for i, s in enumerate(apps[:3], start=1):
            lines.append(f"- {i}. {s}")
    else:
        lines.append("- (원고에 직접적인 권면·적용 문장이 없어 추출하지 못했습니다. AI 요약 버튼을 눌러주세요.)")
    lines.append("")

    lines.append("🙏 마침 기도")
    if a["prayer"]:
        lines.append(a["prayer"][:400])
    else:
        lines.append(f"(원고에 기도문 단락이 없습니다. 본문 {scripture}과 위 적용을 근거로 기도문을 작성해 주세요.)")

    return "\n".join(lines)


# 기존 함수명 호환 유지
def analyze_expository_sermon(title: str, scripture: str, full_text: str) -> str:
    return build_local_summary(title, scripture, full_text)


def build_sermon_context(text: str, max_chars: int = 9000) -> str:
    """
    긴 원고를 프롬프트에 넣을 때 앞부분만 잘라내면 결론/적용이 통째로 날아간다.
    앞 45% / 중간 25% / 뒤 30% 를 균형 있게 샘플링해서 전체 흐름을 살린다.
    """
    text = (text or "").strip()
    if len(text) <= max_chars:
        return text
    head_n = int(max_chars * 0.45)
    mid_n = int(max_chars * 0.25)
    tail_n = max_chars - head_n - mid_n
    mid_start = (len(text) - mid_n) // 2
    return (
        text[:head_n]
        + "\n\n…(중략)…\n\n"
        + text[mid_start:mid_start + mid_n]
        + "\n\n…(중략)…\n\n"
        + text[-tail_n:]
    )


def _get_secret(name: str, default: str = "") -> str:
    """secrets.toml 이 아예 없는 환경에서도 죽지 않도록 감싼다."""
    try:
        v = st.secrets.get(name, default)
        return str(v) if v is not None else default
    except Exception:
        return default


# ==============================================================================
# 영구 저장 설교 DB
# ==============================================================================
SERMON_DB_PATH = "./outputs/sermons_db.json"

# ------------------------------------------------------------------------------
# ★ 영구 보관 엔진 (Streamlit Cloud 대응)
#
#   Streamlit Cloud 의 디스크는 앱이 잠들거나 재배포될 때마다 초기화됩니다.
#   그래서 로컬 파일에만 저장하면 다음 접속 때 설교가 사라집니다.
#   → GitHub Gist(비공개)를 진짜 저장소로 쓰고, 로컬 파일은 캐시로만 씁니다.
#      secrets 에 GITHUB_TOKEN 만 넣어두면 나머지는 앱이 알아서 합니다.
# ------------------------------------------------------------------------------
GIST_FILENAME = "my_sermon_studio_sermons.json"
GIST_DESC = "MY 설교 AI 스튜디오 · 설교 서재 영구 보관"


def _gh_token() -> str:
    return (_get_secret("GITHUB_TOKEN", "") or _get_secret("GH_TOKEN", "")
            or os.environ.get("GITHUB_TOKEN", "")).strip()


def _gh_api(method: str, url: str, payload=None, timeout=20):
    token = _gh_token()
    data = json.dumps(payload).encode("utf-8") if payload is not None else None
    req = urllib.request.Request(url, data=data, method=method, headers={
        "Authorization": f"Bearer {token}",
        "Accept": "application/vnd.github+json",
        "X-GitHub-Api-Version": "2022-11-28",
        "Content-Type": "application/json",
        "User-Agent": "my-sermon-studio",
    })
    with urllib.request.urlopen(req, timeout=timeout) as r:
        body = r.read().decode("utf-8")
    return json.loads(body) if body else {}


def cloud_store_ready() -> bool:
    return bool(_gh_token())


def _gist_find_id():
    """secrets 의 GIST_ID 를 우선 쓰고, 없으면 내 Gist 목록에서 파일명으로 찾는다."""
    fixed = _get_secret("GIST_ID", "").strip()
    if fixed:
        return fixed
    if st.session_state.get("_gist_id"):
        return st.session_state["_gist_id"]
    try:
        for page in (1, 2, 3):
            items = _gh_api("GET", f"https://api.github.com/gists?per_page=100&page={page}")
            if not items:
                break
            for g in items:
                if GIST_FILENAME in (g.get("files") or {}):
                    st.session_state["_gist_id"] = g["id"]
                    return g["id"]
            if len(items) < 100:
                break
    except Exception as e:
        st.session_state["_cloud_error"] = str(e)[:200]
    return None


def cloud_load():
    """Gist 에서 설교 목록을 읽어온다. 실패하면 None."""
    if not cloud_store_ready():
        return None
    gid = _gist_find_id()
    if not gid:
        return []            # 토큰은 있는데 아직 저장소가 없음 → 빈 목록
    try:
        g = _gh_api("GET", f"https://api.github.com/gists/{gid}")
        f = (g.get("files") or {}).get(GIST_FILENAME)
        if not f:
            return []
        if f.get("truncated") and f.get("raw_url"):
            raw = fetch_url_text(f["raw_url"])
        else:
            raw = f.get("content", "")
        data = json.loads(raw) if raw else []
        return data if isinstance(data, list) else []
    except Exception as e:
        st.session_state["_cloud_error"] = str(e)[:200]
        return None


def cloud_save(sermons_list) -> bool:
    if not cloud_store_ready():
        return False
    content = json.dumps(sermons_list, ensure_ascii=False, indent=2)
    payload = {"files": {GIST_FILENAME: {"content": content}}}
    gid = _gist_find_id()
    try:
        if gid:
            _gh_api("PATCH", f"https://api.github.com/gists/{gid}", payload)
        else:
            payload["description"] = GIST_DESC
            payload["public"] = False
            g = _gh_api("POST", "https://api.github.com/gists", payload)
            st.session_state["_gist_id"] = g.get("id")
        st.session_state["_cloud_error"] = ""
        return True
    except Exception as e:
        st.session_state["_cloud_error"] = str(e)[:200]
        return False


def fetch_url_text(url: str) -> str:
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "my-sermon-studio"})
        with urllib.request.urlopen(req, timeout=20) as r:
            return r.read().decode("utf-8")
    except Exception:
        return ""


def _read_local_db():
    try:
        if os.path.exists(SERMON_DB_PATH):
            with open(SERMON_DB_PATH, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, list):
                    return data
    except Exception:
        pass
    return []


def _write_local_db(sermons_list):
    try:
        os.makedirs("./outputs", exist_ok=True)
        with open(SERMON_DB_PATH, "w", encoding="utf-8") as f:
            json.dump(sermons_list, f, ensure_ascii=False, indent=2)
    except Exception:
        pass


def _sig(s: dict) -> str:
    return f"{s.get('title','')}|{s.get('scripture','')}|{len(s.get('text',''))}"


def merge_sermons(a, b):
    """두 목록을 합치되 같은 설교는 한 번만. (사라짐 방지용 안전 병합)"""
    out, seen = [], set()
    for lst in (a or [], b or []):
        for s in lst:
            k = _sig(s)
            if k in seen:
                continue
            seen.add(k)
            out.append(dict(s))
    for i, s in enumerate(out, start=1):
        s["id"] = i
    return out


def get_db_sermons(force_reload: bool = False):
    """
    세션 안에서는 한 번만 원격을 읽고 이후엔 캐시를 쓴다.
    (매 클릭마다 GitHub 를 부르면 느려지므로)
    """
    if not force_reload and st.session_state.get("_db_loaded") and \
            isinstance(st.session_state.get("sermon_library"), list):
        return st.session_state.sermon_library

    local = _read_local_db()
    remote = cloud_load()           # None = 클라우드 미설정/실패

    if remote is None:
        data = local
    else:
        # 클라우드가 진짜 저장소. 다만 로컬에만 있는 설교도 절대 버리지 않는다.
        data = merge_sermons(remote, local)
        if len(data) != len(remote):
            cloud_save(data)

    if not data:
        data = _default_sermons()
        if remote is not None:
            cloud_save(data)

    _write_local_db(data)
    st.session_state.sermon_library = data
    st.session_state._db_loaded = True
    return data


def _default_sermons():
    return [{
        "id": 1,
        "title": "예배와 선교",
        "scripture": "이사야 59:21",
        "testament": "구약",
        "book": "이사야",
        "topic": "예배와 선교",
        "theology": "개혁주의/장로교",
        "date": "2026-08-27",
        "tags": ["선교", "예배", "이사야", "복음"],
        "summary": "",
        "text": """선교전략 용어 중에 1040창(10/40 Window)을 들어보셨을 것입니다. 북위 10도와 40도 사이의 아시아, 북아프리카, 중동 지역을 일컫는 말입니다. 미전도 종족과 빈곤율이 가장 높고 영적 어둠이 짙은 곳입니다.

오늘 본문 이사야 59장 21절은 "내 영과 내 말이 네 입에서 영원토록 떠나지 아니하리라" 말씀하십니다. 참된 예배를 회복할 때 비로소 선교의 문이 열립니다.

우리가 말씀과 성령으로 충만하여 대대손손 복음의 유산을 전수하고 땅끝까지 증인 되는 삶을 살아갑시다."""
    }]


def save_db_sermons(sermons_list):
    """로컬 캐시 + 클라우드(Gist) 양쪽에 저장한다."""
    _write_local_db(sermons_list)
    st.session_state.sermon_library = sermons_list
    st.session_state._db_loaded = True
    if cloud_store_ready():
        if not cloud_save(sermons_list):
            st.warning("⚠️ 클라우드 저장에 실패했습니다. "
                       f"({st.session_state.get('_cloud_error', '')}) "
                       "서재 화면의 [💾 전체 백업] 으로 파일을 꼭 내려받아 두세요.")


def add_sermon_to_db(new_sermon_dict):
    current_list = get_db_sermons()
    existing_ids = [int(s.get("id", 0)) for s in current_list if str(s.get("id", "")).isdigit()]
    new_sermon_dict["id"] = max(existing_ids or [0]) + 1
    current_list.append(new_sermon_dict)
    save_db_sermons(current_list)
    st.session_state.sermon_library = current_list
    return new_sermon_dict


def update_sermon_in_db(sermon_id, updated_summary=None, updated_text=None):
    current_list = get_db_sermons()
    for s in current_list:
        if s.get("id") == sermon_id:
            if updated_summary is not None:
                s["summary"] = updated_summary
            if updated_text is not None:
                s["text"] = updated_text
            break
    save_db_sermons(current_list)
    st.session_state.sermon_library = current_list


# ------------------------------------------------------------------------------
# 설교 전환 시 세션 완전 격리
# ------------------------------------------------------------------------------
DERIVED_KEYS = [
    "small_group_text", "qt5_text", "card_list", "shorts_script_text",
    "sermon_audit_text", "leader_guide_text", "rich_materials",
    "bulletin_column_text", "illustrations", "sermon_quotes", "title_ideas",
    "praise_list", "shorts_rec", "yt_extracted_result",
    "rendered_shorts_out", "vo_audio_path", "verse_card_img",
    "cn_card_idx", "cn_edit_mode", "cn_church_name",
    "temp_generated_sermon", "temp_ai_title", "temp_ai_scrip",
    "ai_fallback_used", "ai_last_error",
]
DERIVED_PREFIXES = (
    "family_worship_", "edit_mode_", "show_copy_", "cn_h_", "cn_b_",
    "edit_sum_area", "edit_grp_area", "edit_qt_area", "edit_sh_area",
    "edit_audit_area", "edit_ldr_area", "edit_fam_", "vc_text_in",
)


def load_sermon_to_workspace(sermon_item, idx=0):
    """서재/등록에서 설교를 불러올 때, 이전 설교의 결과물 캐시를 100% 소각한다."""
    for k in DERIVED_KEYS:
        st.session_state.pop(k, None)
    for k in list(st.session_state.keys()):
        if any(k.startswith(p) for p in DERIVED_PREFIXES):
            st.session_state.pop(k, None)

    st.session_state.current_sermon_id = sermon_item.get("id", 1)
    st.session_state.current_sermon_idx = idx
    st.session_state.sermon_title = sermon_item.get("title", "")
    st.session_state.sermon_scripture = sermon_item.get("scripture", "")
    st.session_state.full_sermon = sermon_item.get("text", "")

    sum_text = sermon_item.get("summary", "")
    if not sum_text or len(sum_text.strip()) < 50:
        sum_text = build_local_summary(
            st.session_state.sermon_title,
            st.session_state.sermon_scripture,
            st.session_state.full_sermon
        )
    st.session_state.sermon_summary_text = sum_text
    st.session_state.dash_active_view = "설교 요약"


# ==============================================================================
# 보안 접속
# ==============================================================================
USER_PIN = _get_secret("APP_PIN", "7777") or "7777"
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False

if not st.session_state.authenticated:
    st.title("🔒 개인 전용 설교 AI 플랫폼")
    pin = st.text_input("접속 비밀번호(PIN)를 입력하세요", type="password")
    if st.button("로그인"):
        if pin == USER_PIN:
            st.session_state.authenticated = True
            st.rerun()
        else:
            st.error("PIN 코드가 올바르지 않습니다.")
    st.stop()


# ==============================================================================
# AI 엔진 (모델 최신화 + 근거 강제 프롬프트 + 실패 투명화)
# ==============================================================================
def get_resolved_api_key():
    sb_k = st.session_state.get("sidebar_api_key_input", "")
    if sb_k and sb_k.strip():
        return sb_k.strip()
    for sec_name in ["GEMINI_API_KEY", "GOOGLE_API_KEY", "API_KEY"]:
        try:
            val = st.secrets.get(sec_name, "")
            if val and str(val).strip():
                return str(val).strip()
        except Exception:
            pass
    for env_name in ["GEMINI_API_KEY", "GOOGLE_API_KEY"]:
        env_val = os.environ.get(env_name, "").strip()
        if env_val:
            return env_val
    return ""


# 2026년 기준 사용 가능한 모델 우선순위 (구형 1.5 계열은 마지막 예비용)
# 목록 조회가 실패했을 때만 쓰는 예비 후보 (최신 세대를 앞에)
FALLBACK_MODELS = [
    "gemini-flash-latest",
    "gemini-3-flash",
    "gemini-3-pro-preview",
    "gemini-pro-latest",
    "gemini-2.0-flash",
    "gemini-flash-lite-latest",
]

# 생성용이 아닌 모델(임베딩·이미지·음성 등)을 걸러내기 위한 키워드
_MODEL_EXCLUDE = ("embedding", "aqa", "vision", "image", "imagen", "tts",
                  "audio", "live", "veo", "learnlm", "gemma")


def _model_rank(name: str):
    """
    모델 이름만 보고 '새롭고 빠른 순'으로 점수를 매긴다.
    구형(2.5 등)이 신형(3)보다 앞서지 않도록 하는 것이 핵심.
    """
    n = name.lower()
    if any(x in n for x in _MODEL_EXCLUDE):
        return None
    if not n.startswith("gemini"):
        return None

    m = re.search(r'gemini-(\d+(?:\.\d+)?)', n)
    if m:
        try:
            ver = float(m.group(1))
        except Exception:
            ver = 0.0
    else:
        # gemini-flash-latest 처럼 버전 없는 별칭 = 항상 최신을 가리키므로 최상위
        ver = 99.0 if "latest" in n else 0.0

    tier = 2.0 if "flash" in n else (1.0 if "pro" in n else 0.0)
    if "lite" in n:
        tier -= 0.6                      # 품질이 낮아 후순위
    if "preview" in n or "-exp" in n or "experimental" in n:
        tier -= 0.4                      # 불안정할 수 있어 후순위
    if "thinking" in n:
        tier -= 0.3
    return (ver, tier)


@st.cache_data(show_spinner=False, ttl=900)
def discover_available_models(key_fingerprint: str):
    """계정이 실제로 호출 가능한 모델을 조회해 '최신 우선'으로 정렬한다."""
    try:
        names = []
        for m in genai.list_models():
            methods = getattr(m, "supported_generation_methods", []) or []
            if "generateContent" in methods:
                names.append(m.name.replace("models/", ""))
        scored = [(n, _model_rank(n)) for n in names]
        scored = [(n, r) for n, r in scored if r is not None]
        scored.sort(key=lambda x: x[1], reverse=True)
        ordered = [n for n, _ in scored]
        return ordered[:10] if ordered else list(FALLBACK_MODELS)
    except Exception:
        return list(FALLBACK_MODELS)


def candidate_models(fingerprint: str):
    """
    실제로 시도할 모델 순서를 만든다.
      ① 사이드바에서 목사님이 직접 고른 모델
      ② 이번 세션에서 마지막으로 성공한 모델
      ③ 계정에서 조회한 최신 모델들
    이미 '없어진 모델'로 판명된 것은 제외한다.
    """
    dead = set(st.session_state.get("_dead_models", []))
    order = []

    picked = st.session_state.get("ai_model_choice", "")
    if picked and picked != "자동 선택 (권장)":
        order.append(picked)

    last_good = st.session_state.get("_last_good_model", "")
    if last_good and last_good not in order:
        order.append(last_good)

    for m in discover_available_models(fingerprint):
        if m not in order:
            order.append(m)

    result = [m for m in order if m not in dead]
    return result or [m for m in FALLBACK_MODELS if m not in dead] or list(FALLBACK_MODELS)


def _classify_api_error(msg: str) -> str:
    m = (msg or "").lower()
    if "429" in m or "quota" in m or "rate limit" in m or "resource_exhausted" in m or "exceeded" in m:
        return "quota"
    if "404" in m or "not found" in m or "no longer available" in m or "not supported" in m:
        return "gone"
    if "403" in m or "permission" in m or "denied" in m:
        return "perm"
    if "api key" in m or "api_key" in m or "invalid argument" in m and "key" in m:
        return "key"
    return "other"


ERROR_HELP = {
    "quota": ("무료 등급 사용량 한도를 넘었습니다 (분당 또는 하루 한도).",
              "1~2분 기다렸다가 다시 눌러 보세요. 계속 뜨면 오늘 하루 한도를 다 쓴 것이니 "
              "내일 다시 시도하시거나, Google AI Studio에서 결제를 연결하면 한도가 크게 늘어납니다. "
              "사이드바에서 **가벼운 모델(flash-lite)** 로 바꾸면 한도가 덜 소모됩니다."),
    "gone": ("선택된 모델이 더 이상 제공되지 않는 구형 모델입니다.",
             "사이드바 [⚙️ AI 연결 설정] → [🔌 연결 테스트] 를 눌러 모델 목록을 새로 고쳐 주세요."),
    "perm": ("API 키에 이 모델을 쓸 권한이 없습니다.",
             "Google AI Studio에서 키를 새로 발급받아 다시 넣어 주세요."),
    "key": ("API 키가 올바르지 않습니다.",
            "사이드바 [⚙️ AI 연결 설정]에서 키를 다시 확인해 주세요."),
    "other": ("AI 서버 호출에 실패했습니다.",
              "잠시 후 다시 시도해 주세요. 계속되면 [🔌 연결 테스트]로 어떤 모델이 되는지 확인해 보세요."),
}


SYSTEM_INSTRUCTION = (
    "당신은 한국 교회 강단 사역을 돕는 최고 수준의 설교 분석 전문가입니다.\n"
    "반드시 지킬 것:\n"
    "1) 사용자가 제공한 <설교원고> 안에 실제로 존재하는 내용만 사용합니다. "
    "원고에 없는 예화·인물·통계·성경구절을 절대 새로 만들어내지 않습니다.\n"
    "2) 어떤 설교에나 붙일 수 있는 상투적 표현(예: '말씀의 깊은 은혜', '믿음의 결단', "
    "'주님의 신실하신 은혜가 충만하기를')만으로 항목을 채우지 않습니다. "
    "반드시 이 원고에만 등장하는 고유한 단어·지명·인물·숫자·예화를 포함시킵니다.\n"
    "3) 영어 사고 과정, 기획 메모, 인사말, 서두 설명 없이 요청된 결과물 본문만 출력합니다.\n"
    "4) 인도자(리더/구역장/셀리더/부모)용 안내는 '[인도자 팁 / 가이드]: ...' 형식으로 씁니다.\n"
    "5) 100% 한국어로 작성합니다.\n"
    "6) 한글 성경 인용은 언제나 **개역개정판**을 기본으로 씁니다. "
    "다른 번역(새번역·쉬운성경·메시지 등)을 쓸 때는 반드시 괄호로 역본명을 밝힙니다."
)


RESEARCH_SYSTEM = (
    "당신은 30년 경력의 한국 장로교 설교자이자 성경 연구자입니다. "
    "히브리어·헬라어 원어, 고대 근동/그레코로만 역사, 교회사, 조직신학, 세계 문학과 예술, "
    "그리고 최근 국내외 시사에 두루 밝습니다."
)

RESEARCH_RULES = """[절대 준수 규칙]
1. 특정 설교 원고에 갇히지 말고, 성경 66권 전체·역사·문학·예술·현대 사건에서 폭넓게 자료를 찾아오십시오.
2. 사실 정확성이 최우선입니다. 인물명·연도·지명·책 제목·구절 장절을 정확히 쓰십시오.
   확신이 없으면 그 항목을 버리고 확실한 다른 항목으로 바꾸십시오.
   그래도 애매하면 문장 끝에 (확인 필요) 라고 반드시 표기하십시오.
3. 존재하지 않는 인용문·저작·사건을 절대 지어내지 마십시오.
4. 어느 본문에나 붙일 수 있는 일반론('은혜가 충만하기를', '믿음의 결단')로 항목을 채우지 마십시오.
   반드시 이 본문·이 주제에만 해당하는 구체적 내용이어야 합니다.
5. 100% 한국어. 영어 사고 과정·머리말·마무리 인사 없이 지정한 형식 그대로만 출력하십시오.
6. 번호는 각 항목(섹션) 안에서 1번부터 다시 시작합니다.
7. 한글 성경 인용은 언제나 **개역개정판**이 기본입니다. 다른 역본을 쓸 경우 반드시
   괄호로 역본명을 밝히십시오. (예: "…" (새번역))"""


def build_research_prompt(task_block: str, scripture: str, topic: str = "",
                          theology: str = "", extra_context: str = "") -> str:
    """
    원고에 매이지 않는 '연구용' 프롬프트.
    설교 본문과 주제만 주고, 성경 전체·역사·문학·현대에서 자료를 끌어오게 한다.
    """
    lens = ""
    if theology and theology in THEOLOGY_LENSES:
        lens = f"\n[신학적 관점]\n{theology}\n{THEOLOGY_LENSES[theology]['guide']}\n"
    topic_line = f"설교 주제 / 강조 포인트: {topic}" if topic.strip() else \
        "설교 주제: (지정 없음 — 본문 자체가 말하는 중심 주제를 스스로 규정하고 그것을 기준으로 작업하십시오)"

    return f"""{RESEARCH_SYSTEM}

[작업 대상]
설교 본문: {scripture}
{topic_line}
{lens}{extra_context}
{RESEARCH_RULES}

{task_block}"""


def build_grounded_prompt(task_block: str, ctx_chars: int = 9000, extra: str = "") -> str:
    """모든 생성 작업에 공통으로 붙는 '근거 강제' 프롬프트 골격"""
    title = st.session_state.get("sermon_title", "")
    scripture = st.session_state.get("sermon_scripture", "")
    raw = st.session_state.get("full_sermon", "")
    a = analyze_manuscript(raw)
    body = build_sermon_context(raw, ctx_chars)

    kw = ", ".join(a["keywords"][:15]) or "(없음)"
    refs = ", ".join(a["refs"][:15]) or "(원고 내 명시적 인용 없음)"

    return f"""[작업 대상 설교]
제목: {title}
설교 본문(대표 성구): {scripture}

[이 원고에서 자동 추출된 고유 키워드]
{kw}

[이 원고에 실제로 인용된 성경 구절]
{refs}

<설교원고>
{body}
</설교원고>

[절대 준수 규칙]
- 위 <설교원고>에 실제로 나오는 내용만 근거로 삼습니다. 원고에 없는 예화/인물/사건/통계/성구를 지어내지 마십시오.
- 위 '고유 키워드' 중 최소 5개 이상을 결과물에 자연스럽게 반영하십시오.
- 모든 항목은 이 설교에만 해당되어야 합니다. 다른 설교에도 그대로 쓸 수 있는 일반론 문장은 금지입니다.
- 대표 성구는 {scripture} 입니다. 원고에 인용되지 않은 다른 성경 구절을 임의로 끌어오지 마십시오.
- 한글 성경 인용은 개역개정판을 기본으로 씁니다. 다른 역본은 괄호로 역본명을 밝히십시오.
- 100% 한국어. 영어 메모·머리말·마무리 인사 금지. 아래 형식 그대로만 출력.

{task_block}
{extra}"""


def clean_korean_output(text: str) -> str:
    if not text:
        return ""
    lines = text.split("\n")
    cleaned = []
    for line in lines:
        s = line.strip()
        if not s:
            cleaned.append("")
            continue
        # 순수 영어 메모 줄만 제거 (한글이 섞인 줄은 절대 건드리지 않음)
        k = len(re.findall(r'[가-힣]', s))
        e = len(re.findall(r'[a-zA-Z]', s))
        if k == 0 and e > 20:
            continue
        if re.match(r'^(Okay|Sure|Certainly|Here is|Here\'s|I will|Let me|Note:)\b', s, re.IGNORECASE):
            continue
        cleaned.append(line.rstrip())
    result = re.sub(r'\n{3,}', '\n\n', "\n".join(cleaned).strip())
    return result if result else text


def extract_json_from_text(text):
    if not text:
        return None
    raw = str(text).strip()
    raw = re.sub(r"^`{1,3}[a-zA-Z0-9_-]*\s*", "", raw)
    raw = re.sub(r"\s*`{1,3}$", "", raw).strip()
    try:
        return json.loads(raw)
    except Exception:
        m = re.search(r"(\{[\s\S]*\})", raw)
        if m:
            try:
                return json.loads(m.group(1))
            except Exception:
                pass
    return None


SEARCH_TOOL_VARIANTS = [
    [{"google_search": {}}],
    [{"google_search_retrieval": {}}],
    "google_search_retrieval",
]


def get_ai_response(prompt: str, is_json: bool = True, temperature: float = 0.35,
                    kind: str = "summary", card_count: int = 7,
                    use_search: bool = False, max_tokens: int = 8192):
    """
    AI 호출. 실패하면 조용히 템플릿을 뱉지 않고,
    세션에 실패 사유를 기록한 뒤 '원고 기반' 대체 결과를 돌려준다.

    use_search=True 이면 Gemini 의 Google 검색 근거(grounding) 를 먼저 시도한다.
    (계정/모델이 지원하지 않으면 자동으로 검색 없이 재시도)
    """
    st.session_state.ai_fallback_used = False
    st.session_state.ai_last_error = ""
    st.session_state.ai_error_kind = ""
    st.session_state.ai_error_detail = []

    active_key = get_resolved_api_key()
    if not active_key:
        st.session_state.ai_fallback_used = True
        st.session_state.ai_error_kind = "key"
        st.session_state.ai_error_detail = ["API 키가 비어 있습니다."]
        st.session_state.ai_last_error = "Gemini API 키가 설정되지 않았습니다. (사이드바 ⚙️ AI 연결 설정)"
        r = grounded_fallback(kind, is_json, card_count)
        return r if is_json else fix_list_numbering(r)

    try:
        genai.configure(api_key=active_key)
        os.environ["GOOGLE_API_KEY"] = active_key
        os.environ["GEMINI_API_KEY"] = active_key
    except Exception as e:
        st.session_state.ai_fallback_used = True
        st.session_state.ai_error_kind = "key"
        st.session_state.ai_error_detail = [str(e)[:200]]
        st.session_state.ai_last_error = f"API 키 설정 오류: {e}"
        r = grounded_fallback(kind, is_json, card_count)
        return r if is_json else fix_list_numbering(r)

    fingerprint = hashlib.sha256(active_key.encode()).hexdigest()[:16]
    models = candidate_models(fingerprint)

    errors, kinds = [], []
    for model_name in models[:5]:                    # 너무 많이 시도하면 한도만 낭비된다
        try:
            try:
                model = genai.GenerativeModel(model_name, system_instruction=SYSTEM_INSTRUCTION)
            except Exception:
                model = genai.GenerativeModel(model_name)

            def _call(cfg, with_search):
                """검색 근거 → 검색 없이 → 토큰 한도 제거, 순서로 물러나며 시도."""
                if with_search:
                    for tv in SEARCH_TOOL_VARIANTS:
                        try:
                            r = model.generate_content(prompt, generation_config=cfg, tools=tv)
                            st.session_state.ai_search_used = True
                            return r
                        except Exception:
                            continue
                st.session_state.ai_search_used = False
                try:
                    return model.generate_content(prompt, generation_config=cfg)
                except Exception as e1:
                    # 모델이 지원하지 않는 옵션(max_output_tokens, response_mime_type)이면 빼고 재시도
                    if _classify_api_error(str(e1)) in ("gone", "quota", "perm", "key"):
                        raise
                    slim = {"temperature": cfg.get("temperature", 0.4)}
                    return model.generate_content(prompt, generation_config=slim)

            if is_json:
                cfg = {"response_mime_type": "application/json", "temperature": temperature,
                       "max_output_tokens": max_tokens}
                res = _call(cfg, False)   # JSON 강제 출력과 검색 도구는 함께 못 쓰는 모델이 많다
                parsed = extract_json_from_text(getattr(res, "text", ""))
                if parsed:
                    st.session_state.ai_model_used = model_name
                    st.session_state["_last_good_model"] = model_name
                    return parsed
                errors.append(f"{model_name}: JSON 형식으로 응답하지 못함")
                kinds.append("other")
            else:
                res = _call({"temperature": temperature, "max_output_tokens": max_tokens}, use_search)
                txt = getattr(res, "text", "") or ""
                cleaned = clean_korean_output(txt)
                if cleaned and len(cleaned.strip()) > 60:
                    st.session_state.ai_model_used = model_name
                    st.session_state["_last_good_model"] = model_name
                    return fix_list_numbering(cleaned)
                errors.append(f"{model_name}: 응답이 비었거나 너무 짧음")
                kinds.append("other")

        except Exception as e:
            msg = str(e)
            k = _classify_api_error(msg)
            kinds.append(k)
            errors.append(f"{model_name} → {msg[:150]}")
            if k == "gone":
                # 폐기된 모델은 이번 세션 동안 다시 시도하지 않는다
                dead = set(st.session_state.get("_dead_models", []))
                dead.add(model_name)
                st.session_state["_dead_models"] = list(dead)
                try:
                    discover_available_models.clear()
                except Exception:
                    pass
            if k in ("key", "perm"):
                break                     # 키 문제면 다른 모델도 소용없다
            if kinds.count("quota") >= 3:
                break                     # 한도 초과가 반복되면 더 부르지 않는다(한도만 더 깎임)
            continue

    st.session_state.ai_fallback_used = True
    # 가장 많이 나온 오류 유형으로 안내 문구를 정한다
    main_kind = "other"
    for k in ("key", "perm", "quota", "gone"):
        if k in kinds:
            main_kind = k
            break
    st.session_state.ai_error_kind = main_kind
    st.session_state.ai_error_detail = errors
    st.session_state.ai_last_error = ERROR_HELP[main_kind][0]
    res = grounded_fallback(kind, is_json, card_count)
    return res if is_json else fix_list_numbering(res)


# ------------------------------------------------------------------------------
# 원고 기반 대체 생성기 (AI 실패 시에도 '이 설교의' 내용이 나오도록)
# ------------------------------------------------------------------------------
def grounded_fallback(kind: str, is_json: bool, card_count: int = 7):
    title = st.session_state.get("sermon_title", "은혜의 말씀")
    scripture = st.session_state.get("sermon_scripture", "본문 말씀")
    raw = st.session_state.get("full_sermon", "")
    a = analyze_manuscript(raw)
    keys = a["keywords"]
    sents = a["key_sentences"]
    secs = a["sections"]
    apps = a["applications"]

    def pick(i, default=""):
        return sents[i] if len(sents) > i else default

    if is_json:
        if kind == "cards":
            n = max(3, int(card_count))
            cards = [{
                "card_number": 1,
                "headline": f"「 {title} 」",
                "body_text": f"{scripture}\n\n{pick(0, title)}"
            }]
            for i, sec in enumerate(secs[:3], start=2):
                cards.append({
                    "card_number": i,
                    "headline": f"0{i-1}. {_trim_title(sec['headline'], 28)}",
                    "body_text": sec["evidence"][:180]
                })
            idx = len(cards) + 1
            for s in sents[3:3 + max(0, n - len(cards) - 2)]:
                cards.append({"card_number": idx, "headline": f"0{idx-1}. 원고 속 핵심 문장",
                              "body_text": s[:180]})
                idx += 1
            if apps:
                cards.append({
                    "card_number": idx,
                    "headline": "💡 삶의 적용",
                    "body_text": "\n".join(f"{k}. {s[:90]}" for k, s in enumerate(apps[:3], start=1))
                })
                idx += 1
            cards.append({
                "card_number": idx,
                "headline": "🙏 함께 드리는 기도",
                "body_text": (a["prayer"][:180] if a["prayer"]
                              else f"{scripture} 말씀을 붙들고 {', '.join(keys[:3])}의 은혜 안에서 살게 하옵소서.")
            })
            for c in cards:
                pass
            cards = cards[:n]
            for i, c in enumerate(cards, start=1):
                c["card_number"] = i
            while len(cards) < n:
                extra_i = len(cards) + 1
                src = sents[(extra_i + 3) % max(1, len(sents))] if sents else title
                cards.append({"card_number": extra_i, "headline": f"0{extra_i-1}. 말씀 되새김",
                              "body_text": src[:180]})
            return {"cards": cards}

        if kind == "praise":
            topic = ", ".join(keys[:3]) or title
            return {
                "hymns": [f"(AI 미연결) '{topic}' 주제 찬송 검색 필요 - 새찬송가 색인 참조"],
                "gospel_songs": [f"(AI 미연결) '{topic}' 주제 복음성가 검색 필요"],
                "ccm": [f"(AI 미연결) '{topic}' 주제 CCM 검색 필요"],
            }

        if kind == "shorts_meta":
            base = [s for s in sents[:5]] or [title]
            return {
                "titles": [f"{i+1}. {_trim_title(s, 34)}" for i, s in enumerate(base)],
                "hashtags": ["#주일설교", "#말씀묵상"] + [f"#{k}" for k in keys[:5]],
            }
        return {}

    # ---- 텍스트 계열 ----
    kw_line = ", ".join(keys[:8]) or "(추출 불가)"
    warn = ("⚠️ AI 서버에 연결되지 않아, 업로드하신 설교 원고에서 직접 추출한 결과를 표시합니다.\n"
            "   사이드바 [⚙️ AI 연결 설정]에서 Gemini API 키를 확인해 주세요.\n")

    if kind == "rich":
        body = [warn, f"[원고 기반 자료 추출: {title} / {scripture}]", "",
                f"🔑 핵심 키워드: {kw_line}",
                f"📖 원고에 인용된 성구: {', '.join(a['refs'][:10]) or '(없음)'}", "",
                "📝 원고에서 뽑은 인용 가치 있는 문장"]
        for i, s in enumerate(sents[:6], start=1):
            body.append(f"- {i}. {s}")
        return "\n".join(body)

    if kind == "leader":
        out = [warn, f"[소그룹 리더 가이드(원고 추출): {title} / {scripture}]", "",
               "1. 🎯 이번 주 모임의 핵심 방향",
               f"- [인도자 팁 / 가이드]: 원고의 중심 문장 — \"{pick(0, title)}\"", "",
               "2. 📖 원고 흐름 요약(리더 숙지용)"]
        for i, sec in enumerate(secs[:3], start=1):
            out.append(f"- {i}. {_trim_title(sec['headline'])}")
            out.append(f"  ▸ 근거: \"{sec['evidence'][:160]}…\"")
        out += ["", "3. 💬 나눔으로 이어갈 원고 속 권면"]
        for s in (apps[:3] or ["(원고에 권면 문장이 없습니다)"]):
            out.append(f"- {s}")
        return "\n".join(out)

    if kind == "smallgroup":
        out = [warn, f"[소그룹 나눔지(원고 추출): {title} / {scripture}]", "",
               "1. 마음 열기",
               f"- [인도자 팁 / 가이드]: 이번 설교의 키워드 '{keys[0] if keys else title}'로 시작하세요.",
               "", "2. 말씀 속으로"]
        for i, sec in enumerate(secs[:3], start=1):
            out.append(f"- {i}. \"{_trim_title(sec['headline'], 60)}\" 이 대목에서 마음에 남은 것은 무엇인가요?")
        out += ["", "3. 삶 속으로"]
        for i, s in enumerate(apps[:2] or [f"{scripture} 말씀을 이번 주 어디에 적용하시겠습니까?"], start=1):
            out.append(f"- {i}. \"{s}\" — 나에게는 어떻게 적용됩니까?")
        out += ["", "4. 마침 기도", a["prayer"][:300] if a["prayer"] else f"{title} 말씀대로 살게 하옵소서. 아멘."]
        return "\n".join(out)

    if kind == "qt":
        days = ["월요일", "화요일", "수요일", "목요일", "금요일"]
        out = [warn, f"[주간 QT 5일치(원고 추출): {title} / {scripture}]", ""]
        pool = sents if sents else [title] * 5
        for i, d in enumerate(days):
            src = pool[i % len(pool)]
            ref = a["refs"][i % len(a["refs"])] if a["refs"] else scripture
            out += [f"📅 {d}",
                    f"- 📖 본문 구절: {ref}",
                    f"- 💡 말씀 묵상: {src}",
                    f"- 🎯 삶의 적용: {(apps[i % len(apps)] if apps else '오늘 이 말씀을 한 가지 행동으로 옮겨 보십시오.')}",
                    f"- 🙏 오늘의 기도: 주님, 이 말씀을 붙들게 하옵소서.", ""]
        return "\n".join(out)

    if kind == "family":
        out = [warn, f"[가정예배 순서지(원고 추출): {title} / {scripture}]", "",
               "1. 찬양 및 신앙고백",
               "- [인도자 팁 / 가이드]: 가족이 함께 아는 찬송으로 시작하세요.", "",
               "2. 함께 읽는 말씀", f"- 본문: {scripture}", "",
               "3. 3분 가족 메시지(원고 근거)"]
        for sec in secs[:2]:
            out.append(f"- {_trim_title(sec['headline'], 60)}")
        out += ["", "4. 나눔 질문",
                f"- 1. 오늘 말씀에서 '{keys[0] if keys else '은혜'}'는 우리 가정에 어떤 의미일까요?",
                "- 2. 이번 주 우리 가족이 실천할 한 가지는 무엇인가요?", "",
                "5. 마무리 기도", a["prayer"][:250] if a["prayer"] else "우리 가정을 지켜주옵소서. 아멘."]
        return "\n".join(out)

    if kind == "audit":
        out = [warn, f"[원고 구조 자동 점검: {title} / {scripture}]", "",
               f"- 원고 분량: {a['char_count']:,}자 (약 {max(1, a['char_count']//350)}분 분량 추정)",
               f"- 문단 수: {len(a['paragraphs'])} / 문장 수: {len(a['sentences'])}",
               f"- 인용 성구: {len(a['refs'])}개 — {', '.join(a['refs'][:8]) or '없음'}",
               f"- 권면·적용 문장: {len(a['applications'])}개",
               f"- 핵심 키워드: {kw_line}", "",
               "🔎 구조 관찰"]
        if len(a['refs']) < 2:
            out.append("- 본문 인용이 적습니다. 대지마다 근거 구절을 명시하면 강해적 설득력이 올라갑니다.")
        if len(a['applications']) < 3:
            out.append("- 구체적 적용 문장이 부족합니다. '이번 주에 ~하십시오' 형태의 실행 문장을 보강해 보십시오.")
        if not a['prayer']:
            out.append("- 마무리 기도 단락이 확인되지 않습니다.")
        if len(out) == 8:
            out.append("- 구조상 큰 결손은 발견되지 않았습니다.")
        return "\n".join(out)

    if kind == "bulletin":
        para = a["paragraphs"]
        out = [warn, f"✍️ 「{title}」", ""]
        out.append(pick(0, title))
        for pp in para[1:4]:
            out.append("")
            out.append(pp[:260])
        if apps:
            out += ["", f"이번 한 주간, {apps[0]}"]
        out += ["", f"📖 이번 주 말씀 — {scripture}", "", "— 드림"]
        return "\n".join(out)

    if kind == "illust":
        out = [warn, f"[예화 자료(원고 추출): {title} / {scripture}]", "",
               "📖 성경 예화 — AI 연결 시 성경 인물·사건 예화가 생성됩니다.",
               f"- 1. 원고에 인용된 구절: {', '.join(a['refs'][:5]) or '(없음)'}", "",
               "🌍 원고에서 발견한 예화성 대목"]
        for i, sname in enumerate(sents[:4], start=1):
            out.append(f"- {i}. {sname}")
        out += ["", "🏛️ 역사·교회사 예화 및 현대 예화는 AI 연결 후 생성됩니다."]
        return "\n".join(out)

    if kind == "quotes":
        return (warn + f"\n[설교 명언: {title}]\n\n"
                "🗣️ 명언 자료는 AI 전용 기능입니다. 인물에게 하지 않은 말을 지어내지 않기 위해\n"
                "AI 연결 없이는 명언을 생성하지 않습니다.\n\n"
                "사이드바 [⚙️ AI 연결 설정]에서 Gemini API 키를 등록한 뒤 다시 눌러 주세요.\n\n"
                f"(원고 핵심 키워드: {kw_line})")

    if kind == "titles5":
        base = [s2 for s2 in sents[:5]] or [title]
        out = [warn, "🏷️ 원고에서 뽑은 설교 제목 후보", ""]
        for i, sname in enumerate(base, start=1):
            out.append(f"- {i}. 「{_trim_title(sname, 22)}」")
            out.append(f"   ▸ 이유: 원고 문장 \"{sname[:60]}…\" 에서 추출")
        out += ["", f"📌 핵심 키워드로 만든 부제: {', '.join(keys[:5])}"]
        return "\n".join(out)

    if kind == "shorts_script":
        pool = sents or [title]
        out = [warn, f"[쇼츠 대본(원고 추출): {title} / {scripture}]", ""]
        labels = ["🎬 1) 감동·위로형", "💡 2) 질문·호기심형", "🔥 3) 결단 선포형"]
        for i, lab in enumerate(labels):
            hook = pool[i % len(pool)]
            mid = pool[(i + 3) % len(pool)]
            end = (apps[i % len(apps)] if apps else f"{scripture} 말씀을 오늘 붙드십시오.")
            out += [lab,
                    f"- 후킹(0~5초): {hook[:60]}",
                    f"- 본론(5~45초): {mid}",
                    f"- 결단(45~60초): {end}",
                    f"- 자막 키워드: {', '.join(keys[i*2:i*2+5]) or ', '.join(keys[:5])}", ""]
        return "\n".join(out)

    if kind in ("manna", "today_verse", "versions"):
        return {} if is_json else ""

    if kind in ("research", "context"):
        return ("⚠️ 본문 연구 도구는 AI 전용 기능입니다.\n\n"
                "성경 66권 전체·역사·문학·현대 자료에서 내용을 찾아와야 하므로, "
                "설교 원고만으로는 만들 수 없습니다.\n"
                "사이드바 [⚙️ AI 연결 설정]에서 Gemini API 키를 등록한 뒤 다시 눌러 주세요.\n\n"
                f"(사유: {st.session_state.get('ai_last_error', '연결 실패')})")

    if kind == "sermon_write":
        return ("⚠️ 강해설교문 자동 작성은 AI 전용 기능입니다.\n"
                "사이드바 [⚙️ AI 연결 설정]에서 Gemini API 키를 등록한 뒤 다시 시도해 주세요.\n"
                f"(사유: {st.session_state.get('ai_last_error', '연결 실패')})")

    return build_local_summary(title, scripture, raw)


# ==============================================================================
# 폰트 (캐싱)
# ==============================================================================
# 시스템에 설치돼 있을 만한 한글 폰트 경로 (우선순위 순)
SYSTEM_FONTS_REGULAR = [
    "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
    "/usr/share/fonts/truetype/nanum/NanumBarunGothic.ttf",
    "/usr/share/fonts/opentype/nanum/NanumGothic.ttf",
    "/usr/share/fonts/nanum/NanumGothic.ttf",
    "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    "C:/Windows/Fonts/malgun.ttf",
    "/System/Library/Fonts/AppleSDGothicNeo.ttc",
]
SYSTEM_FONTS_BOLD = [
    "/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf",
    "/usr/share/fonts/opentype/nanum/NanumGothicBold.ttf",
    "/usr/share/fonts/truetype/noto/NotoSansCJK-Bold.ttc",
    "C:/Windows/Fonts/malgunbd.ttf",
] + SYSTEM_FONTS_REGULAR

# 폰트가 없을 때 내려받을 미러 (User-Agent 없으면 403 이 나므로 반드시 헤더를 붙인다)
FONT_MIRRORS = {
    "regular": [
        "https://cdn.jsdelivr.net/gh/google/fonts@main/ofl/nanumgothic/NanumGothic-Regular.ttf",
        "https://raw.githubusercontent.com/google/fonts/main/ofl/nanumgothic/NanumGothic-Regular.ttf",
        "https://cdn.jsdelivr.net/gh/projectnoonnu/noonfonts@master/NanumGothic.ttf",
    ],
    "bold": [
        "https://cdn.jsdelivr.net/gh/google/fonts@main/ofl/nanumgothic/NanumGothic-Bold.ttf",
        "https://raw.githubusercontent.com/google/fonts/main/ofl/nanumgothic/NanumGothic-Bold.ttf",
    ],
}
FONT_DIR = "./fonts"


def _download_font(kind: str, dest: str) -> bool:
    for url in FONT_MIRRORS.get(kind, []):
        try:
            req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
            with urllib.request.urlopen(req, timeout=20) as r:
                data = r.read()
            if len(data) > 200000:                       # 정상 TTF 인지 최소 검증
                os.makedirs(FONT_DIR, exist_ok=True)
                with open(dest, "wb") as f:
                    f.write(data)
                return True
        except Exception:
            continue
    return False


@st.cache_resource(show_spinner=False)
def ensure_korean_fonts():
    """(regular_path, bold_path) 반환. 확보 실패 시 (None, None)."""
    reg = next((p for p in SYSTEM_FONTS_REGULAR if os.path.exists(p)), None)
    bold = next((p for p in SYSTEM_FONTS_BOLD if os.path.exists(p)), None)

    if not reg:
        cand = os.path.join(FONT_DIR, "NanumGothic-Regular.ttf")
        if os.path.exists(cand) or _download_font("regular", cand):
            reg = cand
    if not bold:
        cand = os.path.join(FONT_DIR, "NanumGothic-Bold.ttf")
        if os.path.exists(cand) or _download_font("bold", cand):
            bold = cand
    return reg, (bold or reg)


@st.cache_resource(show_spinner=False)
def init_korean_font():
    """
    reportlab(PDF)용 한글 폰트 등록.
    reportlab은 .ttc(폰트 컬렉션)나 CFF 계열을 제대로 못 읽는 경우가 많아
    반드시 순수 .ttf 를 우선 시도하고, 없으면 NanumGothic 을 내려받는다.
    """
    reg, _ = ensure_korean_fonts()

    tries = []
    for p in SYSTEM_FONTS_REGULAR:
        if os.path.exists(p) and p.lower().endswith(".ttf"):
            tries.append((p, None))

    dl = os.path.join(FONT_DIR, "NanumGothic-Regular.ttf")
    if os.path.exists(dl) or _download_font("regular", dl):
        tries.append((dl, None))

    if reg and reg.lower().endswith((".ttc", ".otc")):
        tries.append((reg, 0))          # 컬렉션은 첫 번째 서브폰트로 시도

    for path, sub in tries:
        try:
            if sub is None:
                pdfmetrics.registerFont(TTFont("NanumKorean", path))
            else:
                pdfmetrics.registerFont(TTFont("NanumKorean", path, subfontIndex=sub))
            return "NanumKorean"
        except Exception:
            continue
    return "Helvetica"


PDF_FONT_NAME = init_korean_font()
KOREAN_FONT_OK = (PDF_FONT_NAME != "Helvetica")


@st.cache_resource(show_spinner=False)
def get_pil_font(size: int):
    _, bold = ensure_korean_fonts()
    if bold:
        try:
            return PIL.ImageFont.truetype(bold, size)
        except Exception:
            pass
    return PIL.ImageFont.load_default()


# ==============================================================================
# 배경 이미지
# ==============================================================================
@st.cache_data(show_spinner=False, ttl=86400)
def fetch_image_bytes(url: str):
    try:
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=6) as response:
            return response.read()
    except Exception:
        return None


# ------------------------------------------------------------------------------
# 무한 배경 엔진
#   - 고정된 8장이 아니라, 시드(설교 제목 + 카드 번호 + 새로고침 횟수)로
#     매번 다른 이미지를 만들어 낸다.
#   - 1순위: 사진 CDN(키 불필요) / 2순위: 다른 CDN / 3순위: 절대 실패하지 않는
#     프로시저럴 그라데이션(무한 조합)
# ------------------------------------------------------------------------------
BG_THEMES = {
    "자연 · 풍경": ["nature", "mountain", "forest", "sunrise", "field", "lake"],
    "하늘 · 빛": ["sky", "clouds", "light", "sunbeam", "sunset", "stars"],
    "바다 · 물": ["ocean", "sea", "wave", "river", "water", "shore"],
    "길 · 여정": ["path", "road", "journey", "bridge", "desert", "trail"],
    "예배 · 경건": ["church", "cathedral", "candle", "cross", "chapel", "stained-glass"],
    "도시 · 일상": ["city", "window", "cafe", "street", "home", "book"],
    "추상 · 그라데이션": [],
}

# 그라데이션 팔레트(무한 조합의 씨앗)
BG_PALETTES = [
    ((14, 22, 58), (76, 29, 149), (14, 116, 144)),
    ((3, 26, 42), (12, 74, 110), (13, 148, 136)),
    ((30, 12, 44), (109, 40, 217), (219, 39, 119)),
    ((7, 20, 34), (30, 64, 175), (56, 189, 248)),
    ((28, 16, 8), (146, 64, 14), (245, 158, 11)),
    ((10, 30, 24), (6, 95, 70), (52, 211, 153)),
    ((26, 8, 22), (157, 23, 77), (251, 113, 133)),
    ((10, 12, 30), (49, 46, 129), (129, 140, 248)),
    ((18, 24, 12), (63, 98, 18), (163, 230, 53)),
    ((24, 10, 30), (91, 33, 182), (232, 121, 249)),
    ((6, 18, 28), (7, 89, 133), (125, 211, 252)),
    ((32, 14, 6), (154, 52, 18), (253, 186, 116)),
]


def _seed_int(seed: str) -> int:
    return int(hashlib.sha256(str(seed).encode("utf-8")).hexdigest()[:12], 16)


@st.cache_data(show_spinner=False, max_entries=256)
def make_gradient_bg(seed: str, size=(1080, 1080)) -> bytes:
    """네트워크 없이도 항상 성공하는 프로시저럴 배경 (조합 사실상 무한)"""
    W, H = size
    n = _seed_int(seed)
    c0, c1, c2 = BG_PALETTES[n % len(BG_PALETTES)]
    angle = (n >> 5) % 4          # 대각선 방향 4종
    swirl = ((n >> 9) % 3)        # 광원 위치 3종

    # 색상 미세 변주 — 같은 팔레트라도 매번 다른 색감이 나오도록
    def jitter(c, shift):
        return tuple(max(0, min(255, v + ((n >> shift) % 46) - 22)) for v in c)

    c0, c1, c2 = jitter(c0, 13), jitter(c1, 17), jitter(c2, 21)

    # 작은 캔버스에 그린 뒤 확대 — 훨씬 빠르고 결과는 더 부드럽다
    S = 96
    small = PIL.Image.new("RGB", (S, S))
    px = small.load()
    for y in range(S):
        for x in range(S):
            if angle == 0:
                t = (x / S) * 0.5 + (y / S) * 0.5
            elif angle == 1:
                t = (1 - x / S) * 0.5 + (y / S) * 0.5
            elif angle == 2:
                t = y / S
            else:
                t = x / S
            if t < 0.5:
                k = t * 2
                a, b_ = c0, c1
            else:
                k = (t - 0.5) * 2
                a, b_ = c1, c2
            px[x, y] = (int(a[0] + (b_[0] - a[0]) * k),
                        int(a[1] + (b_[1] - a[1]) * k),
                        int(a[2] + (b_[2] - a[2]) * k))

    # 부드러운 광원 한 점 → 사진 같은 깊이감
    glow = PIL.Image.new("L", (S, S), 0)
    gd = PIL.ImageDraw.Draw(glow)
    cx = int(S * [0.25, 0.72, 0.5][swirl])
    cy = int(S * [0.22, 0.30, 0.78][swirl])
    rad = int(S * (0.30 + (n % 7) * 0.025))
    gd.ellipse([cx - rad, cy - rad, cx + rad, cy + rad], fill=140)
    glow = glow.filter(PIL.ImageFilter.GaussianBlur(radius=S * 0.18))

    light = PIL.Image.new("RGB", (S, S), (255, 246, 214))
    small = PIL.Image.composite(light, small, glow.point(lambda v: int(v * 0.55)))

    base = small.resize((W, H), PIL.Image.LANCZOS)
    base = base.filter(PIL.ImageFilter.GaussianBlur(radius=1.2))

    # 아주 옅은 비네팅으로 텍스트 가독성 확보
    vig = PIL.Image.new("L", (S, S), 0)
    vd = PIL.ImageDraw.Draw(vig)
    vd.ellipse([-S * 0.2, -S * 0.2, S * 1.2, S * 1.2], fill=255)
    vig = vig.filter(PIL.ImageFilter.GaussianBlur(radius=S * 0.12)).resize((W, H), PIL.Image.LANCZOS)
    dark = PIL.Image.new("RGB", (W, H), (0, 0, 0))
    base = PIL.Image.composite(base, dark, vig)

    out = io.BytesIO()
    base.save(out, format="JPEG", quality=90)
    return out.getvalue()


@st.cache_data(show_spinner=False, max_entries=256, ttl=86400)
def get_background_bytes(seed: str, theme: str = "자연 · 풍경", size=(1080, 1080)) -> bytes:
    """
    시드가 달라지면 무조건 다른 배경이 나온다.
    사진 CDN → 실패 시 그라데이션으로 자동 대체(앱은 절대 멈추지 않음).
    """
    W, H = size
    n = _seed_int(seed)

    if theme != "추상 · 그라데이션":
        kws = BG_THEMES.get(theme) or BG_THEMES["자연 · 풍경"]
        kw = kws[n % len(kws)]
        candidates = [
            f"https://picsum.photos/seed/{n % 10_000_000}/{W}/{H}",
            f"https://loremflickr.com/{W}/{H}/{kw}?lock={n % 100000}",
            f"https://source.unsplash.com/random/{W}x{H}/?{kw}&sig={n % 100000}",
        ]
        for url in candidates:
            b = fetch_image_bytes(url)
            if b and len(b) > 20000:
                try:
                    PIL.Image.open(io.BytesIO(b)).verify()
                    return b
                except Exception:
                    continue

    return make_gradient_bg(f"{seed}|{theme}", size=size)


def bg_seed(index: int = 0) -> str:
    """설교 제목 + 인덱스 + '배경 새로고침' 횟수로 시드 생성"""
    shuffle = st.session_state.get("bg_shuffle", 0)
    title = st.session_state.get("sermon_title", "")
    scrip = st.session_state.get("sermon_scripture", "")
    return f"{title}|{scrip}|{index}|{shuffle}"


def current_bg_theme() -> str:
    return st.session_state.get("bg_theme", "자연 · 풍경")


EMOJI_RE = re.compile(
    "[\U0001F000-\U0001FAFF\U00002600-\U000027BF\U0001F1E6-\U0001F1FF⬀-⯿️‍]+"
)


def strip_emoji(text: str) -> str:
    """한글 폰트에는 이모지 글리프가 없어 □(두부)로 찍히므로 이미지에서는 제거한다."""
    return re.sub(r'[ \t]{2,}', ' ', EMOJI_RE.sub('', text or "")).strip()


def wrap_korean_text(text: str, font, max_width: int, draw) -> str:
    if not text:
        return ""
    wrapped = []
    for p in str(text).split('\n'):
        words = p.split(' ')
        curr = ""
        for w in words:
            test = f"{curr} {w}".strip()
            bbox = draw.textbbox((0, 0), test, font=font)
            if (bbox[2] - bbox[0]) > max_width and curr:
                wrapped.append(curr)
                curr = w
            else:
                curr = test
        wrapped.append(curr)
    return "\n".join(wrapped)


# ==============================================================================
# PPTX 투명도 헬퍼 (배경 이미지가 검게 덮이던 버그 수정)
# ==============================================================================
def set_shape_fill_alpha(shape, alpha: float):
    """alpha 0.0(투명) ~ 1.0(불투명)"""
    try:
        spPr = shape.fill._xPr
        solid = spPr.find(qn('a:solidFill'))
        if solid is None:
            return
        clr = solid.find(qn('a:srgbClr'))
        if clr is None:
            return
        for old in clr.findall(qn('a:alpha')):
            clr.remove(old)
        el = etree.SubElement(clr, qn('a:alpha'))
        el.set('val', str(int(max(0.0, min(1.0, alpha)) * 100000)))
    except Exception:
        pass


# ==============================================================================
# 카드 이미지 생성 (캐싱)
# ==============================================================================
@st.cache_data(show_spinner=False, max_entries=64)
def generate_single_card_png_bytes(card_json: str, idx: int, scripture_str: str, church_name: str,
                                   bg_seed_str: str = "", bg_theme: str = "자연 · 풍경") -> bytes:
    card_item = json.loads(card_json)
    img_b = get_background_bytes(bg_seed_str or f"card-{idx}", bg_theme)
    try:
        base_img = PIL.Image.open(io.BytesIO(img_b)).convert("RGBA").resize((1080, 1080))
    except Exception:
        base_img = PIL.Image.new("RGBA", (1080, 1080), (15, 23, 42, 255))

    overlay = PIL.Image.new("RGBA", (1080, 1080), (8, 12, 26, 158))
    combined = PIL.Image.alpha_composite(base_img, overlay)
    # 위·아래에 부드러운 스크림 → 배경 색은 살리고 글자는 또렷하게
    scrim = PIL.Image.new("RGBA", (1080, 1080), (0, 0, 0, 0))
    sd = PIL.ImageDraw.Draw(scrim)
    for yy in range(1080):
        if yy < 300:
            a = int(96 * (1 - yy / 300))
        elif yy > 780:
            a = int(120 * ((yy - 780) / 300))
        else:
            a = 0
        if a:
            sd.line([(0, yy), (1080, yy)], fill=(4, 7, 18, a))
    combined = PIL.Image.alpha_composite(combined, scrim)
    draw = PIL.ImageDraw.Draw(combined)

    M = 100                       # 좌우 여백
    MAXW = 1080 - M * 2
    head_txt = strip_emoji(str(card_item.get("headline", "")))
    body_txt = strip_emoji(str(card_item.get("body_text", "")))

    # 글자 수에 따라 폰트 크기를 자동 조절해 넘침/여백을 줄인다
    h_size = 52 if len(head_txt) <= 18 else (46 if len(head_txt) <= 28 else 38)
    b_size = 34 if len(body_txt) <= 90 else (30 if len(body_txt) <= 160 else 26)
    font_b, font_t, font_s = get_pil_font(h_size), get_pil_font(b_size), get_pil_font(26)

    head = wrap_korean_text(head_txt, font_b, MAXW, draw)
    body = wrap_korean_text(body_txt, font_t, MAXW, draw)

    hb = draw.multiline_textbbox((0, 0), head, font=font_b, spacing=14) if head else (0, 0, 0, 0)
    bb = draw.multiline_textbbox((0, 0), body, font=font_t, spacing=18) if body else (0, 0, 0, 0)
    h_h, b_h = hb[3] - hb[1], bb[3] - bb[1]
    gap = 46 if head and body else 0
    block = h_h + gap + b_h

    # 상단 배지/하단 출처 영역을 뺀 공간의 세로 중앙에 배치
    top_limit, bottom_limit = 200, 860
    y = max(top_limit, top_limit + ((bottom_limit - top_limit) - block) // 2)

    draw.text((M, 86), f"CARD {card_item.get('card_number', idx + 1)}",
              fill=(129, 140, 248, 255), font=font_s)
    draw.line([(M, 132), (M + 70, 132)], fill=(253, 224, 71, 255), width=4)

    if head:
        draw.multiline_text((M, y), head, fill=(253, 224, 71, 255), font=font_b, spacing=14)
    if body:
        draw.multiline_text((M, y + h_h + gap), body, fill=(241, 245, 249, 255), font=font_t, spacing=18)

    if scripture_str:
        draw.text((M, 906), f"「 {scripture_str} 」", fill=(253, 224, 71, 255), font=get_pil_font(30))
    if church_name:
        draw.text((M, 960), church_name, fill=(147, 197, 253, 255), font=font_s)

    out = io.BytesIO()
    combined.convert("RGB").save(out, format="PNG")
    return out.getvalue()


@st.cache_data(show_spinner=False, max_entries=16)
def generate_cardnews_zip_bytes(cards_json: str, scripture_str: str, church_name: str,
                                seed_base: str = "", bg_theme: str = "자연 · 풍경") -> bytes:
    cards = json.loads(cards_json)
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for i, card in enumerate(cards):
            png = generate_single_card_png_bytes(json.dumps(card, ensure_ascii=False), i,
                                                 scripture_str, church_name,
                                                 f"{seed_base}|{i}", bg_theme)
            zf.writestr(f"cardnews_{i+1:02d}.png", png)
    return buf.getvalue()


# ==============================================================================
# 문서 내보내기 — 워드 / PDF / PPT / TXT 모두 '색이 있는' 문서로
#   · 인도자 가이드 = 파란 박스 + 파란 글씨
#   · 섹션 제목    = 보라 배경 강조
#   · 번호 항목    = 번호만 강조색
#   · 원고 근거    = 인디고 인용 블록
# ==============================================================================
DOC_COLORS = {
    "title":  (76, 29, 149),      # 보라 (문서 제목)
    "head":   (91, 33, 182),      # 섹션 제목
    "headbg": "EDE9FE",
    "leader": (7, 89, 133),       # 인도자 가이드 글씨(진한 파랑)
    "leaderlbl": (2, 132, 199),   # 인도자 가이드 라벨
    "leaderbg": "E0F2FE",         # 인도자 가이드 배경(하늘)
    "num":    (219, 39, 119),     # 번호 배지색
    "quote":  (67, 56, 202),      # 원고 근거
    "quotebg": "EEF2FF",
    "body":   (30, 41, 59),
    "meta":   (100, 116, 139),
}

# 구조 인식용 정규식 (렌더링·내보내기 공용)
LEADER_RE = re.compile(r'^\s*[-•]?\s*\[?\s*인도자\s*(팁|가이드)[^\]\n]*\]?\s*[:：]?\s*(.*)$')
SEC_HEAD_RE = re.compile(
    r'^\s*(?:[0-9]{1,2}\s*[\.\)]\s*)?[🎯📌💡🙏📖💬⚠️🏡🎵📝🔎✨🔥📅🧭🔑💎🗣️🏷️✍️📊🎬🏗️🎙️🌍🏛️❓]'
)
LIST_ITEM_RE = re.compile(r'^(\s*)[-•]\s*(\d{1,2})\s*[\.\)]\s*(.*)$')
TOP_NUM_RE = re.compile(r'^(\d{1,2})\s*[\.\)]\s*(.+)$')
QUOTE_RE = re.compile(r'^\s*▸\s*(.*)$')


def parse_doc_blocks(content: str):
    """텍스트를 (종류, 내용) 블록으로 분해한다. 모든 내보내기가 이 결과를 공유한다."""
    blocks = []
    for raw in (content or "").split("\n"):
        line = raw.rstrip()
        if not line.strip():
            blocks.append(("blank", ""))
            continue
        m = LEADER_RE.match(line)
        if m:
            blocks.append(("leader", m.group(2).strip()))
            continue
        m = QUOTE_RE.match(line)
        if m:
            blocks.append(("quote", m.group(1).strip()))
            continue
        m = LIST_ITEM_RE.match(line)
        if m:
            blocks.append(("item", (m.group(2), m.group(3).strip())))
            continue
        if SEC_HEAD_RE.match(line) or (TOP_NUM_RE.match(line) and len(line.strip()) < 60):
            blocks.append(("head", line.strip()))
            continue
        blocks.append(("plain", line.strip()))
    return blocks


# ---------------------------- WORD ----------------------------
def _docx_shade(par, hex_fill):
    try:
        pPr = par._p.get_or_add_pPr()
        shd = OxmlElement('w:shd')
        shd.set(docx_qn('w:val'), 'clear')
        shd.set(docx_qn('w:color'), 'auto')
        shd.set(docx_qn('w:fill'), hex_fill)
        pPr.append(shd)
    except Exception:
        pass


def _docx_left_border(par, hex_color, sz=24):
    try:
        pPr = par._p.get_or_add_pPr()
        pbdr = OxmlElement('w:pBdr')
        left = OxmlElement('w:left')
        left.set(docx_qn('w:val'), 'single')
        left.set(docx_qn('w:sz'), str(sz))
        left.set(docx_qn('w:space'), '8')
        left.set(docx_qn('w:color'), hex_color)
        pbdr.append(left)
        pPr.append(pbdr)
    except Exception:
        pass


@st.cache_data(show_spinner=False, max_entries=48)
def create_docx_bytes(title: str, content: str) -> bytes:
    try:
        doc = Document()

        tp = doc.add_paragraph()
        r = tp.add_run(title)
        r.font.size, r.font.bold = DocxPt(20), True
        r.font.color.rgb = DocxRGB(*DOC_COLORS["title"])
        _docx_shade(tp, "F5F3FF")
        _docx_left_border(tp, "7C3AED", 36)

        mp = doc.add_paragraph()
        mr = mp.add_run(f"작성일 {datetime.now().strftime('%Y-%m-%d')}  ·  MY 설교 AI 스튜디오")
        mr.font.size = DocxPt(8.5)
        mr.font.color.rgb = DocxRGB(*DOC_COLORS["meta"])

        for kind, val in parse_doc_blocks(content):
            if kind == "blank":
                doc.add_paragraph()
            elif kind == "head":
                p = doc.add_paragraph()
                run = p.add_run(val)
                run.font.size, run.font.bold = DocxPt(13), True
                run.font.color.rgb = DocxRGB(*DOC_COLORS["head"])
                _docx_shade(p, DOC_COLORS["headbg"])
                _docx_left_border(p, "7C3AED", 28)
            elif kind == "leader":
                p = doc.add_paragraph()
                lbl = p.add_run("💡 인도자 가이드   ")
                lbl.font.size, lbl.font.bold = DocxPt(10), True
                lbl.font.color.rgb = DocxRGB(*DOC_COLORS["leaderlbl"])
                body = p.add_run(val)
                body.font.size, body.font.bold = DocxPt(10.5), True
                body.font.color.rgb = DocxRGB(*DOC_COLORS["leader"])
                _docx_shade(p, DOC_COLORS["leaderbg"])
                _docx_left_border(p, "0284C7", 36)
            elif kind == "quote":
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = DocxPt(16)
                run = p.add_run("▸ " + val)
                run.font.size, run.font.italic = DocxPt(9.5), True
                run.font.color.rgb = DocxRGB(*DOC_COLORS["quote"])
                _docx_shade(p, DOC_COLORS["quotebg"])
                _docx_left_border(p, "6366F1", 20)
            elif kind == "item":
                num, text = val
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = DocxPt(14)
                nr = p.add_run(f"{num}.  ")
                nr.font.size, nr.font.bold = DocxPt(11), True
                nr.font.color.rgb = DocxRGB(*DOC_COLORS["num"])
                tr = p.add_run(text)
                tr.font.size = DocxPt(10.5)
                tr.font.color.rgb = DocxRGB(*DOC_COLORS["body"])
            else:
                p = doc.add_paragraph()
                run = p.add_run(val)
                run.font.size = DocxPt(10.5)
                run.font.color.rgb = DocxRGB(*DOC_COLORS["body"])

        bio = io.BytesIO()
        doc.save(bio)
        return bio.getvalue()
    except Exception:
        return (content or "").encode("utf-8")


# ---------------------------- PDF ----------------------------
@st.cache_data(show_spinner=False, max_entries=48)
def create_pdf_bytes(title: str, content: str) -> bytes:
    try:
        f = init_korean_font()
        bio = io.BytesIO()
        doc = SimpleDocTemplate(bio, pagesize=letter, rightMargin=38, leftMargin=38,
                                topMargin=38, bottomMargin=38)

        S = dict(
            title=ParagraphStyle("T", fontName=f, fontSize=16, leading=22, textColor="#4C1D95",
                                 backColor="#F5F3FF", borderColor="#7C3AED", borderWidth=0,
                                 leftIndent=8, borderPadding=8, spaceAfter=4),
            meta=ParagraphStyle("M", fontName=f, fontSize=8, leading=12, textColor="#64748B",
                                spaceAfter=12),
            head=ParagraphStyle("H", fontName=f, fontSize=12.5, leading=19, textColor="#5B21B6",
                                backColor="#EDE9FE", borderPadding=(6, 8, 6, 10),
                                spaceBefore=10, spaceAfter=6),
            leader=ParagraphStyle("L", fontName=f, fontSize=10, leading=16, textColor="#075985",
                                  backColor="#E0F2FE", borderColor="#0284C7", borderWidth=1,
                                  borderPadding=(7, 9, 7, 11), leftIndent=2,
                                  spaceBefore=4, spaceAfter=6),
            quote=ParagraphStyle("Q", fontName=f, fontSize=9, leading=14, textColor="#4338CA",
                                 backColor="#EEF2FF", borderPadding=(5, 7, 5, 9),
                                 leftIndent=14, spaceAfter=4),
            item=ParagraphStyle("I", fontName=f, fontSize=10, leading=16, textColor="#1E293B",
                                leftIndent=16, spaceAfter=3),
            body=ParagraphStyle("B", fontName=f, fontSize=10, leading=16, textColor="#1E293B",
                                spaceAfter=4),
        )

        def esc(t):
            return (t or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

        story = [Paragraph(f"<b>{esc(title)}</b>", S["title"]),
                 Paragraph(f"작성일 {datetime.now().strftime('%Y-%m-%d')} · MY 설교 AI 스튜디오", S["meta"])]

        for kind, val in parse_doc_blocks(content):
            if kind == "blank":
                story.append(Spacer(1, 5))
            elif kind == "head":
                story.append(Paragraph(f"<b>{esc(val)}</b>", S["head"]))
            elif kind == "leader":
                story.append(Paragraph(
                    f'<font color="#0284C7"><b>💡 인도자 가이드</b></font><br/>'
                    f'<b>{esc(val)}</b>', S["leader"]))
            elif kind == "quote":
                story.append(Paragraph(f"<i>▸ {esc(val)}</i>", S["quote"]))
            elif kind == "item":
                num, text = val
                story.append(Paragraph(
                    f'<font color="#DB2777"><b>{esc(num)}.</b></font>&nbsp;&nbsp;{esc(text)}',
                    S["item"]))
            else:
                story.append(Paragraph(esc(val), S["body"]))

        doc.build(story)
        return bio.getvalue()
    except Exception:
        return (content or "").encode("utf-8")


# ---------------------------- TXT ----------------------------
def create_txt_bytes(title: str, content: str) -> bytes:
    """색은 못 넣지만, 기호로 인도자 가이드가 한눈에 구분되게 한다."""
    lines = ["═" * 58, f"  {title}",
             f"  작성일 {datetime.now().strftime('%Y-%m-%d')} · MY 설교 AI 스튜디오",
             "═" * 58, ""]
    circled = "①②③④⑤⑥⑦⑧⑨⑩⑪⑫"
    for kind, val in parse_doc_blocks(content):
        if kind == "blank":
            lines.append("")
        elif kind == "head":
            lines += ["", "━━━ " + val + " " + "━" * max(3, 46 - len(val)), ""]
        elif kind == "leader":
            lines += ["┃ 💡 인도자 가이드", f"┃ {val}", "┗" + "━" * 40]
        elif kind == "quote":
            lines.append(f"    ▸ {val}")
        elif kind == "item":
            num, text = val
            try:
                mark = circled[int(num) - 1]
            except Exception:
                mark = f"{num}."
            lines.append(f"   {mark} {text}")
        else:
            lines.append(val)
    return "\n".join(lines).encode("utf-8")


# ==============================================================================
# ★ 설교 PPT 디자인 시스템 (업로드해 주신 「믿음의 경주를 달려가라」 템플릿 기준)
#   10 × 5.625 inch(16:9) · 밝은 아이보리 배경 · 네이비 세리프 제목 · 블루 라벨
# ==============================================================================
TPL = {
    "W": 10.0, "H": 5.625,
    "M": 0.71, "CW": 8.58,          # 좌우 여백 / 본문 폭
    "bg":      "FBFCFD",            # 배경
    "white":   "FFFFFF",
    "ink":     "1B2B3D",            # 제목 네이비
    "body":    "455567",            # 본문 회색
    "muted":   "8A96A3",            # 보조 텍스트
    "accent":  "2B527A",            # 포인트 블루 (라벨·번호)
    "accent2": "7B9CC0",            # 연블루 (숫자·구분선)
    "line":    "DCE2EB",            # 구분선
    "card":    "F5F6F9",            # 연회색 카드
    "gold":    "C4A26B",            # 기도 금색
    "cream":   "F1E2C6",            # 기도 본문
    "sand":    "EEE8DA",
    "serif":   "Noto Serif KR",     # 제목용 세리프
    "sans":    "Pretendard",        # 라벨·본문용 산세리프
}


def _c(hex_str):
    return RGBColor.from_string(hex_str)


def _set_font(run, name):
    """라틴·한글(EastAsia) 글꼴을 함께 지정해야 한글에도 서체가 먹는다."""
    try:
        run.font.name = name
        rPr = run.font._rPr
        for tag in ("a:ea", "a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = etree.SubElement(rPr, qn(tag))
            el.set("typeface", name)
    except Exception:
        pass


def _letter_spacing(run, hundredths_pt: int):
    try:
        run.font._rPr.set("spc", str(int(hundredths_pt)))
    except Exception:
        pass


def _slide(prs, bg=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                              Inches(TPL["W"]), Inches(TPL["H"]))
    rect.fill.solid()
    rect.fill.fore_color.rgb = _c(bg or TPL["bg"])
    rect.line.fill.background()
    rect.shadow.inherit = False
    return s


def _bar(s, x, y, w, h, color):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid()
    r.fill.fore_color.rgb = _c(color)
    r.line.fill.background()
    r.shadow.inherit = False
    return r


def _oval(s, x, y, d, color):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    o.fill.solid()
    o.fill.fore_color.rgb = _c(color)
    o.line.fill.background()
    o.shadow.inherit = False
    return o


def _tb(s, x, y, w, h, lines, size=16, bold=False, color=None, font=None,
        align=PP_ALIGN.LEFT, line_spacing=1.35, space_after=6, spc=None):
    """lines: 문자열 또는 문자열 리스트(문단별)"""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = [lines]
    lines = [l for l in lines if str(l).strip() != ""] or [""]
    color = color or TPL["body"]
    font = font or TPL["sans"]
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = str(txt)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _c(color)
        _set_font(run, font)
        if spc:
            _letter_spacing(run, spc)
        p.alignment = align
        try:
            p.line_spacing = line_spacing
            p.space_after = Pt(space_after)
        except Exception:
            pass
    return box


def _eyebrow(s, label):
    """상단 작은 라벨 + 짧은 파란 밑줄"""
    _tb(s, TPL["M"], 0.55, TPL["CW"], 0.28, label, size=11, bold=True,
        color=TPL["accent"], font=TPL["sans"], space_after=0)
    _bar(s, TPL["M"], 0.87, 0.44, 0.02, TPL["accent"])


def _headline(s, text, size=32, y=1.15, h=0.98):
    _tb(s, TPL["M"], y, TPL["CW"], h, text, size=size, bold=True,
        color=TPL["ink"], font=TPL["serif"], line_spacing=1.25, space_after=0)


def _photo_bg(s, seed_key, theme, alpha=0.72):
    """사진 배경 + 흰색 반투명 오버레이 (템플릿의 표지·핵심말씀 스타일)"""
    try:
        b = get_background_bytes(seed_key, theme, size=(1280, 720))
        if b:
            s.shapes.add_picture(io.BytesIO(b), 0, 0,
                                 width=Inches(TPL["W"]), height=Inches(TPL["H"]))
    except Exception:
        pass
    ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(TPL["W"]), Inches(TPL["H"]))
    ov.fill.solid()
    ov.fill.fore_color.rgb = _c(TPL["white"])
    set_shape_fill_alpha(ov, alpha)
    ov.line.fill.background()
    ov.shadow.inherit = False


def _wrap_for_slide(text, per_line=46, max_lines=9):
    """슬라이드에 넣기 좋게 문단을 자른다."""
    t = re.sub(r'\s+', ' ', str(text or "")).strip()
    limit = per_line * max_lines
    return t[:limit] + ("…" if len(t) > limit else "")


# ---------- 슬라이드 종류별 빌더 ----------
def sl_cover(prs, title, subtitle, scripture, seed, theme, kicker="주 일 설 교"):
    s = _slide(prs, TPL["bg"])
    _photo_bg(s, f"{seed}|cover", theme, 0.70)
    _tb(s, 0, 1.20, TPL["W"], 0.33, kicker, size=12, bold=True, color=TPL["accent"],
        font=TPL["sans"], align=PP_ALIGN.CENTER, space_after=0, spc=400)
    _tb(s, TPL["M"], 1.86, TPL["CW"], 1.31, _wrap_for_slide(title, 22, 2),
        size=42 if len(str(title)) <= 18 else 34, bold=True, color=TPL["ink"],
        font=TPL["serif"], align=PP_ALIGN.CENTER, line_spacing=1.25, space_after=0)
    if subtitle:
        _tb(s, TPL["M"], 3.28, TPL["CW"], 0.62, _wrap_for_slide(subtitle, 32, 2),
            size=17, color=TPL["body"], font=TPL["serif"], align=PP_ALIGN.CENTER,
            line_spacing=1.35, space_after=0)
    _bar(s, 4.67, 3.99, 0.66, 0.012, TPL["accent2"])
    _tb(s, TPL["M"], 4.21, TPL["CW"], 0.44, f"본문 · {scripture}", size=13,
        color=TPL["body"], font=TPL["sans"], align=PP_ALIGN.CENTER, space_after=0)
    return s


def sl_text(prs, eyebrow, headline, body, head_size=34):
    s = _slide(prs)
    _eyebrow(s, eyebrow)
    _headline(s, headline, size=head_size, y=1.20)
    paras = [p for p in str(body or "").split("\n") if p.strip()][:4]
    _tb(s, TPL["M"], 2.30, TPL["CW"], 3.06,
        [_wrap_for_slide(p, 46, 4) for p in paras] or [""],
        size=17, color=TPL["body"], font=TPL["sans"], line_spacing=1.5, space_after=10)
    return s


def sl_outline(prs, points, eyebrow="오늘의 말씀", headline="설교의 흐름"):
    s = _slide(prs)
    _eyebrow(s, eyebrow)
    _headline(s, headline, size=32, y=1.18, h=0.79)
    y = 2.35
    for i, t in enumerate(points[:4], start=1):
        _tb(s, TPL["M"], y, 0.87, 0.62, f"{i:02d}", size=26, bold=True,
            color=TPL["accent2"], font=TPL["serif"], space_after=0)
        _tb(s, 1.75, y + 0.05, 7.54, 0.62, _wrap_for_slide(t, 34, 1), size=18, bold=True,
            color=TPL["ink"], font=TPL["serif"], space_after=0)
        if i < min(4, len(points)):
            _bar(s, TPL["M"], y + 0.64, TPL["CW"], 0.012, TPL["line"])
        y += 0.69
    return s


def sl_scripture(prs, scripture, verses):
    s = _slide(prs)
    _eyebrow(s, f"본문 말씀  ·  {scripture}")
    lines = [v for v in verses if str(v).strip()][:5]
    _tb(s, TPL["M"], 1.48, TPL["CW"], 3.49,
        [_wrap_for_slide(v, 44, 3) for v in lines] or [scripture],
        size=16, color=TPL["ink"], font=TPL["serif"], line_spacing=1.6, space_after=12)
    return s


def sl_point(prs, order, label, headline, quote, quote_ref, body, cards=None):
    s = _slide(prs)
    _eyebrow(s, f"{order:02d}   ·   {label}")
    _headline(s, _wrap_for_slide(headline, 24, 2), size=30, y=1.15)

    top = 2.41
    if quote:
        _bar(s, TPL["M"], top, 0.014, 0.77, TPL["accent2"])
        _tb(s, 0.93, top, 8.36, 0.49, f"“{_wrap_for_slide(quote, 44, 1)}”",
            size=15, color=TPL["body"], font=TPL["serif"], space_after=0)
        _tb(s, 0.93, top + 0.52, 8.36, 0.27, f"— {quote_ref}" if quote_ref else "",
            size=11, color=TPL["muted"], font=TPL["sans"], space_after=0)
        top = 3.34
    else:
        top = 2.50

    cards = [c for c in (cards or []) if c and str(c[0]).strip()]
    if len(cards) >= 2:
        n = min(3, len(cards))
        gap = 0.22
        cw = (TPL["CW"] - gap * (n - 1)) / n
        for i, (ct, cd) in enumerate(cards[:n]):
            x = TPL["M"] + i * (cw + gap)
            dark = (i == 0)
            box = _bar(s, x, top, cw, 1.80, TPL["accent"] if dark else TPL["white"])
            if not dark:
                box.line.color.rgb = _c(TPL["line"])
                box.line.width = Pt(0.75)
            _bar(s, x, top, cw, 0.035, TPL["white"] if dark else TPL["accent"])
            _tb(s, x + 0.27, top + 0.72, cw - 0.54, 0.46, _wrap_for_slide(ct, 14, 1),
                size=17, bold=True, color=TPL["white"] if dark else TPL["accent"],
                font=TPL["serif"], space_after=0)
            _tb(s, x + 0.27, top + 1.20, cw - 0.54, 0.50, _wrap_for_slide(cd, 20, 2),
                size=10.5, color="D8DEE7" if dark else TPL["body"],
                font=TPL["sans"], line_spacing=1.35, space_after=0)
    else:
        paras = [p for p in str(body or "").split("\n") if p.strip()][:3]
        _tb(s, TPL["M"], top, TPL["CW"], 5.30 - top,
            [_wrap_for_slide(p, 46, 3) for p in paras] or [""],
            size=16, color=TPL["body"], font=TPL["sans"], line_spacing=1.5, space_after=10)
    return s


def sl_keyverse(prs, quote, ref, caption, seed, theme):
    s = _slide(prs, TPL["bg"])
    _photo_bg(s, f"{seed}|key", theme, 0.74)
    _tb(s, 0, 0.87, TPL["W"], 0.33, "핵  심  말  씀", size=12, bold=True,
        color=TPL["accent"], font=TPL["sans"], align=PP_ALIGN.CENTER, space_after=0, spc=400)
    _tb(s, TPL["M"], 1.75, TPL["CW"], 2.19, f"“{_wrap_for_slide(quote, 30, 3)}”",
        size=32 if len(str(quote)) <= 46 else 26, bold=True, color=TPL["ink"],
        font=TPL["serif"], align=PP_ALIGN.CENTER, line_spacing=1.45, space_after=0)
    _bar(s, 4.67, 4.16, 0.66, 0.012, TPL["accent2"])
    _tb(s, TPL["M"], 4.37, TPL["CW"], 0.38, ref, size=16, color=TPL["ink"],
        font=TPL["serif"], align=PP_ALIGN.CENTER, space_after=0)
    if caption:
        _tb(s, TPL["M"], 4.95, TPL["CW"], 0.33, _wrap_for_slide(caption, 44, 1),
            size=11, color=TPL["muted"], font=TPL["sans"], align=PP_ALIGN.CENTER, space_after=0)
    return s


def sl_numbered(prs, eyebrow, headline, items):
    s = _slide(prs)
    _eyebrow(s, eyebrow)
    _headline(s, headline, size=32, y=1.18, h=0.79)
    y = 2.32
    for i, t in enumerate(items[:3], start=1):
        _bar(s, TPL["M"], y, TPL["CW"], 0.75, TPL["card"])
        _tb(s, 0.84, y + 0.10, 0.77, 0.55, str(i), size=24, bold=True,
            color=TPL["accent"], font=TPL["serif"], align=PP_ALIGN.CENTER, space_after=0)
        _tb(s, 1.72, y + 0.10, 7.40, 0.60, _wrap_for_slide(t, 52, 2), size=14,
            color=TPL["body"], font=TPL["sans"], line_spacing=1.4, space_after=0)
        y += 0.85
    return s


def sl_prayer(prs, prayer_text, heading="기  도"):
    s = _slide(prs, TPL["ink"])
    _tb(s, TPL["M"], 1.31, 2.41, 0.55, heading, size=30, bold=True,
        color=TPL["gold"], font=TPL["serif"], space_after=0, spc=300)
    _bar(s, TPL["M"], 2.05, 0.44, 0.02, TPL["gold"])
    lines = [l.strip() for l in re.split(r'(?<=[.!?])\s+|\n', str(prayer_text or "")) if l.strip()][:8]
    _tb(s, 3.12, 0.61, 6.34, 4.47, [_wrap_for_slide(l, 40, 3) for l in lines] or ["아멘."],
        size=14, color=TPL["cream"], font=TPL["serif"], line_spacing=1.85, space_after=8)
    return s


def _new_deck():
    prs = Presentation()
    prs.slide_width = Inches(TPL["W"])
    prs.slide_height = Inches(TPL["H"])
    return prs


# ---------------------------- PPT (문서형) ----------------------------
@st.cache_data(show_spinner=False, max_entries=32)
def create_document_pptx_bytes(title: str, content: str, seed_base: str = "",
                               bg_theme: str = "자연 · 풍경") -> bytes:
    """
    일반 문서형 PPT (QT·나눔지·가이드·칼럼 등) — 설교 PPT와 같은 디자인 언어.
    표지 → 내용 슬라이드(라벨 + 세리프 제목 + 본문/카드/인용) 구조.
    """
    try:
        prs = _new_deck()
        blocks = [b for b in parse_doc_blocks(content) if b[0] != "blank"]

        # 표지
        sl_cover(prs, title, "", datetime.now().strftime("%Y-%m-%d"),
                 seed_base or title, bg_theme, kicker="사 역 자 료")

        # 섹션(head) 단위로 슬라이드를 나눈다
        sections, cur_head, cur = [], "", []
        for kind, val in blocks:
            if kind == "head":
                if cur or cur_head:
                    sections.append((cur_head, cur))
                cur_head, cur = val, []
            else:
                cur.append((kind, val))
        if cur or cur_head:
            sections.append((cur_head, cur))
        if not sections:
            sections = [("", blocks)]

        def weight(kind, val):
            txt = val[1] if kind == "item" else str(val)
            return max(1, len(txt) // 42 + 1) + (1 if kind == "leader" else 0)

        sections = [(h, b) for h, b in sections if b] or sections
        idx = 0
        for head, body in sections:
            pages, cur_p, used = [], [], 0
            for kind, val in body:
                w = weight(kind, val)
                if used + w > 8 and cur_p:
                    pages.append(cur_p)
                    cur_p, used = [], 0
                cur_p.append((kind, val))
                used += w
            if cur_p or not pages:
                pages.append(cur_p)

            for pi, pg in enumerate(pages):
                idx += 1
                s2 = _slide(prs)
                raw_head = strip_emoji(re.sub(r'^\s*\d+[\.\)]\s*', '', head)) or title
                raw_head = re.split(r'\s*·\s*', raw_head)[0].strip() or title
                label = _wrap_for_slide(raw_head, 26, 1)
                _eyebrow(s2, f"{idx:02d}   ·   {label}")
                _headline(s2, _wrap_for_slide(label, 20, 2)
                          + (f"  ({pi+1})" if len(pages) > 1 else ""),
                          size=30 if len(label) <= 16 else 25, y=1.12, h=0.9)

                y = 2.22
                for kind, val in pg:
                    if kind == "leader":
                        h = 0.30 + 0.22 * max(1, len(str(val)) // 50 + 1)
                        card = _bar(s2, TPL["M"], y, TPL["CW"], h, "EAF4FB")
                        card.line.color.rgb = _c("9FC7E3")
                        card.line.width = Pt(0.75)
                        _bar(s2, TPL["M"], y, 0.05, h, TPL["accent"])
                        _tb(s2, TPL["M"] + 0.22, y + 0.09, TPL["CW"] - 0.44, 0.22,
                            "인도자 가이드", size=10.5, bold=True, color=TPL["accent"],
                            font=TPL["sans"], space_after=0)
                        _tb(s2, TPL["M"] + 0.22, y + 0.33, TPL["CW"] - 0.44, h - 0.42,
                            _wrap_for_slide(val, 54, 3), size=13, bold=True,
                            color="1D4F73", font=TPL["sans"], line_spacing=1.4, space_after=0)
                        y += h + 0.14
                    elif kind == "quote":
                        _bar(s2, TPL["M"], y, 0.014, 0.42, TPL["accent2"])
                        _tb(s2, TPL["M"] + 0.22, y, TPL["CW"] - 0.22, 0.42,
                            _wrap_for_slide(val, 54, 2), size=12.5, color=TPL["muted"],
                            font=TPL["serif"], line_spacing=1.35, space_after=0)
                        y += 0.52
                    elif kind == "item":
                        num, text = val
                        h = 0.42 + 0.18 * (len(text) // 52)
                        _bar(s2, TPL["M"], y, TPL["CW"], h, TPL["card"])
                        _tb(s2, 0.84, y + 0.06, 0.55, 0.34, str(num), size=16, bold=True,
                            color=TPL["accent"], font=TPL["serif"], align=PP_ALIGN.CENTER,
                            space_after=0)
                        _tb(s2, 1.50, y + 0.07, 7.60, h - 0.14, _wrap_for_slide(text, 54, 2),
                            size=13.5, color=TPL["body"], font=TPL["sans"],
                            line_spacing=1.4, space_after=0)
                        y += h + 0.10
                    else:
                        h = 0.30 + 0.22 * (len(str(val)) // 50)
                        _tb(s2, TPL["M"], y, TPL["CW"], h + 0.2, _wrap_for_slide(val, 52, 3),
                            size=14, color=TPL["body"], font=TPL["sans"],
                            line_spacing=1.5, space_after=0)
                        y += h + 0.14
                    if y > 5.05:
                        break

        bio = io.BytesIO()
        prs.save(bio)
        return bio.getvalue()
    except Exception:
        return (content or "").encode("utf-8")


# ------------------------------------------------------------------------------
# 설교 구조 PPT — 이제 요약 텍스트를 '진짜로' 파싱해서 슬라이드에 넣는다
# ------------------------------------------------------------------------------
def parse_sermon_content(title, scripture, summary_content, full_sermon=""):
    """요약문에서 명제/대지/적용/기도를 구조적으로 뽑아낸다. 실패 시 원고 분석으로 대체."""
    text = summary_content if summary_content and len(summary_content) > 40 else full_sermon
    text = text or ""

    # 명제
    prop = ""
    m = re.search(r'🎯[^\n]*\n(.+?)(?=\n\s*(?:🔑|📖|📌|💡|🙏|$))', text, re.DOTALL)
    if m:
        prop = m.group(1).strip()
    if not prop:
        m = re.search(r'(?:핵심\s*명제|중심\s*사상)[:：]?\s*(.+?)(?=\n\s*\n|\n[🔑📖📌💡🙏]|$)', text, re.DOTALL)
        prop = m.group(1).strip() if m else ""

    # 대지: "1. ..." ~ "3. ..." 블록 (📌 섹션 우선)
    points = []
    seg = text
    ms = re.search(r'📌(.+?)(?=\n\s*💡|\n\s*🙏|$)', text, re.DOTALL)
    if ms:
        seg = ms.group(1)
    for pm in re.finditer(r'(?m)^\s*([1-4])[\.\)]\s*(.+?)(?=\n\s*[1-4][\.\)]\s|\Z)', seg, re.DOTALL):
        block = pm.group(2).strip()
        if len(block) > 4:
            points.append(block)

    a = analyze_manuscript(full_sermon)
    if len(points) < 3:
        points = []
        for sec in a["sections"][:3]:
            points.append(f"{_trim_title(sec['headline'])}\n▸ 원고 근거: {sec['evidence'][:200]}")

    if not prop:
        prop = a["key_sentences"][0] if a["key_sentences"] else title

    # 적용
    app_text = ""
    m = re.search(r'💡[^\n]*\n(.+?)(?=\n\s*🙏|\Z)', text, re.DOTALL)
    if m:
        app_text = m.group(1).strip()
    if not app_text:
        if a["applications"]:
            app_text = "\n".join(f"{i}. {s}" for i, s in enumerate(a["applications"][:3], start=1))
        else:
            app_text = f"{scripture} 말씀을 이번 한 주 삶의 자리에서 한 가지 행동으로 옮기십시오."

    # 기도
    prayer_text = ""
    m = re.search(r'🙏[^\n]*\n(.+)$', text, re.DOTALL)
    if m:
        prayer_text = m.group(1).strip()
    if not prayer_text:
        prayer_text = a["prayer"] or (
            f"주님, 오늘 [{title}] 말씀({scripture})을 마음에 새기고 순종하게 하옵소서. "
            f"예수님의 이름으로 기도드립니다. 아멘.")

    # 대지가 하나도 없을 때만 원고에서 보충한다 (적용문을 대지로 잘못 올리지 않도록)
    if not points:
        for sec in a["sections"][:3]:
            points.append(f"{_trim_title(sec['headline'])}\n\u25b8 원고 근거: {sec['evidence'][:200]}")
    if not points:
        points = [prop]

    # 본문 말씀 인용 (📜 섹션)
    verses = []
    mv = re.search(r'📜[^\n]*\n(.+?)(?=\n\s*[🎯🔑📌💡🙏]|\Z)', text, re.DOTALL)
    if mv:
        for ln in mv.group(1).split("\n"):
            ln = ln.strip()
            if len(ln) > 5:
                verses.append(ln)
    if not verses:
        for sent in a["sentences"]:
            if ('"' in sent or '“' in sent) and SCRIPTURE_REF_RE.search(sent):
                verses.append(sent)
            if len(verses) >= 3:
                break

    return {"prop": prop, "points": points, "app": app_text,
            "prayer": prayer_text, "verses": verses[:5]}


def _split_point(raw: str):
    """대지 한 덩어리에서 제목 / 인용구 / 출처 / 본문 / 카드를 뽑아낸다."""
    lines = [l.strip() for l in str(raw or "").split("\n") if l.strip()]
    if not lines:
        return "", "", "", "", []
    head = lines[0]
    ref = ""
    m = re.search(r'\[\s*근거\s*성구\s*[:：]\s*([^\]]+)\]', head)
    if m:
        ref = m.group(1).strip()
        head = head[:m.start()].strip()
    head = re.sub(r'^\s*\d+[\.\)]\s*', '', head).strip(' .')

    quote, body_parts, cards = "", [], []
    for l in lines[1:]:
        t = re.sub(r'^▸\s*', '', l)
        if t.startswith("원고 근거"):
            q = t.split(":", 1)[-1].strip().strip('"“”')
            if not quote:
                quote = q
        elif t.startswith("주해") or t.startswith("설명"):
            body_parts.append(t.split(":", 1)[-1].strip())
        else:
            mm = re.match(r'^(.{2,16}?)\s*[—\-–:]\s*(.+)$', t)
            if mm and len(t) < 90:
                cards.append((mm.group(1).strip(), mm.group(2).strip()))
            else:
                body_parts.append(t)

    if not quote and body_parts:
        quote = body_parts[0][:70]
    body = "\n".join(body_parts) or head
    return head, quote, ref, body, cards


@st.cache_data(show_spinner=False, max_entries=24)
def generate_sermon_structure_pptx_bytes(title: str, scripture: str, summary_content: str,
                                         full_sermon: str = "", seed_base: str = "",
                                         bg_theme: str = "자연 · 풍경") -> bytes:
    """
    설교 요약 PPT — 업로드해 주신 템플릿과 같은 구성.
    표지 → 들어가며 → 설교의 흐름 → 본문 말씀 → 대지 1~3 → 핵심 말씀 → 삶의 적용 → 나가며 → 기도
    """
    try:
        prs = _new_deck()
        parsed = parse_sermon_content(title, scripture, summary_content, full_sermon)
        a = analyze_manuscript(full_sermon)
        points = parsed["points"]
        seed = seed_base or f"{title}|{scripture}"

        heads, blocks = [], []
        for p_raw in points[:4]:
            h, q, r, b, cs = _split_point(p_raw)
            heads.append(h or title)
            blocks.append((h or title, q, r or scripture, b, cs))

        # 1. 표지
        sub = re.split(r'(?<=[.!?])\s+', str(parsed.get("prop") or ""))[0]
        sl_cover(prs, title, sub[:56], scripture, seed, bg_theme)

        # 2. 들어가며
        intro = (a["paragraphs"][0] if a["paragraphs"] else parsed["prop"]) or parsed["prop"]
        sl_text(prs, "들어가며", "말씀 앞에 서며", intro, head_size=32)

        # 3. 설교의 흐름
        sl_outline(prs, heads or [title])

        # 4. 본문 말씀
        verses = parsed.get("verses") or []
        if not verses:
            verses = [s2 for s2 in a["key_sentences"][:3] if s2] or [parsed["prop"]]
        sl_scripture(prs, scripture, verses)

        # 5~7. 대지
        labels = ["첫 번째 대지", "두 번째 대지", "세 번째 대지", "네 번째 대지"]
        for i, (h, q, r, b, cs) in enumerate(blocks[:4]):
            sl_point(prs, i + 1, labels[i], h, q, r, b, cs)

        # 8. 핵심 말씀
        vlist = [re.sub(r'^\s*\d+\s+', '', v).strip() for v in (parsed.get("verses") or [])]
        vlist = [v for v in vlist if v]
        keyq = min(vlist, key=lambda v: abs(len(v) - 45)) if vlist else parsed["prop"]
        if len(keyq) > 90:
            keyq = re.split(r'(?<=[.!?])\s+|,\s+', keyq)[0][:90]
        sl_keyverse(prs, _wrap_for_slide(keyq, 30, 3), scripture,
                    (a["applications"][0] if a["applications"] else ""), seed, bg_theme)

        # 9. 삶의 적용
        apps = [re.sub(r'^\s*[-•]?\s*\d+[\.\)]\s*', '', x).strip()
                for x in str(parsed["app"]).split("\n") if x.strip()]
        sl_numbered(prs, "삶의 적용", "이렇게 살아갑시다", apps or [parsed["prop"]])

        # 10. 나가며
        outro = ""
        for para in reversed(a["paragraphs"][-4:] if a["paragraphs"] else []):
            if not PRAYER_HINT_RE.search(para):
                outro = para
                break
        sl_text(prs, "나가며", "말씀을 품고 나아가며", outro or parsed["prop"], head_size=32)

        # 11. 기도
        sl_prayer(prs, parsed["prayer"])

        bio = io.BytesIO()
        prs.save(bio)
        return bio.getvalue()
    except Exception:
        return create_document_pptx_bytes(title, summary_content)


@st.cache_data(show_spinner=False, max_entries=16)
def generate_cardnews_pptx_bytes(cards_json: str, church_name: str = "",
                                 seed_base: str = "", bg_theme: str = "자연 · 풍경",
                                 scripture: str = "") -> bytes:
    """카드뉴스 PPT — 설교 PPT와 같은 디자인 언어(밝은 배경 · 네이비 세리프 · 블루 라벨)"""
    cards = json.loads(cards_json)
    prs = Presentation()
    prs.slide_width = prs.slide_height = Inches(10)          # 1:1 정사각형
    blank = prs.slide_layouts[6]

    for idx, item in enumerate(cards):
        s = prs.slides.add_slide(blank)
        base = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(10))
        base.fill.solid()
        base.fill.fore_color.rgb = _c(TPL["bg"])
        base.line.fill.background()
        base.shadow.inherit = False

        b = get_background_bytes(f"{seed_base}|{idx}", bg_theme)
        if b:
            s.shapes.add_picture(io.BytesIO(b), 0, 0, width=Inches(10), height=Inches(10))
        ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(10))
        ov.fill.solid()
        ov.fill.fore_color.rgb = _c(TPL["white"])
        set_shape_fill_alpha(ov, 0.80)
        ov.line.fill.background()
        ov.shadow.inherit = False

        head = strip_emoji(str(item.get("headline", "")))
        body = strip_emoji(str(item.get("body_text", "")))

        _tb(s, 1.0, 1.0, 8.0, 0.35, f"CARD {item.get('card_number', idx + 1):02d}",
            size=13, bold=True, color=TPL["accent"], font=TPL["sans"], space_after=0, spc=400)
        _bar(s, 1.0, 1.48, 0.5, 0.02, TPL["accent"])

        h_size = 36 if len(head) <= 16 else (30 if len(head) <= 26 else 25)
        h_lines = min(3, max(1, len(head) // 18 + 1))
        paras = [p for p in body.split("\n") if p.strip()][:4]
        b_lines = sum(min(3, max(1, len(p) // 28 + 1)) for p in paras) or 1

        h_h = h_lines * (h_size / 72 * 1.35)
        b_h = b_lines * (18 / 72 * 1.75)
        gap = 0.62
        top = 2.15 + max(0.0, ((7.55 - 2.15) - (h_h + gap + b_h)) / 2)

        _tb(s, 1.0, top, 8.0, h_h + 0.3, _wrap_for_slide(head, 18, 3),
            size=h_size, bold=True, color=TPL["ink"], font=TPL["serif"],
            line_spacing=1.35, space_after=0)
        _tb(s, 1.0, top + h_h + gap, 8.0, b_h + 0.4,
            [_wrap_for_slide(p, 28, 3) for p in paras] or [""],
            size=18, color=TPL["body"], font=TPL["sans"], line_spacing=1.75, space_after=14)

        _bar(s, 1.0, 8.10, 0.66, 0.012, TPL["accent2"])
        if scripture:
            _tb(s, 1.0, 8.32, 8.0, 0.4, f"「 {scripture} 」", size=15, bold=True,
                color=TPL["accent"], font=TPL["serif"], space_after=0)
        if church_name:
            _tb(s, 1.0, 8.80, 8.0, 0.4, church_name, size=13, color=TPL["muted"],
                font=TPL["sans"], space_after=0)

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


# ==============================================================================
# 말씀카드
# ==============================================================================
# ==============================================================================
# 오늘의 만나 / 오늘의 말씀 — 이미지 렌더러
# ==============================================================================
SERIF_FONT_PATHS = [
    "/usr/share/fonts/opentype/noto/NotoSerifCJK-Regular.ttc",
    "/usr/share/fonts/opentype/noto/NotoSerifCJK-Medium.ttc",
    "/usr/share/fonts/truetype/nanum/NanumMyeongjo.ttf",
    "C:/Windows/Fonts/batang.ttc",
    "/System/Library/Fonts/AppleMyungjo.ttf",
]


@st.cache_resource(show_spinner=False)
def get_serif_font(size: int):
    for p in SERIF_FONT_PATHS:
        if os.path.exists(p):
            try:
                return PIL.ImageFont.truetype(p, size)
            except Exception:
                continue
    return get_pil_font(size)


MONTH_EN = ["January", "February", "March", "April", "May", "June",
            "July", "August", "September", "October", "November", "December"]


@st.cache_data(show_spinner=False, max_entries=48)
def render_manna_png(date_iso: str, ko_verse: str, ko_ref: str, en_verse: str, en_ref: str,
                     story_title: str, story_body: str, seed: str = "",
                     theme: str = "하늘 · 빛", org_name: str = "") -> bytes:
    """
    「365일 성경묵상」 지면과 같은 구성의 '오늘의 만나' 이미지.
    좌측 큰 박스: 개역개정 + NIV / 우측 패널: 예화 / 상단: 월 이름 + 월/일
    """
    W = 1600
    d0 = datetime.strptime(date_iso, "%Y-%m-%d")

    # ── 먼저 빈 캔버스에 재보고 필요한 높이만큼만 그림을 만든다
    _m = PIL.ImageDraw.Draw(PIL.Image.new("RGB", (10, 10)))
    px, pw = W - 430, 360
    f_st = get_pil_font(27)
    f_sb = get_pil_font(21)
    story_t = wrap_korean_text(strip_emoji(story_title or ""), f_st, pw - 70, _m)
    story_b = wrap_korean_text(strip_emoji(story_body or ""), f_sb, pw - 56, _m)
    tb_h = _m.multiline_textbbox((0, 0), story_t, font=f_st, spacing=8)[3]
    bb_h = _m.multiline_textbbox((0, 0), story_b, font=f_sb, spacing=13)[3]
    panel_h = 60 + tb_h + 26 + bb_h + 40

    bx0, bx1 = 60, px - 60
    inner_w = bx1 - bx0 - 110
    ko = str(ko_verse or "").strip()
    size = 46 if len(ko) <= 60 else (40 if len(ko) <= 100 else 34)
    f_ko = get_serif_font(size)
    f_ref = get_serif_font(26)
    f_en = get_serif_font(27)
    f_enref = get_serif_font(23)
    ko_wrapped = wrap_korean_text(ko, f_ko, inner_w, _m)
    en_wrapped = wrap_korean_text(str(en_verse or "").strip(), f_en, inner_w, _m)
    ko_h = _m.multiline_textbbox((0, 0), ko_wrapped, font=f_ko, spacing=int(size * 0.55))[3]
    en_h = _m.multiline_textbbox((0, 0), en_wrapped, font=f_en, spacing=14)[3]
    ref_h = _m.textbbox((0, 0), "고후 5:17", font=f_ref)[3]
    content_h = ko_h + 22 + ref_h + 62 + en_h + 16 + ref_h
    box_h = content_h + 130

    TOP = 250
    H = TOP + max(box_h, panel_h) + 110

    # ── 배경: 사진을 아주 옅은 푸른 톤으로 깔기
    img = PIL.Image.new("RGB", (W, H), (245, 248, 253))
    try:
        b = get_background_bytes(seed or date_iso, theme, size=(1080, 1080))
        if b:
            photo = PIL.Image.open(io.BytesIO(b)).convert("RGB").resize((W, H))
            photo = photo.filter(PIL.ImageFilter.GaussianBlur(radius=3))
            tint = PIL.Image.new("RGB", (W, H), (226, 236, 250))
            img = PIL.Image.blend(photo, tint, 0.80)
    except Exception:
        pass

    d = PIL.ImageDraw.Draw(img, "RGBA")

    NAVY = (36, 58, 96)
    BLUE = (78, 108, 168)
    SOFT = (120, 146, 194)
    PANEL = (222, 232, 247, 235)

    # ── 상단: 월 이름 + 날짜
    f_month = get_serif_font(72)
    f_day = get_serif_font(88)
    mtxt = MONTH_EN[d0.month - 1]
    mb = d.textbbox((0, 0), mtxt, font=f_month)
    d.text((W - 560 - (mb[2] - mb[0]) // 2, 66), mtxt, fill=(140, 166, 210), font=f_month)
    daytxt = f"{d0.month}/{d0.day}"
    db = d.textbbox((0, 0), daytxt, font=f_day)
    d.text((W - 90 - (db[2] - db[0]), 58), daytxt, fill=(150, 174, 214), font=f_day)
    d.line([(70, 210), (W - 70, 210)], fill=(180, 200, 228, 180), width=2)

    # ── 우측 예화 패널
    d.rounded_rectangle([px, TOP, px + pw, TOP + panel_h], radius=10, fill=PANEL)
    d.rectangle([px + 24, TOP + 38, px + 33, TOP + 47], fill=BLUE)
    d.multiline_text((px + 44, TOP + 30), story_t, fill=NAVY, font=f_st, spacing=8)
    d.multiline_text((px + 28, TOP + 30 + tb_h + 26), story_b, fill=(58, 78, 118),
                     font=f_sb, spacing=13)

    # ── 좌측 말씀 박스
    by0, by1 = TOP, TOP + box_h
    d.rectangle([bx0, by0, bx1, by1], fill=(255, 255, 255, 168), outline=BLUE, width=4)

    y = by0 + max(50, (box_h - content_h) // 2)
    d.multiline_text((bx0 + 55, y), ko_wrapped, fill=(26, 34, 52),
                     font=f_ko, spacing=int(size * 0.55))
    y += ko_h + 22
    d.text((bx0 + 55, y), f"– {ko_ref}", fill=(70, 92, 132), font=f_ref)
    y += ref_h + 62
    d.multiline_text((bx0 + 55, y), en_wrapped, fill=SOFT, font=f_en, spacing=14)
    y += en_h + 16
    d.text((bx0 + 55, y), f"– {en_ref}", fill=SOFT, font=f_enref)

    # ── 하단 날짜 · 단체명
    f_small = get_pil_font(22)
    d.text((70, H - 58), d0.strftime("%Y년 %m월 %d일") + "  ·  오늘의 만나",
           fill=(130, 152, 190), font=f_small)
    if str(org_name or "").strip():
        f_org = get_serif_font(26)
        ob = d.textbbox((0, 0), org_name, font=f_org)
        d.text((W - 70 - (ob[2] - ob[0]), H - 62), org_name,
               fill=(96, 124, 172), font=f_org)

    out = io.BytesIO()
    img.save(out, format="PNG")
    return out.getvalue()


@st.cache_data(show_spinner=False, max_entries=48)
def render_today_verse_png(date_iso: str, verse: str, ref: str, church: str = "",
                           seed: str = "", theme: str = "자연 · 풍경",
                           size_key: str = "1:1 정사각형") -> bytes:
    """오늘의 말씀 카드 — 배경 사진 + 성경 구절 + 날짜"""
    W, H = (1080, 1080) if size_key.startswith("1:1") else (1080, 1350)
    d0 = datetime.strptime(date_iso, "%Y-%m-%d")

    base = PIL.Image.new("RGBA", (W, H), (15, 23, 42, 255))
    try:
        b = get_background_bytes(seed or f"today|{date_iso}", theme, size=(1080, 1080))
        if b:
            base = PIL.Image.open(io.BytesIO(b)).convert("RGBA").resize((W, H))
    except Exception:
        pass

    overlay = PIL.Image.new("RGBA", (W, H), (7, 12, 28, 130))
    img = PIL.Image.alpha_composite(base, overlay)

    # 위·아래 스크림
    scrim = PIL.Image.new("RGBA", (W, H), (0, 0, 0, 0))
    sd = PIL.ImageDraw.Draw(scrim)
    for yy in range(H):
        if yy < int(H * 0.30):
            a = int(120 * (1 - yy / (H * 0.30)))
        elif yy > int(H * 0.72):
            a = int(140 * ((yy - H * 0.72) / (H * 0.28)))
        else:
            a = 0
        if a:
            sd.line([(0, yy), (W, yy)], fill=(4, 8, 20, a))
    img = PIL.Image.alpha_composite(img, scrim)
    d = PIL.ImageDraw.Draw(img)

    # 상단 날짜
    f_date = get_pil_font(30)
    dtxt = d0.strftime("%Y. %m. %d")
    wtxt = ["월", "화", "수", "목", "금", "토", "일"][d0.weekday()] + "요일"
    line = f"{dtxt}  {wtxt}"
    bb = d.textbbox((0, 0), line, font=f_date)
    d.text(((W - (bb[2] - bb[0])) // 2, 92), line, fill=(253, 224, 71, 255), font=f_date)
    d.line([(W // 2 - 46, 148), (W // 2 + 46, 148)], fill=(253, 224, 71, 220), width=3)

    # 본문
    txt = strip_emoji(str(verse or "").strip())
    fs = 54 if len(txt) <= 45 else (46 if len(txt) <= 80 else 38)
    f_v = get_serif_font(fs)
    wrapped = wrap_korean_text(txt, f_v, W - 190, d)
    vb = d.multiline_textbbox((0, 0), wrapped, font=f_v, spacing=int(fs * 0.62))
    vw, vh = vb[2] - vb[0], vb[3] - vb[1]
    vy = max(230, (H - vh) // 2 - 20)
    for dx in (-2, -1, 0, 1, 2):
        for dy in (-2, -1, 0, 1, 2):
            if dx or dy:
                d.multiline_text(((W - vw) // 2 + dx, vy + dy), wrapped, font=f_v,
                                 fill=(0, 0, 0, 210), align="center", spacing=int(fs * 0.62))
    d.multiline_text(((W - vw) // 2, vy), wrapped, font=f_v, fill=(255, 255, 255, 255),
                     align="center", spacing=int(fs * 0.62))

    # 출처
    f_r = get_serif_font(34)
    rtxt = f"「 {ref} 」"
    rb = d.textbbox((0, 0), rtxt, font=f_r)
    d.text(((W - (rb[2] - rb[0])) // 2, min(vy + vh + 60, H - 200)), rtxt,
           fill=(253, 224, 71, 255), font=f_r)

    if church:
        f_c = get_pil_font(26)
        cb = d.textbbox((0, 0), church, font=f_c)
        d.text(((W - (cb[2] - cb[0])) // 2, H - 96), church,
               fill=(190, 210, 245, 255), font=f_c)

    out = io.BytesIO()
    img.convert("RGB").save(out, format="PNG")
    return out.getvalue()


# ==============================================================================
# 문맥 연구용 도해(圖解) 이미지 — 본문 단락 구조 · 연대표를 그림으로
# ==============================================================================
@st.cache_data(show_spinner=False, max_entries=32)
def render_outline_diagram(scripture: str, outline_json: str, timeline_json: str = "[]") -> bytes:
    """단락 구조 + 연대표를 한 장의 도해 이미지(PNG)로 만든다. 슬라이드에 바로 붙일 수 있다."""
    try:
        outline = json.loads(outline_json) or []
    except Exception:
        outline = []
    try:
        timeline = json.loads(timeline_json) or []
    except Exception:
        timeline = []

    W = 1500
    row_h = 108
    head_h = 130
    tl_h = 240 if timeline else 0
    H = head_h + max(1, len(outline)) * row_h + tl_h + 70

    img = PIL.Image.new("RGB", (W, H), (250, 250, 255))
    d = PIL.ImageDraw.Draw(img)

    f_title = get_pil_font(38)
    f_head = get_pil_font(27)
    f_body = get_pil_font(22)
    f_small = get_pil_font(19)

    # 상단 그라데이션 띠
    for y in range(96):
        t = y / 96
        d.line([(0, y), (W, y)], fill=(int(76 + 48 * t), int(29 + 100 * t), int(149 + 80 * t)))
    d.text((44, 28), f"본문 구조 도해 · {scripture}", fill=(255, 255, 255), font=f_title)

    palette = [(124, 58, 237), (14, 165, 233), (219, 39, 119), (16, 185, 129),
               (245, 158, 11), (99, 102, 241)]

    y = head_h
    for i, sec in enumerate(outline[:8]):
        c = palette[i % len(palette)]
        d.rounded_rectangle([40, y, W - 40, y + row_h - 18], radius=16,
                            fill=(255, 255, 255), outline=c, width=3)
        d.rounded_rectangle([40, y, 44 + 150, y + row_h - 18], radius=16, fill=c)
        rng = str(sec.get("range", f"{i+1}"))[:14]
        d.text((60, y + 14), rng, fill=(255, 255, 255), font=f_head)
        d.text((215, y + 10), str(sec.get("title", ""))[:44], fill=(30, 27, 75), font=f_head)
        d.text((215, y + 50), str(sec.get("summary", ""))[:72], fill=(71, 85, 105), font=f_small)
        y += row_h

    # 연대표
    if timeline:
        ty = y + 20
        d.rounded_rectangle([40, ty, 200, ty + 44], radius=12, fill=(237, 233, 254))
        d.text((60, ty + 8), "연대표", fill=(76, 29, 149), font=f_head)
        ty += 90
        d.line([(70, ty + 26), (W - 70, ty + 26)], fill=(148, 163, 184), width=4)
        n = min(len(timeline), 6)
        step = (W - 200) // max(1, n - 1) if n > 1 else 0
        for i, ev in enumerate(timeline[:n]):
            x = 100 + step * i
            d.ellipse([x - 11, ty + 15, x + 11, ty + 37], fill=palette[i % len(palette)])
            when = str(ev.get("when", ""))[:16]
            what = str(ev.get("what", ""))[:20]
            wb = d.textbbox((0, 0), when, font=f_body)
            d.text((x - (wb[2] - wb[0]) // 2, ty - 18), when, fill=(30, 27, 75), font=f_body)
            tb = d.textbbox((0, 0), what, font=f_small)
            d.text((x - (tb[2] - tb[0]) // 2, ty + 48), what, fill=(71, 85, 105), font=f_small)

    d.text((44, H - 40), "MY 설교 AI 스튜디오 · 본문 구조 도해",
           fill=(148, 163, 184), font=f_small)

    out = io.BytesIO()
    img.save(out, format="PNG")
    return out.getvalue()


def bible_place_links(names):
    """지명·유물에 대해 지도/사진/고고학 자료 바로가기 링크를 만든다."""
    rows = []
    for nm in names:
        q = urllib.parse.quote(str(nm))
        rows.append({
            "이름": nm,
            "지도": f"https://www.openbible.info/geo/search?q={q}",
            "사진·유물": f"https://commons.wikimedia.org/w/index.php?search={q}&title=Special:MediaSearch&type=image",
            "이미지 검색": f"https://www.google.com/search?tbm=isch&q={q}",
        })
    return rows


def generate_verse_card_png(text_str, scripture_str, bg_option="사진", custom_bg_file=None,
                            font_size=42, line_spacing=18, font_color="#FDE047",
                            stroke_color="#000000", overlay_opacity=0.6, church_name="",
                            bg_index=0):
    W = H = 1080
    if bg_option == "직접 업로드" and custom_bg_file is not None:
        try:
            base = PIL.Image.open(custom_bg_file).convert("RGBA").resize((W, H))
        except Exception:
            base = PIL.Image.new("RGBA", (W, H), (15, 23, 42, 255))
    elif bg_option == "기본":
        base = PIL.Image.new("RGBA", (W, H), (15, 23, 42, 255))
    else:
        b = get_background_bytes(f"verse|{bg_index}|{st.session_state.get('bg_shuffle', 0)}",
                                 st.session_state.get("vc_theme", "자연 · 풍경"))
        try:
            base = PIL.Image.open(io.BytesIO(b)).convert("RGBA").resize((W, H)) if b \
                else PIL.Image.new("RGBA", (W, H), (15, 23, 42, 255))
        except Exception:
            base = PIL.Image.new("RGBA", (W, H), (15, 23, 42, 255))

    overlay = PIL.Image.new("RGBA", (W, H), (10, 15, 30, int(255 * overlay_opacity)))
    combined = PIL.Image.alpha_composite(base, overlay)
    draw = PIL.ImageDraw.Draw(combined)

    font_main = get_pil_font(int(font_size))
    font_sub = get_pil_font(30)

    def parse_hex(c):
        try:
            h = c.lstrip('#')
            return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4)) + (255,)
        except Exception:
            return (255, 255, 255, 255)

    f_color = parse_hex(font_color)
    s_color = parse_hex(stroke_color) if stroke_color else None

    wrapped = wrap_korean_text(strip_emoji(text_str), font_main, 880, draw)
    bbox = draw.multiline_textbbox((0, 0), wrapped, font=font_main, spacing=line_spacing)
    t_w, t_h = bbox[2] - bbox[0], bbox[3] - bbox[1]
    x = (W - t_w) // 2
    y = max(120, (H - t_h) // 2 - 50)

    if s_color:
        for dx in (-2, -1, 0, 1, 2):
            for dy in (-2, -1, 0, 1, 2):
                if dx or dy:
                    draw.multiline_text((x + dx, y + dy), wrapped, font=font_main,
                                        fill=s_color, align="center", spacing=line_spacing)
    draw.multiline_text((x, y), wrapped, font=font_main, fill=f_color, align="center", spacing=line_spacing)

    if scripture_str:
        s_txt = f"「 {scripture_str} 」"
        sb = draw.textbbox((0, 0), s_txt, font=font_sub)
        draw.text(((W - (sb[2] - sb[0])) // 2, min(y + t_h + 50, H - 170)),
                  s_txt, fill=(253, 224, 71, 255), font=font_sub)
    if church_name:
        cb = draw.textbbox((0, 0), church_name, font=font_sub)
        draw.text(((W - (cb[2] - cb[0])) // 2, H - 100), church_name,
                  fill=(147, 197, 253, 255), font=font_sub)

    out = io.BytesIO()
    combined.convert("RGB").save(out, format="PNG")
    out.seek(0)
    return out


# ==============================================================================
# TTS
# ==============================================================================
async def _tts(text: str, voice: str, out_path: str):
    communicate = edge_tts.Communicate(text, voice)
    await communicate.save(out_path)


def generate_voiceover_audio(text: str, voice: str = "ko-KR-InJoonNeural") -> str:
    if not HAS_TTS:
        raise RuntimeError("edge-tts 패키지가 설치되어 있지 않습니다. (pip install edge-tts)")
    os.makedirs("./outputs", exist_ok=True)
    out_path = f"./outputs/voiceover_{int(datetime.now().timestamp())}.mp3"
    try:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        loop.run_until_complete(_tts(text, voice, out_path))
        loop.close()
    except RuntimeError:
        asyncio.run(_tts(text, voice, out_path))
    return out_path


# ==============================================================================
# 유튜브 → 9:16 숏츠
# ==============================================================================
def extract_youtube_to_shorts(yt_url, start_sec, duration_sec, title, subtitle_text, church_name=""):
    if not (HAS_YTDLP and HAS_MOVIEPY and HAS_VIDEO_ENGINE):
        raise Exception("영상 처리 모듈(yt-dlp / moviepy / video_engine)이 준비되지 않았습니다.")

    out_dir = "./outputs"
    os.makedirs(out_dir, exist_ok=True)
    src_video = os.path.join(out_dir, "yt_raw_source.mp4")

    vid = re.search(r'(?:v=|/live/|/shorts/|youtu\.be/)([a-zA-Z0-9_-]{11})', yt_url.strip())
    clean_url = f"https://www.youtube.com/watch?v={vid.group(1)}" if vid else yt_url.strip()

    ok = False
    for clients in (['ios'], ['android_creator'], ['mweb'], ['tv_embedded'], ['web']):
        if os.path.exists(src_video):
            try:
                os.remove(src_video)
            except Exception:
                pass
        opts = {
            'format': 'best[ext=mp4][height<=720]/bestvideo[ext=mp4][height<=720]+bestaudio[ext=m4a]/best',
            'outtmpl': src_video, 'overwrites': True, 'quiet': True, 'no_warnings': True,
            'nocheckcertificate': True, 'geo_bypass': True,
            'http_headers': {'User-Agent': 'Mozilla/5.0 (iPhone; CPU iPhone OS 17_4_1 like Mac OS X)'},
            'extractor_args': {'youtube': {'player_client': clients}},
        }
        try:
            with yt_dlp.YoutubeDL(opts) as ydl:
                ydl.download([clean_url])
            if os.path.exists(src_video) and os.path.getsize(src_video) > 100000:
                ok = True
                break
        except Exception:
            continue

    if not ok:
        raise Exception("유튜브 서버가 클라우드 IP를 차단했습니다. "
                        "[🎨 AI 나레이션 & 템플릿 숏츠 제작] 탭에서 영상을 직접 업로드하시면 즉시 9:16 렌더링이 가능합니다.")

    raw = VideoFileClip(src_video)
    start_sec = max(0, min(start_sec, int(raw.duration) - 5))
    end_sec = min(start_sec + duration_sec, int(raw.duration))
    sub = raw.subclip(start_sec, end_sec)

    w, h = sub.size
    target_w = int(h * 9 / 16)
    cropped = sub.crop(x1=w / 2 - target_w / 2, y1=0, x2=w / 2 + target_w / 2, y2=h) if w > target_w else sub
    final = cropped.resize((1080, 1920))
    dur = final.duration

    layers = [final,
              ColorClip(size=(1080, 240), color=(0, 0, 0), duration=dur).set_opacity(0.45).set_position(('center', 100)),
              create_pil_text_clip(title, fontsize=48, color="#FDE047", stroke_color="black",
                                   stroke_width=2, size=(920, None), duration=dur).set_position(("center", 140))]
    if subtitle_text:
        layers.append(create_pil_text_clip(subtitle_text, fontsize=40, color="white", stroke_color="black",
                                           stroke_width=2, size=(900, None), duration=dur).set_position(("center", 1400)))
    if church_name:
        layers.append(create_pil_text_clip(church_name, fontsize=28, color="#93C5FD", stroke_color="black",
                                           stroke_width=1, size=(800, None), duration=dur).set_position(("center", 1780)))

    comp = CompositeVideoClip(layers)
    out_file = os.path.join(out_dir, f"yt_shorts_{int(datetime.now().timestamp())}.mp4")
    comp.write_videofile(out_file, fps=24, codec="libx264", audio_codec="aac", threads=4, preset="ultrafast")
    raw.close(); sub.close(); comp.close()
    return out_file


# ==============================================================================
# 공통 툴바
# ==============================================================================
def render_section_top_toolbar(title: str, content: str, state_key: str, ppt_mode: str = "doc"):
    content = content or "내용 없음"
    c_title, c_btn = st.columns([1.2, 2.8])
    with c_title:
        st.markdown(f"<h3 style='margin:0;font-size:20px;font-weight:800;line-height:1.3;'>{title}</h3>",
                    unsafe_allow_html=True)
    with c_btn:
        b1, b2, b3, b4, b5, b6 = st.columns([1, 1, 1.1, 1.1, 1.1, 1])
        with b1:
            if st.button("✏️ 수정", key=f"edit_btn_{state_key}"):
                st.session_state[f"edit_mode_{state_key}"] = not st.session_state.get(f"edit_mode_{state_key}", False)
        with b2:
            if st.button("📋 복사", key=f"copy_btn_{state_key}"):
                st.session_state[f"show_copy_{state_key}"] = not st.session_state.get(f"show_copy_{state_key}", False)
        with b3:
            st.download_button("📥 워드", data=create_docx_bytes(title, content), file_name=f"{title}.docx",
                               mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                               key=f"dl_docx_{state_key}")
        with b4:
            st.download_button("📥 PDF", data=create_pdf_bytes(title, content), file_name=f"{title}.pdf",
                               mime="application/pdf", key=f"dl_pdf_{state_key}")
        with b5:
            if ppt_mode == "sermon":
                ppt_bytes = generate_sermon_structure_pptx_bytes(
                    st.session_state.get("sermon_title", title),
                    st.session_state.get("sermon_scripture", ""),
                    content,
                    st.session_state.get("full_sermon", ""),
                    bg_seed(0), current_bg_theme()
                )
            else:
                ppt_bytes = create_document_pptx_bytes(title, content,
                                                       bg_seed(0), current_bg_theme())
            st.download_button("📥 PPT", data=ppt_bytes, file_name=f"{title}.pptx",
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                               key=f"dl_ppt_{state_key}")
        with b6:
            st.download_button("📥 txt", data=create_txt_bytes(title, content), file_name=f"{title}.txt",
                               mime="text/plain", key=f"dl_txt_{state_key}")

    if st.session_state.get(f"show_copy_{state_key}", False):
        st.info("💡 아래 상자의 텍스트를 복사하여 사용하세요:")
        st.code(content, language="text")


def st_image_full(data, caption=None):
    """streamlit 버전에 따라 이미지 폭 옵션이 달라 호환 처리"""
    try:
        st.image(data, caption=caption, width='stretch')
    except Exception:
        st.image(data, caption=caption, use_container_width=True)


def show_ai_status():
    """AI 호출이 실패했으면 원인과 해결책을 한국어로 정확히 알려준다."""
    if st.session_state.get("ai_fallback_used"):
        kind = st.session_state.get("ai_error_kind", "other")
        title, howto = ERROR_HELP.get(kind, ERROR_HELP["other"])
        icon = {"quota": "⏳", "gone": "🔄", "perm": "🔒", "key": "🔑"}.get(kind, "⚠️")
        st.warning(f"{icon} **{title}**\n\n{howto}")
        with st.expander("🔎 자세한 오류 내용 보기"):
            for line in st.session_state.get("ai_error_detail", []) or ["(상세 없음)"]:
                st.code(line, language="text")
    elif st.session_state.get("ai_model_used"):
        st.caption(f"✅ 생성 모델: `{st.session_state.ai_model_used}`")


# ------------------------------------------------------------------------------
# 결과물 후처리 : 섹션별 번호 재시작 + 인도자 가이드 블록 강조
# ------------------------------------------------------------------------------
def fix_list_numbering(text: str) -> str:
    """
    번호가 문서 전체에 걸쳐 이어지는 문제(말씀나눔 1,2 → 기도제목 3,4)를 고친다.
    '- 1. ...' 형태의 항목만, 섹션(헤더)이 바뀔 때마다 1번부터 다시 매긴다.
    """
    if not text:
        return text
    out, cnt = [], 0
    for line in text.split("\n"):
        m = LIST_ITEM_RE.match(line)
        if m:
            cnt += 1
            out.append(f"{m.group(1)}- {cnt}. {m.group(3)}")
            continue
        # 인도자 팁·빈 줄은 번호 흐름을 끊지 않는다
        if not line.strip() or LEADER_RE.match(line):
            out.append(line)
            continue
        cnt = 0          # 그 외의 줄(=새 섹션 제목)을 만나면 카운터 리셋
        out.append(line)
    return "\n".join(out)


def _esc(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def render_body(text: str):
    """
    - 📌/💡/🙏 등으로 시작하는 섹션 제목 → 그라데이션 헤더
    - [인도자 팁 / 가이드] 줄 → 제목과 내용 전체가 파란 블록
    - '- 1.' 항목 → 번호 배지
    """
    lines = (text or "").split("\n")
    html_parts = []
    for raw in lines:
        line = raw.rstrip()
        if not line.strip():
            html_parts.append("<div style='height:9px'></div>")
            continue

        m = LEADER_RE.match(line)
        if m:
            body = _esc(m.group(2)).strip() or "&nbsp;"
            html_parts.append(
                f"<span class='leader-block'><b>💡 인도자 가이드</b><br>{body}</span>")
            continue

        if re.match(r'^\s*▸', line):
            html_parts.append(f"<span class='ground-quote'>{_esc(line.strip())}</span>")
            continue

        mi = LIST_ITEM_RE.match(line)
        if mi:
            html_parts.append(
                f"<span class='num-item'><span class='num-badge'>{_esc(mi.group(2))}</span>"
                f"{_esc(mi.group(3))}</span>")
            continue

        if SEC_HEAD_RE.match(line) or (TOP_NUM_RE.match(line) and len(line) < 60):
            html_parts.append(f"<span class='sec-head'>{_esc(line.strip())}</span>")
            continue

        html_parts.append(f"<div class='p-line'>{_esc(line)}</div>")

    st.markdown("<div class='content-box'>" + "".join(html_parts) + "</div>",
                unsafe_allow_html=True)


# ==============================================================================
# 오늘의 만나 · 오늘의 말씀 (설교 대시보드 하단 섹션)
# ==============================================================================
MANNA_THEMES = [
    "새로운 시작과 변화", "하나님을 신뢰함", "고난 중의 연단", "말이 아닌 행함의 사랑",
    "감사와 기쁨", "인내와 견딤", "겸손", "용서", "말의 능력", "지혜",
    "기도의 능력", "이웃 사랑", "정직과 성실", "두려움 대신 담대함", "섬김",
    "가정과 세대", "물질과 청지기", "회개와 돌이킴", "성령의 인도", "소망과 부활",
]


def _today_str():
    """한국 시간 기준 오늘 날짜"""
    try:
        return (datetime.utcnow() + timedelta(hours=9)).strftime("%Y-%m-%d")
    except Exception:
        return datetime.now().strftime("%Y-%m-%d")


def render_manna_section():
    with st.expander("🍞 오늘의 만나 — 날짜별 말씀 · NIV · 예화 한 장", expanded=False):
        st.caption("버튼을 누른 날짜 기준으로 개역개정 구절, NIV 영어 구절, "
                   "그리고 그 구절과 맞물리는 실제 예화(역사·문학·사건·인물·성경)를 한 장으로 만듭니다.")

        today = _today_str()
        c1, c2, c3 = st.columns([1, 1.2, 1.6])
        with c1:
            sel_date = st.date_input("날짜", value=datetime.strptime(today, "%Y-%m-%d").date(),
                                     key="manna_date")
            date_iso = sel_date.strftime("%Y-%m-%d")
        with c2:
            manna_theme = st.selectbox("배경 분위기", list(BG_THEMES.keys()),
                                       index=1, key="manna_bg_theme")
        with c3:
            manna_org = st.text_input(
                "교회 · 공동체 · 단체 이름 (자유 기입)",
                value=st.session_state.get("manna_org_name",
                                           st.session_state.get("cn_church_name", "")),
                placeholder="예: 화광교회 / 사랑의공동체 / 새벽기도회 — 비워두면 표시 안 함",
                key="manna_org_input")
            st.session_state["manna_org_name"] = manna_org

        key = f"manna::{date_iso}"

        b1, b2 = st.columns([1, 1])
        with b1:
            go = st.button("✨ AI로 오늘의 만나 생성", type="primary", key="btn_gen_manna")
        with b2:
            write_mode = st.toggle("✍️ 직접 작성 / 수정하기", key="manna_write_mode",
                                   value=False)

        # ── AI 생성
        if go:
            with st.spinner("오늘의 말씀과 예화를 준비하는 중..."):
                d0 = datetime.strptime(date_iso, "%Y-%m-%d")
                seed_topic = MANNA_THEMES[(d0.timetuple().tm_yday) % len(MANNA_THEMES)]
                task = f"""[출력 형식 — 아래 JSON 하나만. 설명 문장 금지]
{{
 "topic": "오늘의 주제를 한 낱말~한 구절로",
 "ko_verse": "개역개정 성경 구절 본문 (2~4줄 분량, 40~90자)",
 "ko_ref": "예: 고후 5:17",
 "en_verse": "같은 구절의 NIV 영어 본문",
 "en_ref": "예: 2Co 5:17",
 "story_title": "예화 제목 (8~14자)",
 "story_body": "예화 본문 (280~380자, 한 문단). 실제로 있었던 일만 쓸 것."
}}

[오늘의 조건]
- 날짜: {d0.strftime('%Y년 %m월 %d일')}
- 오늘의 주제 방향: {seed_topic}

[반드시 지킬 것]
1. ko_verse 는 **개역개정판** 본문을 정확히 인용하십시오. 다른 역본을 쓰지 마십시오.
   확실히 아는 구절만 쓰고, 장절(ko_ref)을 정확히 표기하십시오.
2. en_verse 는 같은 구절의 NIV 본문입니다. 한국어 번역이 아니라 실제 영어 본문이어야 합니다.
3. story_body 는 **실제 예화**여야 합니다. 다음 중 하나에서 고르십시오.
   ① 실제 역사적 사건·전쟁·사고  ② 실존 인물의 일화(연도·이름 명시)
   ③ 문학·예술 작품과 그 창작 배경  ④ 성경 속 실제 에피소드
   지어낸 '어떤 성도의 이야기'는 절대 금지입니다.
4. 예화는 반드시 오늘의 구절이 말하는 바와 맞물려야 합니다.
5. 예화 마지막 한 문장은 구절과의 연결을 담담하게 맺으십시오. 설교조로 훈계하지 마십시오.
6. 100% 한국어(en_verse 제외)."""
                data = get_ai_response(
                    build_research_prompt(task, "오늘의 말씀", seed_topic),
                    is_json=True, temperature=0.85, kind="manna", max_tokens=4000)
                if isinstance(data, dict) and data:
                    st.session_state[key] = data
            st.rerun()

        data = st.session_state.get(key) or {}

        # ── 직접 작성 / 수정 폼
        if write_mode:
            st.markdown("#### ✍️ 직접 작성 · 수정")
            st.caption("AI로 만든 내용을 고치거나, 목사님이 준비하신 말씀과 예화를 직접 적어 넣으세요.")
            with st.form(f"manna_form_{date_iso}"):
                f1, f2 = st.columns([2.4, 1])
                with f1:
                    kv = st.text_area("개역개정 본문", value=data.get("ko_verse", ""), height=100)
                with f2:
                    kr = st.text_input("성경 구절 (예: 고후 5:17)", value=data.get("ko_ref", ""))
                g1, g2 = st.columns([2.4, 1])
                with g1:
                    ev = st.text_area("NIV 영어 본문", value=data.get("en_verse", ""), height=100)
                with g2:
                    er = st.text_input("영문 약어 (예: 2Co 5:17)", value=data.get("en_ref", ""))
                sttl = st.text_input("예화 제목", value=data.get("story_title", ""))
                sbody = st.text_area("예화 본문", value=data.get("story_body", ""), height=190)
                saved = st.form_submit_button("💾 저장하고 다시 그리기", type="primary")
            if saved:
                st.session_state[key] = {
                    "ko_verse": kv, "ko_ref": kr, "en_verse": ev, "en_ref": er,
                    "story_title": sttl, "story_body": sbody,
                    "topic": data.get("topic", ""),
                }
                st.success("저장했습니다. 아래 이미지가 새로 그려집니다.")
                st.rerun()

        if not data:
            if st.session_state.get("ai_fallback_used"):
                show_ai_status()
            else:
                st.caption("위 버튼을 눌러 오늘의 만나를 만들거나, "
                           "[✍️ 직접 작성 / 수정하기]를 켜서 직접 적어 넣으세요.")
            return

        show_ai_status()
        png = render_manna_png(date_iso, data.get("ko_verse", ""), data.get("ko_ref", ""),
                               data.get("en_verse", ""), data.get("en_ref", ""),
                               data.get("story_title", ""), data.get("story_body", ""),
                               seed=f"manna|{date_iso}|{st.session_state.get('bg_shuffle', 0)}",
                               theme=manna_theme, org_name=manna_org.strip())
        st_image_full(png, caption=f"오늘의 만나 · {date_iso}")

        dl1, dl2 = st.columns(2)
        with dl1:
            st.download_button("📥 이미지(PNG) 내려받기", data=png,
                               file_name=_safe_filename("오늘의만나", date_iso, manna_org, "png"),
                               mime="image/png", key=f"dl_manna_{date_iso}")
        with dl2:
            if st.button("🔀 배경 이미지 바꾸기", key=f"manna_shuffle_{date_iso}"):
                st.session_state.bg_shuffle = int(st.session_state.get("bg_shuffle", 0)) + 1
                st.rerun()

        txt = (f"🍞 오늘의 만나 · {date_iso}"
               + (f" · {manna_org}" if manna_org.strip() else "") + "\n\n"
               f"📖 {data.get('ko_ref','')} (개역개정)\n{data.get('ko_verse','')}\n\n"
               f"🌍 {data.get('en_ref','')} (NIV)\n{data.get('en_verse','')}\n\n"
               f"▪ {data.get('story_title','')}\n{data.get('story_body','')}")
        render_section_top_toolbar(f"오늘의만나_{date_iso}", txt, f"manna_{date_iso}")
        render_body(txt)


def _safe_filename(prefix: str, date_iso: str, name: str = "", ext: str = "png") -> str:
    """파일명에 쓸 수 없는 문자를 정리한다."""
    tail = re.sub(r'[\\/:*?"<>|\s]+', "", str(name or ""))[:20]
    return f"{prefix}_{date_iso}" + (f"_{tail}" if tail else "") + f".{ext}"


def render_today_word_section():
    with st.expander("📖 오늘의 말씀 — 말씀카드 이미지 자동 생성", expanded=False):
        st.caption("버튼을 누르면 그날 날짜가 적힌 말씀카드가 배경 사진과 함께 만들어집니다. "
                   "교회 SNS나 단톡방에 그대로 올리실 수 있습니다.")

        today = _today_str()
        c1, c2, c3 = st.columns([1, 1.2, 1])
        with c1:
            sel_date = st.date_input("날짜", value=datetime.strptime(today, "%Y-%m-%d").date(),
                                     key="tw_date")
            date_iso = sel_date.strftime("%Y-%m-%d")
        with c2:
            tw_theme = st.selectbox("배경 분위기", list(BG_THEMES.keys()), index=0, key="tw_theme")
        with c3:
            tw_size = st.selectbox("규격", ["1:1 정사각형", "4:5 세로형"], key="tw_size")

        n1, n2 = st.columns([2.2, 1])
        with n1:
            tw_church = st.text_input(
                "카드 하단에 넣을 이름 (교회 · 단체 · 부서 · 개인 무엇이든)",
                value=st.session_state.get("tw_church_name",
                                           st.session_state.get("cn_church_name", "")),
                placeholder="예: 화광교회 / 청년부 / 사랑의교회 새벽기도 / 비워두면 표시 안 함",
                key="tw_church_input")
            st.session_state["tw_church_name"] = tw_church
        with n2:
            st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)
            go = st.button("📖 오늘의 말씀 생성", type="primary", key="btn_gen_today_word")

        src = st.radio("구절 가져오기", ["AI 추천 구절", "오늘의 만나 구절 사용", "직접 입력"],
                       horizontal=True, key="tw_src")

        key = f"todayword::{date_iso}"
        manna = st.session_state.get(f"manna::{date_iso}") or {}

        if src == "직접 입력":
            v = st.text_area("말씀 문구", value=st.session_state.get(f"{key}_v", ""),
                             height=90, key=f"{key}_vin")
            r = st.text_input("성경 구절", value=st.session_state.get(f"{key}_r", ""),
                              key=f"{key}_rin")
            if go:
                st.session_state[key] = {"verse": v, "ref": r}
                st.rerun()
        elif src == "오늘의 만나 구절 사용":
            if go:
                if manna.get("ko_verse"):
                    st.session_state[key] = {"verse": manna["ko_verse"], "ref": manna.get("ko_ref", "")}
                else:
                    st.warning("먼저 위의 [🍞 오늘의 만나]를 생성해 주세요.")
                st.rerun()
        else:
            if go:
                with st.spinner("오늘 나눌 말씀을 고르는 중..."):
                    d0 = datetime.strptime(date_iso, "%Y-%m-%d")
                    seed_topic = MANNA_THEMES[(d0.timetuple().tm_yday + 7) % len(MANNA_THEMES)]
                    task = f"""[출력 형식 — 아래 JSON 하나만]
{{"verse": "개역개정 성경 구절 본문 (30~70자, 카드에 넣기 좋은 길이)",
  "ref": "예: 시편 23:1",
  "why": "오늘 이 구절을 나누는 이유 한 문장"}}

- 날짜: {d0.strftime('%Y년 %m월 %d일')}
- 오늘의 주제 방향: {seed_topic}
- 개역개정 본문을 정확히 인용하십시오. 확실히 아는 구절만 쓰십시오.
- 너무 긴 구절은 피하고, 한 화면에 들어오는 길이로 고르십시오."""
                    data = get_ai_response(
                        build_research_prompt(task, "오늘의 말씀", seed_topic),
                        is_json=True, temperature=0.9, kind="today_verse", max_tokens=1500)
                    st.session_state[key] = data if isinstance(data, dict) else {}
                st.rerun()

        data = st.session_state.get(key) or {}
        verse, ref = data.get("verse", ""), data.get("ref", "")
        if not verse:
            if st.session_state.get("ai_fallback_used"):
                show_ai_status()
            else:
                st.caption("위 버튼을 눌러 오늘의 말씀 카드를 만들어 보세요.")
            return

        show_ai_status()
        if data.get("why"):
            st.caption(f"오늘 이 말씀을 나누는 이유 — {data['why']}")

        png = render_today_verse_png(
            date_iso, verse, ref, tw_church.strip(),
            seed=f"today|{date_iso}|{st.session_state.get('bg_shuffle', 0)}",
            theme=tw_theme, size_key=tw_size)

        p1, p2 = st.columns([1.2, 1])
        with p1:
            st_image_full(png, caption=f"오늘의 말씀 · {date_iso}")
        with p2:
            st.download_button("📥 말씀카드 PNG 내려받기", data=png,
                               file_name=_safe_filename("오늘의말씀", date_iso, tw_church, "png"),
                               mime="image/png",
                               key=f"dl_tw_{date_iso}")
            if st.button("🔀 배경 사진 바꾸기", key=f"tw_shuffle_{date_iso}"):
                st.session_state.bg_shuffle = int(st.session_state.get("bg_shuffle", 0)) + 1
                st.rerun()
            st.markdown(f"<div class='lib-card' style='padding:14px;margin-top:10px;'>"
                        f"<div style='color:#fde047;font-weight:800;'>「 {_esc(ref)} 」</div>"
                        f"<div style='color:#e8ecff;margin-top:8px;line-height:1.7;'>"
                        f"{_esc(verse)}</div></div>", unsafe_allow_html=True)

        # ── 문서 내려받기 (수정 · 복사 · 워드 · PDF · PPT · txt)
        d0 = datetime.strptime(date_iso, "%Y-%m-%d")
        wk = ["월", "화", "수", "목", "금", "토", "일"][d0.weekday()] + "요일"
        doc_field = f"todayword_doc::{date_iso}"
        if not st.session_state.get(doc_field):
            lines = [f"📖 오늘의 말씀 · {d0.strftime('%Y년 %m월 %d일')} {wk}"]
            if tw_church.strip():
                lines.append(f"　{tw_church.strip()}")
            lines += ["", f"1. 오늘의 성구 (개역개정)", f"- 1. {ref}",
                      f"- 2. {verse}", ""]
            if data.get("why"):
                lines += ["2. 오늘 이 말씀을 나누는 이유", f"- 1. {data['why']}", ""]
            lines += ["3. 함께 드리는 기도",
                      f"- 1. 주님, 오늘 주신 {ref} 말씀을 마음에 새기고 하루를 살아가게 하옵소서."]
            st.session_state[doc_field] = "\n".join(lines)

        st.write("")
        render_section_top_toolbar(
            _safe_filename("오늘의말씀", date_iso, tw_church, "").rstrip("."),
            st.session_state[doc_field], f"tw_doc_{date_iso}")
        if editable_section(f"tw_doc_{date_iso}", doc_field, "오늘의 말씀 내용 편집", height=260):
            render_body(st.session_state[doc_field])


# ==============================================================================
# 본문 연구 도구 7종 (설교 등록 → AI 강해설교문 생성 탭의 하위 항목)
# ==============================================================================
def _rkey(name: str, scripture: str) -> str:
    """본문이 바뀌면 이전 결과가 자동으로 무효가 되도록 본문을 키에 넣는다."""
    return f"res_{name}::{scripture}"


def _research_block(label: str, name: str, scripture: str, task: str, topic: str,
                    theology: str, temp: float = 0.55, tokens: int = 8192,
                    use_search: bool = True, expanded: bool = False):
    """연구 도구 한 칸(생성 버튼 + 결과 + 내려받기)을 그린다."""
    key = _rkey(name, scripture)
    with st.expander(label, expanded=expanded):
        c1, c2 = st.columns([1, 2])
        with c1:
            go = st.button("✨ 생성 / 다시 생성", key=f"btn_{name}_{abs(hash(scripture)) % 99999}",
                           type="primary")
        with c2:
            st.caption(f"대상 본문: {scripture}" + (f"　·　주제: {topic}" if topic.strip() else ""))
        if go:
            with st.spinner("자료를 찾아 정리하는 중입니다... (30초~1분)"):
                st.session_state[key] = get_ai_response(
                    build_research_prompt(task, scripture, topic, theology),
                    is_json=False, temperature=temp, kind="research",
                    use_search=use_search, max_tokens=tokens)
            st.rerun()

        val = st.session_state.get(key, "")
        if val:
            show_ai_status()
            if st.session_state.get("ai_search_used"):
                st.caption("🌐 웹 검색 근거를 사용해 작성했습니다.")
            render_section_top_toolbar(f"{scripture}_{name}", val, f"{name}_{abs(hash(scripture)) % 9999}")
            render_body(val)
        else:
            st.caption("위 버튼을 눌러 자료를 만들어 보세요.")
    return st.session_state.get(key, "")


def context_to_text(d: dict, scripture: str) -> str:
    """문맥 연구 JSON 을 내려받기 좋은 텍스트로 조립"""
    L = [f"🧭 본문 문맥 연구 · {scripture}", ""]

    def sec(title, body):
        if body:
            L.extend([title, str(body), ""])

    L.append("📖 기본 정보")
    for lab, k in [("저자", "author"), ("저작 시기", "date_written"),
                   ("청중(대상)", "audience"), ("기록 목적", "purpose")]:
        if d.get(k):
            L.append(f"- {lab}: {d[k]}")
    L.append("")
    sec("🏛️ 역사적 배경", d.get("historical"))
    sec("📜 성경 문맥적 배경", d.get("biblical_context"))
    sec("👥 사회·문화적 배경", d.get("socio_cultural"))
    sec("⚖️ 정치·경제적 배경", d.get("political_economic"))

    if d.get("outline"):
        L.append("🧩 단락 나누기")
        for i, o in enumerate(d["outline"], start=1):
            L.append(f"- {i}. [{o.get('range','')}] {o.get('title','')}")
            if o.get("summary"):
                L.append(f"   ▸ {o['summary']}")
        L.append("")

    if d.get("key_verses"):
        L.append("💎 본문의 핵심 구절")
        for i, v in enumerate(d["key_verses"], start=1):
            L.append(f"- {i}. {v.get('ref','')} “{v.get('text','')}”")
            if v.get("why"):
                L.append(f"   ▸ {v['why']}")
        L.append("")

    if d.get("homiletic"):
        L.append("🎙️ 본문의 설교적 원리")
        for i, x in enumerate(d["homiletic"], start=1):
            L.append(f"- {i}. {x}")
        L.append("")
    if d.get("expository"):
        L.append("📐 강해적 원리")
        for i, x in enumerate(d["expository"], start=1):
            L.append(f"- {i}. {x}")
        L.append("")

    sec("⛪ 개혁주의 신학에서 본문의 위치와 해석", d.get("reformed"))
    sec("🌍 선교적 교회의 관점으로 본 본문", d.get("missional"))

    if d.get("timeline"):
        L.append("🗓️ 연대표")
        for i, t in enumerate(d["timeline"], start=1):
            L.append(f"- {i}. {t.get('when','')} — {t.get('what','')}")
        L.append("")
    if d.get("compare_table"):
        L.append("📊 본문 당시 vs 오늘")
        for i, r in enumerate(d["compare_table"], start=1):
            L.append(f"- {i}. {r.get('항목','')} | 당시: {r.get('본문 당시','')} | 오늘: {r.get('오늘 우리','')}")
        L.append("")
    if d.get("places"):
        L.append("🗺️ 관련 지명: " + ", ".join(map(str, d["places"])))
    if d.get("artifacts"):
        L.append("🏺 관련 유물·고고학 자료: " + ", ".join(map(str, d["artifacts"])))
    return "\n".join(L)


def render_context_section(scripture: str, topic: str, theology: str):
    """문맥 연구 — 표 · 도해 이미지 · 지도/유물 링크까지"""
    key = _rkey("context", scripture)
    with st.expander("🧭 문맥 연구 — 배경 · 단락 · 도표 · 지도 · 도해", expanded=False):
        c1, c2 = st.columns([1, 2])
        with c1:
            go = st.button("✨ 문맥 연구 생성", type="primary",
                           key=f"btn_ctx_{abs(hash(scripture)) % 99999}")
        with c2:
            st.caption("역사적·문맥적·사회문화적·정치경제적 배경, 저자·시기·청중·목적, "
                       "단락 나누기, 핵심 구절, 설교적·강해적 원리, 개혁주의·선교적 관점까지 한 번에.")
        if go:
            with st.spinner("본문 배경을 조사하고 도표를 만드는 중입니다... (1분 내외)"):
                task = """[출력 형식 — 아래 JSON 하나만 출력. 설명 문장 금지]
{
 "author": "저자와 그 근거(2~3문장)",
 "date_written": "저작 시기와 근거(2~3문장)",
 "audience": "1차 청중이 누구였는지, 그들의 상황(3문장)",
 "purpose": "기록 목적(3문장)",
 "historical": "역사적 배경 — 당시 제국/왕조, 사건, 연대를 구체적으로 (6~8문장)",
 "biblical_context": "성경 문맥적 배경 — 앞뒤 문맥, 이 책 전체 구조 속 위치, 구속사적 자리 (6~8문장)",
 "socio_cultural": "사회·문화적 배경 — 관습, 신분, 가정, 종교 관행 (6~8문장)",
 "political_economic": "정치·경제적 배경 — 통치 체제, 세금, 생업, 계층 (5~7문장)",
 "outline": [{"range":"1-3절","title":"단락 제목","summary":"한 문장 요약"}],
 "key_verses": [{"ref":"장:절","text":"구절 전문(개역개정)","why":"왜 핵심인지 2문장"}],
 "homiletic": ["본문의 설교적 원리 4가지"],
 "expository": ["강해적 원리 4가지 — 본문을 어떻게 열어야 하는가"],
 "reformed": "개혁주의 신학에서 이 본문의 위치와 해석 (6문장). 언약·주권·은혜 관점 포함",
 "missional": "선교적 교회의 관점으로 본 본문 (6문장). 보내심받은 공동체로서의 적용",
 "timeline": [{"when":"주전/주후 연대","what":"사건"}],
 "compare_table": [{"항목":"예: 감옥","본문 당시":"당시 상황","오늘 우리":"오늘의 대응물"}],
 "places": ["본문과 관련된 실제 지명 4~6개"],
 "artifacts": ["관련 유물·고고학 자료·유적 이름 3~5개"]
}
※ outline 은 3~6개, key_verses 는 3개, timeline 은 4~6개, compare_table 은 4~5행.
※ 연대·지명·인물은 정확하게. 학자 간 견해가 갈리면 '다수설/소수설'을 함께 적으십시오."""
                data = get_ai_response(
                    build_research_prompt(task, scripture, topic, theology),
                    is_json=True, temperature=0.35, kind="context", max_tokens=12000)
                st.session_state[key] = data if isinstance(data, dict) else {}
            st.rerun()

        d = st.session_state.get(key) or {}
        if not d:
            if st.session_state.get("ai_fallback_used"):
                show_ai_status()
            else:
                st.caption("위 버튼을 눌러 문맥 연구를 시작하세요.")
            return

        show_ai_status()
        text_all = context_to_text(d, scripture)
        render_section_top_toolbar(f"{scripture}_문맥연구", text_all,
                                   f"ctx_{abs(hash(scripture)) % 9999}")

        m1, m2, m3, m4 = st.columns(4)
        for col, lab, k in ((m1, "저자", "author"), (m2, "저작 시기", "date_written"),
                            (m3, "청중", "audience"), (m4, "목적", "purpose")):
            with col:
                st.markdown(
                    f"<div class='lib-card' style='padding:14px;'>"
                    f"<div style='color:#fde047;font-weight:800;font-size:13px;'>{lab}</div>"
                    f"<div style='color:#e8ecff;font-size:13px;margin-top:6px;line-height:1.6;'>"
                    f"{_esc(str(d.get(k,'—'))[:180])}</div></div>", unsafe_allow_html=True)

        st.markdown("#### 🖼️ 본문 구조 도해 (슬라이드에 바로 사용)")
        png = render_outline_diagram(scripture,
                                     json.dumps(d.get("outline", []), ensure_ascii=False),
                                     json.dumps(d.get("timeline", []), ensure_ascii=False))
        st_image_full(png, caption=f"{scripture} 단락 구조 · 연대표")
        st.download_button("📥 도해 이미지(PNG) 내려받기", data=png,
                           file_name=f"{scripture}_구조도해.png", mime="image/png",
                           key=f"dl_diag_{abs(hash(scripture)) % 9999}")

        if d.get("outline"):
            st.markdown("#### 🧩 단락 나누기")
            st.table([{"구간": o.get("range", ""), "단락 제목": o.get("title", ""),
                       "요약": o.get("summary", "")} for o in d["outline"]])
        if d.get("compare_table"):
            st.markdown("#### 📊 본문 당시 vs 오늘")
            st.table(d["compare_table"])
        if d.get("timeline"):
            st.markdown("#### 🗓️ 연대표")
            st.table([{"연대": t.get("when", ""), "사건": t.get("what", "")} for t in d["timeline"]])

        places = list(d.get("places", [])) + list(d.get("artifacts", []))
        if places:
            st.markdown("#### 🗺️ 지도 · 유물 사진 바로가기")
            st.caption("AI는 사진을 직접 가져올 수 없어, 신뢰할 수 있는 자료 사이트로 바로 연결합니다.")
            for row in bible_place_links(places[:10]):
                st.markdown(
                    f"**{row['이름']}** — [🗺️ 성경 지도]({row['지도']}) · "
                    f"[🏺 유물·사진(위키미디어)]({row['사진·유물']}) · "
                    f"[🔍 이미지 검색]({row['이미지 검색']})")

        st.markdown("#### 📜 배경 상세")
        render_body(text_all)


def render_research_tools(scripture: str, topic: str, theology: str):
    st.markdown("### 🔬 본문 연구 도구")
    st.caption("아래 항목들은 설교 원고가 없어도, 선택한 본문과 주제만으로 자료를 만들어 옵니다.")

    # 1) 문맥
    render_context_section(scripture, topic, theology)

    # 2) 성경 여러 번역본 비교
    render_version_compare_section(scripture, topic, theology)

    # 3) 원어 주해
    _research_block(
        "🔤 원어 주해 — 절별 히브리어/헬라어 연구", "original", scripture,
        """[출력 형식 — 본문의 각 절마다 아래 틀 그대로 반복]

📖 (장:절)
원문: (히브리어 또는 헬라어 원문 그대로)
음역: (한글 음역)
직역: (한국어 직역 — 어순을 살려 딱딱하게)

- 1. 핵심 단어 ①: 원어(음역, 스트롱번호) — 기본 뜻 / 이 문맥에서의 뜻
   ▸ 문법: (품사·시제·태·법·인칭·수, 명사면 격·수·성)
   ▸ 신학적 함의: (2문장)
- 2. 핵심 단어 ②: (동일 형식)
- 3. 구문 분석: (주절과 종속절 관계, 접속사·전치사·분사구문의 기능을 2~3문장)
- 4. 주해: (이 절이 말하는 바를 3~4문장. 오해하기 쉬운 지점이 있으면 지적)

(본문의 모든 절을 위 형식으로 다룬 뒤, 마지막에)

🔑 본문 전체 원어 요약
- 1. 반복되는 핵심 어휘와 그 효과
- 2. 문학적 구조(교차대구·수미상관·점층 등)가 있으면 지적
- 3. 번역본 간 차이가 큰 지점과 그 이유

※ 원어 철자와 스트롱번호는 정확하게. 확신이 없으면 스트롱번호를 생략하십시오.
※ 절 수가 많으면 의미 단위로 2~3절씩 묶어도 됩니다.""",
        topic, theology, temp=0.3, tokens=14000)

    # 4) 예화 4종
    _research_block(
        "💎 설교 예화 12개 — 성경 · 역사 · 문학예술 · 현대", "illust4", scripture,
        """[출력 형식 — 네 갈래, 각 3개씩 총 12개. 아래 틀 그대로]

📖 성경 예화 (성경 66권 전체에서)
- 1. 「제목」 (본문: 성경 장절)
   ▸ 내용: (5~7문장. 인물·상황·전환점·결말을 이야기로)
   ▸ 오늘 본문과의 연결: (왜 이 본문과 맞물리는지 2~3문장)
- 2. (동일 형식)
- 3. (동일 형식)

🏛️ 역사 예화 (실제 역사적 사건·사고·인물의 일화)
- 1. 「제목」 (연도 / 장소 / 인물)
   ▸ 내용: (5~7문장. 연도·지명·이름을 정확히)
   ▸ 오늘 본문과의 연결: ...
- 2. (동일 형식)
- 3. (동일 형식)

🎼 문학·예술 예화 (소설·시·회화·음악·영화에서)
- 1. 「작품명」 (작가/작곡가/화가, 발표 연도, 장르)
   ▸ 내용: (5~7문장. 어떤 장면·악장·구절인지 구체적으로)
   ▸ 오늘 본문과의 연결: ...
- 2. (동일 형식)
- 3. (동일 형식)

🌍 현대 실제 사례 (2018년 이후, 실제로 일어난 일)
- 1. 「제목」 (연도, 국가/지역, 인물 또는 사건명)
   ▸ 내용: (5~7문장. 실제 보도·기록에 근거한 사건만. 지어낸 미담 금지)
   ▸ 오늘 본문과의 연결: ...
- 2. (동일 형식)
- 3. (동일 형식)

[이 항목의 특별 규칙]
- 성경 예화는 오늘 본문 자체를 다시 말하지 말고, 성경 66권의 '다른' 본문에서 가져오십시오.
- 역사 예화는 교회사에 국한하지 말고 세계사·한국사의 사건·사고·일화까지 포함하십시오.
- 문학·예술 예화는 작품명과 창작자를 정확히 쓰고, 줄거리를 지어내지 마십시오.
- 현대 사례는 반드시 실제로 일어난 일이어야 합니다. 만들어낸 '어떤 성도의 이야기'는 절대 금지입니다.
  확실한 사건이 3개가 안 되면, 확실한 것만 쓰고 부족한 자리에 (추가 확인 필요)라고 적으십시오.""",
        topic, theology, temp=0.6, tokens=14000)

    # 5) 명언
    _research_block(
        "🗣️ 설교 명언 — 주제 관련 신학자·목회자 5인 이상", "quotes5", scripture,
        """[출력 형식 — 최소 5명, 서로 다른 인물. 설교 원고가 아니라 '본문과 주제'에 맞춰 고르십시오]

🗣️ 이 본문·주제를 살리는 명언

- 1. "명언 원문(한국어 번역)" — 인물 이름 (생몰연대 / 직함)
   ▸ 원문 출처: (책 제목·설교 제목·문헌. 불확실하면 '출처 확인 필요'라고 정직하게)
   ▸ 왜 이 본문에 맞는가: (2문장)
   ▸ 활용 위치: (서론 / 제1대지 / 제2대지 / 결론 중 어디에, 어떻게 인용할지)
- 2. ~ 5. (동일 형식)

📌 인물 구성 규칙
- 교부·종교개혁자(어거스틴, 칼빈, 루터 등) 최소 1명
- 근현대 설교자·신학자(스펄전, 로이드 존스, 본회퍼, 존 스토트, 팀 켈러, 유진 피터슨 등) 최소 2명
- 한국 교회 목회자 또는 한국 기독교 사상가 최소 1명

⚠️ 그 인물이 하지 않은 말을 절대 지어내지 마십시오. 확신이 없으면 다른 인물로 바꾸십시오.""",
        topic, theology, temp=0.45, tokens=8192)

    # 6) 주석가 관점
    _research_block(
        "📚 주석가 관점 — 본문에 대한 주요 주석가 5인 이상 해석", "commentary", scripture,
        """[출력 형식 — 이 본문을 실제로 다룬 주석가·신학자 5인 이상]

📚 주석가별 해석 비교

- 1. 존 칼빈 (Calvin) — 『기독교 강요』/『성경주석』
   ▸ 대상 범위: (전체 / 몇 절 / 어느 단락)
   ▸ 핵심 주장: (3~4문장. 이 주석가가 이 본문에서 무엇을 강조했는지)
   ▸ 대표적 논지: (가능하면 요지를 인용하듯 제시하되, 지어낸 직접 인용문은 금지)
- 2. ~ 5. (동일 형식. 서로 다른 신학 전통에서 고르십시오)

⚖️ 쟁점별 견해 대조
- 1. (쟁점 1 — 예: 이 절의 주어가 누구인가)
   ▸ 다수설: ... / 소수설: ... / 판단 근거: ...
- 2. (쟁점 2)
- 3. (쟁점 3)

🧭 설교자를 위한 정리
- 1. 강단에서 취할 해석과 그 이유
- 2. 피해야 할 해석과 그 이유
- 3. 청중에게 굳이 설명하지 않아도 될 학문적 쟁점

📌 주석가 구성 규칙
- 개혁주의 계열(칼빈, 매튜 헨리, 헨드릭슨, 존 머레이 등) 최소 2명
- 현대 비평·주해 계열(F.F. 브루스, 고든 피, 더글라스 무, N.T. 라이트, 크레이그 키너 등) 최소 2명
- 설교자 계열(스펄전, 로이드 존스, 존 스토트 등) 최소 1명

⚠️ 실제로 그 주석가가 그 본문을 다루었는지 확신이 없으면 다른 주석가로 대체하십시오.
⚠️ 존재하지 않는 책 제목이나 직접 인용문을 만들지 마십시오. 요지 서술로 대신하십시오.""",
        topic, theology, temp=0.35, tokens=12000)

    # 7) 추천 찬양
    render_praise_section(scripture, topic, theology)

    # 8) 설교 제목
    _research_block(
        "🏷️ 설교 제목 추천 5개 — 참신하고 인상적인 제목", "titles5b", scripture,
        """[출력 형식 — 아래 틀 그대로]

🏷️ 설교 제목 추천 5개

- 1. 「제목」  (유형: 은유·이미지형)
   ▸ 근거 구절: (본문 안의 어느 구절에서 나왔는지)
   ▸ 왜 참신한가: (한 문장)
   ▸ 부제: (선택)
- 2. 「제목」  (유형: 역설·반전형)
- 3. 「제목」  (유형: 질문형)
- 4. 「제목」  (유형: 명령·초청형)
- 5. 「제목」  (유형: 본문 핵심어 인용형)

🎯 가장 추천하는 제목 1개와 그 이유 (3문장)

[제목 작성 규칙 — 반드시 지킬 것]
- 8~14자 사이. 길면 강단 스크린에서 힘을 잃습니다.
- 본문에 실제로 등장하는 단어·이미지·동사를 한 개 이상 살릴 것.
- 다음과 같은 상투적 제목은 절대 금지: '은혜의 삶', '축복의 통로', '믿음의 승리',
  '주님과 동행하는 삶', '감사의 능력', '기도의 힘'.
- 낯설게 하기: 익숙한 신앙 용어 대신 구체적 명사·동사를 쓰십시오.
  (예: '기도합시다' 대신 '무릎이 먼저 도착했다')
- 다섯 제목이 서로 다른 각도여야 합니다. 같은 말의 변주 금지.""",
        topic, theology, temp=0.9, tokens=6000)


BIBLE_VERSIONS = [
    ("개역개정", "한국어 · 개역개정판 (기본 역본)"),
    ("표준새번역(새번역)", "한국어 · 대한성서공회 새번역"),
    ("새한글성경", "한국어 · 새한글성경"),
    ("메시지(한글)", "메시지성경 영문을 한국어로 옮긴 의역"),
    ("메시지(영문)", "The Message (Eugene Peterson)"),
    ("NIV", "New International Version"),
    ("NASB", "New American Standard Bible"),
    ("RSV", "Revised Standard Version"),
    ("원문(히브리어/헬라어)", "구약=히브리어, 신약=헬라어 원문"),
    ("원문 음역", "원문의 한글 음역"),
]


def render_version_compare_section(scripture: str, topic: str, theology: str):
    """성경 여러 번역본 비교 — 행=번역본 / 열=절 표"""
    key = _rkey("versions", scripture)
    with st.expander("📚 성경 여러 번역본 비교 — 개역개정 · 새번역 · 메시지 · NIV · 원문", expanded=False):
        c1, c2 = st.columns([1, 2])
        with c1:
            go = st.button("✨ 번역본 비교표 생성", type="primary",
                           key=f"btn_ver_{abs(hash(scripture)) % 99999}")
        with c2:
            st.caption("개역개정 · 표준새번역 · 새한글성경 · 메시지(영/한) · NIV · NASB · RSV · "
                       "원문(히브리어/헬라어) · 음역을 절별로 나란히 비교합니다.")

        if go:
            with st.spinner("번역본을 절별로 모으는 중입니다... (1분 내외)"):
                vlist = "\n".join(f'   - "{n}"' for n, _ in BIBLE_VERSIONS)
                task = f"""[출력 형식 — 아래 JSON 하나만. 설명 문장 금지]
{{
 "verses": ["1", "2", "3"],
 "rows": [
   {{"version": "개역개정", "texts": ["1절 본문", "2절 본문", "3절 본문"]}},
   {{"version": "표준새번역(새번역)", "texts": ["...", "...", "..."]}}
 ]
}}

[작업]
{scripture} 본문을 절 단위로 나누고, 아래 역본을 **정확히 이 순서대로** 모두 채우십시오.
{vlist}

[반드시 지킬 것]
1. "verses" 에는 본문에 실제로 포함된 절 번호만 순서대로 넣으십시오 (문자열).
2. 모든 row 의 "texts" 배열 길이는 "verses" 길이와 정확히 같아야 합니다.
3. 개역개정은 한 글자도 바꾸지 말고 정확히 인용하십시오. 이 역본이 기준입니다.
4. "메시지(한글)" 은 메시지성경 영문("메시지(영문)")을 한국어로 자연스럽게 옮긴 의역입니다.
   직역이 아니라 메시지성경 특유의 구어체 어감을 살리십시오.
5. "원문(히브리어/헬라어)" 은 구약이면 히브리어(BHS), 신약이면 헬라어(NA28) 본문을 그대로 쓰십시오.
6. "원문 음역" 은 그 원문을 한글 발음으로 옮긴 것입니다. (예: 아쉬레이 하이쉬 아쉐르)
7. 확실하지 않은 역본의 절은 지어내지 말고 "(확인 필요)" 라고만 쓰십시오.
8. 절이 12개를 넘으면 앞의 12절까지만 다루십시오."""
                data = get_ai_response(
                    build_research_prompt(task, scripture, topic, theology),
                    is_json=True, temperature=0.2, kind="versions", max_tokens=16000)
                st.session_state[key] = data if isinstance(data, dict) else {}
            st.rerun()

        d = st.session_state.get(key) or {}
        if not d.get("rows"):
            if st.session_state.get("ai_fallback_used"):
                show_ai_status()
            else:
                st.caption("위 버튼을 눌러 번역본 비교표를 만들어 보세요.")
            return

        show_ai_status()
        verses = [str(v) for v in d.get("verses", [])]
        rows = d.get("rows", [])

        # ── 표 (행 = 번역본, 열 = 절)
        table = []
        for r in rows:
            row = {"번역본": r.get("version", "")}
            for i, vno in enumerate(verses):
                txts = r.get("texts", [])
                row[f"{vno}절"] = txts[i] if i < len(txts) else ""
            table.append(row)
        st.markdown("#### 📊 절별 비교표")
        st.caption("표 오른쪽 끝으로 스크롤하면 나머지 절이 보입니다. 셀을 클릭하면 전체 문장이 펼쳐집니다.")
        try:
            st.dataframe(table, width='stretch', hide_index=True)
        except Exception:
            st.dataframe(table, use_container_width=True, hide_index=True)

        # ── 절별로 읽기 좋은 형태
        st.markdown("#### 📖 절별로 나란히 읽기")
        for i, vno in enumerate(verses):
            with st.expander(f"{scripture} · {vno}절", expanded=(i == 0)):
                for r in rows:
                    txts = r.get("texts", [])
                    body = txts[i] if i < len(txts) else ""
                    if not str(body).strip():
                        continue
                    st.markdown(
                        f"<div style='margin-bottom:9px;'>"
                        f"<span style='display:inline-block;min-width:150px;color:#fde047;"
                        f"font-weight:800;font-size:13px;'>{_esc(r.get('version',''))}</span>"
                        f"<span style='color:#e8ecff;font-size:14.5px;line-height:1.75;'>"
                        f"{_esc(str(body))}</span></div>", unsafe_allow_html=True)

        # ── 문서 내려받기용 텍스트
        lines = [f"📚 성경 여러 번역본 비교 · {scripture}", ""]
        for i, vno in enumerate(verses):
            lines.append(f"{i+1}. {scripture} {vno}절")
            n = 0
            for r in rows:
                txts = r.get("texts", [])
                body = txts[i] if i < len(txts) else ""
                if str(body).strip():
                    n += 1
                    lines.append(f"- {n}. [{r.get('version','')}] {body}")
            lines.append("")
        doc_text = "\n".join(lines)

        st.write("")
        render_section_top_toolbar(f"{scripture}_번역본비교", doc_text,
                                   f"ver_{abs(hash(scripture)) % 9999}")

        # 엑셀로도 받을 수 있게 CSV
        try:
            import csv as _csv
            buf = io.StringIO()
            w = _csv.DictWriter(buf, fieldnames=list(table[0].keys()))
            w.writeheader()
            w.writerows(table)
            st.download_button("📥 표 그대로 내려받기 (CSV · 엑셀에서 열기)",
                               data=("﻿" + buf.getvalue()).encode("utf-8"),
                               file_name=f"{scripture}_번역본비교.csv", mime="text/csv",
                               key=f"dl_vercsv_{abs(hash(scripture)) % 9999}")
        except Exception:
            pass


def render_praise_section(scripture: str, topic: str, theology: str):
    """추천 찬양 15곡 + 곡마다 인도용 멘트"""
    key = _rkey("praise15", scripture)
    with st.expander("🎵 추천 찬양 15곡 — 곡마다 인도용 멘트 포함", expanded=False):
        if st.button("🎶 찬양 15곡 + 인도 멘트 생성", type="primary",
                     key=f"btn_praise_{abs(hash(scripture)) % 99999}"):
            with st.spinner("본문에 맞는 찬양을 고르고 인도 멘트를 쓰는 중..."):
                task = """[출력 형식 — 아래 JSON 하나만. 설명 문장 금지]
{
 "hymns": [{"title":"새찬송가 000장 - 제목","ment":"인도자가 이 곡 직전에 할 멘트 2문장"}],
 "gospel_songs": [{"title":"복음성가 제목 - 가수/작곡가","ment":"인도 멘트 2문장"}],
 "ccm": [{"title":"CCM 제목 - 아티스트","ment":"인도 멘트 2문장"}]
}
각 갈래 정확히 5곡씩, 총 15곡.
- 곡은 이 본문·주제의 정서와 실제로 맞아야 합니다.
- 새찬송가는 장수를 정확히 쓰십시오. 확실하지 않으면 장수를 빼고 제목만 쓰십시오.
- 인도 멘트는 회중에게 말하듯 구어체로, 본문 내용을 한 줄 얹어 자연스럽게 곡으로 넘어가게 쓰십시오.
- 멘트는 곡마다 서로 달라야 합니다. 같은 문장을 반복하지 마십시오."""
                st.session_state[key] = get_ai_response(
                    build_research_prompt(task, scripture, topic, theology),
                    is_json=True, temperature=0.55, kind="praise", max_tokens=10000)
            st.rerun()

        data = st.session_state.get(key)
        if not data:
            st.caption("위 버튼을 눌러 찬양을 추천받으세요.")
            return

        show_ai_status()
        lines = [f"🎵 추천 찬양 15곡 · {scripture}", ""]
        for label, k in (("📖 새찬송가", "hymns"), ("🕊️ 복음성가", "gospel_songs"), ("🎸 현대 CCM", "ccm")):
            items = data.get(k, []) or []
            st.markdown(f"#### {label}")
            lines.append(label)
            for i, it in enumerate(items, start=1):
                if isinstance(it, str):
                    it = {"title": it, "ment": ""}
                t = str(it.get("title", ""))
                m = str(it.get("ment", ""))
                q = urllib.parse.quote(t)
                st.markdown(
                    f"<div class='lib-card' style='padding:12px 16px;margin-bottom:8px;'>"
                    f"<div style='font-weight:800;color:#fde047;'>{i}. {_esc(t)}</div>"
                    f"<div class='leader-block' style='margin-top:8px;'>"
                    f"<b>💡 인도 멘트</b><br>{_esc(m)}</div>"
                    f"<div style='margin-top:6px;font-size:12.5px;'>"
                    f"<a href='https://www.youtube.com/results?search_query={q}' target='_blank'>▶️ 듣기</a>　"
                    f"<a href='https://www.google.com/search?q={q}' target='_blank'>🔍 악보·가사</a>"
                    f"</div></div>", unsafe_allow_html=True)
                lines.append(f"- {i}. {t}")
                if m:
                    lines.append(f"   [인도자 팁 / 가이드]: {m}")
            lines.append("")
        render_section_top_toolbar(f"{scripture}_추천찬양15곡", "\n".join(lines),
                                   f"praise_{abs(hash(scripture)) % 9999}")


def editable_section(state_key: str, session_field: str, label: str, height: int = 350,
                     persist_summary: bool = False):
    """수정 모드 공통 처리"""
    val = st.session_state.get(session_field, "")
    if st.session_state.get(f"edit_mode_{state_key}", False):
        edited = st.text_area(label, value=val, height=height, key=f"ta_{state_key}")
        if st.button("💾 저장", key=f"save_{state_key}"):
            st.session_state[session_field] = edited
            st.session_state[f"edit_mode_{state_key}"] = False
            if persist_summary:
                update_sermon_in_db(st.session_state.get("current_sermon_id", 1), updated_summary=edited)
            st.success("저장되었습니다.")
            st.rerun()
        return False
    return True


# ==============================================================================
# 세션 초기화
# ==============================================================================
if "sermon_library" not in st.session_state:
    st.session_state.sermon_library = get_db_sermons()
if "current_sermon_idx" not in st.session_state:
    st.session_state.current_sermon_idx = 0

_cur = st.session_state.sermon_library[st.session_state.current_sermon_idx] \
    if st.session_state.sermon_library else {}

st.session_state.setdefault("current_sermon_id", _cur.get("id", 1))
st.session_state.setdefault("sermon_title", _cur.get("title", "예배와 선교"))
st.session_state.setdefault("sermon_scripture", _cur.get("scripture", "이사야 59:21"))
st.session_state.setdefault("full_sermon", _cur.get("text", ""))
st.session_state.setdefault("sermon_summary_text", _cur.get("summary", ""))
st.session_state.setdefault("preacher_name", "김세훈목사")
st.session_state.setdefault("dash_active_view", "설교 요약")
st.session_state.setdefault("cn_church_name", "화광교회")
st.session_state.setdefault("ai_model_used", "")


# ==============================================================================
# 사이드바
# ==============================================================================
with st.sidebar.expander("⚙️ AI 연결 설정", expanded=not bool(get_resolved_api_key())):
    st.text_input("🔑 Gemini API Key", value=get_resolved_api_key(), type="password",
                  key="sidebar_api_key_input",
                  help="Google AI Studio(aistudio.google.com)에서 발급받은 키를 입력하세요.")
    if get_resolved_api_key():
        st.markdown("<span class='badge-ok'>키 등록됨</span>", unsafe_allow_html=True)

        _fp = hashlib.sha256(get_resolved_api_key().encode()).hexdigest()[:16]
        try:
            genai.configure(api_key=get_resolved_api_key())
        except Exception:
            pass
        _avail = discover_available_models(_fp)
        st.selectbox("사용할 모델", ["자동 선택 (권장)"] + _avail, key="ai_model_choice",
                     help="자동 선택이 안 될 때만 직접 고르세요. 목록이 이상하면 아래 [모델 목록 새로고침]을 누르세요.")

        cga, cgb = st.columns(2)
        with cga:
            if st.button("🔌 실제 작동 테스트", key="btn_test_api"):
                st.session_state["_dead_models"] = []
                try:
                    discover_available_models.clear()
                except Exception:
                    pass
                with st.spinner("모델마다 실제로 호출해 보는 중... (20초)"):
                    rows = []
                    for mname in discover_available_models(_fp)[:6]:
                        try:
                            mm = genai.GenerativeModel(mname)
                            r = mm.generate_content(
                                "한국어로 '연결됨' 이라고만 답하세요.",
                                generation_config={"temperature": 0, "max_output_tokens": 20})
                            ok = bool(getattr(r, "text", ""))
                            rows.append((mname, "✅ 사용 가능" if ok else "⚠️ 빈 응답"))
                        except Exception as e:
                            k = _classify_api_error(str(e))
                            label = {"quota": "⏳ 한도 초과", "gone": "❌ 폐기된 모델",
                                     "perm": "🔒 권한 없음", "key": "🔑 키 오류"}.get(k, "⚠️ 오류")
                            rows.append((mname, label))
                    st.session_state["_model_test_rows"] = rows
                st.rerun()
        with cgb:
            if st.button("🔄 목록 새로고침", key="btn_refresh_models"):
                st.session_state["_dead_models"] = []
                st.session_state.pop("_last_good_model", None)
                try:
                    discover_available_models.clear()
                except Exception:
                    pass
                st.rerun()

        if st.session_state.get("_model_test_rows"):
            st.markdown("**테스트 결과**")
            for mname, status in st.session_state["_model_test_rows"]:
                st.markdown(f"- `{mname}` → {status}")
            usable = [m for m, s in st.session_state["_model_test_rows"] if s.startswith("✅")]
            if usable:
                st.success(f"사용 가능한 모델: {usable[0]}")
            else:
                st.error("지금은 쓸 수 있는 모델이 없습니다. 대부분 ⏳ 한도 초과라면 "
                         "1~2분 뒤 다시 시도해 주세요.")
        if st.session_state.get("_last_good_model"):
            st.caption(f"최근 성공 모델: {st.session_state['_last_good_model']}")
    else:
        st.markdown("<span class='badge-bad'>키 없음 · AI 기능 제한</span>", unsafe_allow_html=True)

with st.sidebar.expander("🗄️ 설교 서재 저장소", expanded=not cloud_store_ready()):
    if cloud_store_ready():
        gid = _gist_find_id()
        if gid:
            st.markdown("<span class='badge-ok'>영구 보관 작동 중</span>", unsafe_allow_html=True)
            st.caption(f"GitHub 비공개 Gist에 저장됩니다.\n\n보관함 ID: `{gid[:10]}…`")
            st.link_button("🔗 보관함 열어보기", f"https://gist.github.com/{gid}")
        else:
            st.markdown("<span class='badge-ok'>토큰 확인됨</span>", unsafe_allow_html=True)
            st.caption("첫 설교를 등록하면 비공개 보관함이 자동으로 만들어집니다.")
        if st.session_state.get("_cloud_error"):
            st.error(f"클라우드 오류: {st.session_state['_cloud_error']}")
        if st.button("🔄 보관함에서 다시 불러오기", key="btn_cloud_reload"):
            get_db_sermons(force_reload=True)
            st.success("보관함과 동기화했습니다.")
            st.rerun()
    else:
        st.markdown("<span class='badge-bad'>임시 저장 · 사라질 수 있음</span>", unsafe_allow_html=True)
        st.error(
            "지금은 설교가 **서버 임시 폴더**에만 저장됩니다.\n\n"
            "Streamlit Cloud는 앱이 잠들거나 재배포될 때 이 폴더를 지웁니다. "
            "그래서 다음 접속 때 설교가 사라진 것입니다.\n\n"
            "**영구 보관 켜는 법 (5분)**\n"
            "1. github.com → 오른쪽 위 프로필 → Settings\n"
            "2. 맨 아래 Developer settings → Personal access tokens → **Tokens (classic)**\n"
            "3. Generate new token (classic) → Note에 `sermon`, Expiration은 **No expiration**\n"
            "4. 체크박스 중 **gist** 하나만 체크 → Generate token\n"
            "5. 나온 긴 문자열(ghp_… )을 복사\n"
            "6. Streamlit 앱 화면 → Manage app → ⋮ → Settings → **Secrets** 에 아래 한 줄 붙여넣고 Save\n"
            "```\nGITHUB_TOKEN = \"ghp_여기에붙여넣기\"\n```\n"
            "7. 앱이 자동 재시작되면 끝입니다. 이후 등록한 설교는 영원히 보관됩니다."
        )
        st.info("설정 전까지는 서재 화면의 **[💾 전체 백업]** 버튼으로 파일을 꼭 내려받아 두세요.")

if not KOREAN_FONT_OK:
    st.sidebar.error(
        "⚠️ 서버에 한글 폰트가 없어 PDF·카드 이미지의 한글이 깨질 수 있습니다.\n\n"
        "Streamlit Cloud라면 저장소 루트에 `packages.txt` 파일을 만들고 한 줄만 넣어 재배포하세요:\n\n"
        "```\nfonts-nanum\n```"
    )

st.sidebar.markdown(
    f"**📖 현재 작업 설교**\n\n{st.session_state.sermon_title}\n\n`{st.session_state.sermon_scripture}` "
    f"· 원고 {len(st.session_state.full_sermon):,}자"
)

app_mode = st.sidebar.radio(
    "🕊️ 플랫폼 대메뉴",
    ["📊 설교 대시보드 (메인 작업실)",
     "📤 새 설교 등록/원고작성",
     "🎙️ AI 보이스오버 스튜디오",
     "🎬 쇼츠 만들기 (스튜디오)",
     "📷 말씀카드 이미지",
     "📚 설교 서재 (Sermon Library)"]
)


# ==============================================================================
# 1) 설교 대시보드
# ==============================================================================
if app_mode == "📊 설교 대시보드 (메인 작업실)":
    st.markdown(
        f"""<div class="hero">
        <h1>{st.session_state.sermon_title}</h1>
        <div style="margin-top:10px;">
          <span class="chip chip-gold">📖 {st.session_state.sermon_scripture}</span>
          <span class="chip chip-vio">✍️ {st.session_state.get('preacher_name','')}</span>
          <span class="chip chip-mint">📝 원고 {len(st.session_state.full_sermon):,}자</span>
        </div></div>""",
        unsafe_allow_html=True
    )

    _an = analyze_manuscript(st.session_state.full_sermon)
    if _an["char_count"] < 80:
        st.error("설교 원고가 비어 있습니다. [📤 새 설교 등록/원고작성]에서 원고를 먼저 등록해 주세요.")
    else:
        st.caption(
            f"🔑 원고 자동 추출 키워드: {', '.join(_an['keywords'][:10]) or '(없음)'}　|　"
            f"📖 원고 인용 성구: {', '.join(_an['refs'][:6]) or '(없음)'}"
        )

    with st.expander("🖼️ 카드뉴스 · PPT 배경 이미지 설정 (무한 생성)", expanded=False):
        bg1, bg2 = st.columns([2, 1])
        with bg1:
            st.selectbox("이미지 분위기", list(BG_THEMES.keys()), key="bg_theme",
                         help="분위기를 바꾸면 카드뉴스·PPT·말씀카드 배경이 통째로 달라집니다.")
        with bg2:
            st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)
            if st.button("🔀 배경 전부 새로 뽑기", key="btn_bg_shuffle"):
                st.session_state.bg_shuffle = int(st.session_state.get("bg_shuffle", 0)) + 1
                st.rerun()
        st.caption(f"현재 배경 세트 #{st.session_state.get('bg_shuffle', 0)} · "
                   "누를 때마다 완전히 새로운 이미지 조합이 나옵니다. "
                   "사진을 받아오지 못하는 환경에서는 자동으로 고급 그라데이션 배경이 생성됩니다.")

    render_manna_section()
    render_today_word_section()

    st.write("---")
    left, right = st.columns([1, 2.5])

    MENU = ["설교 요약", "✍️ 주보 칼럼", "소그룹 나눔", "QT 5일치", "카드뉴스", "쇼츠 대본",
            "🏡 세대별 가정예배지", "🔍 설교 점검 및 제안", "📖 소그룹 리더가이드"]

    with left:
        st.markdown("<p style='font-size:12px;font-weight:bold;color:#94a3b8;'>사역 메뉴 선택</p>",
                    unsafe_allow_html=True)
        sel = st.radio("사역 메뉴", MENU,
                       index=MENU.index(st.session_state.dash_active_view)
                       if st.session_state.dash_active_view in MENU else 0,
                       key="dash_menu_selector", label_visibility="collapsed")
        st.session_state.dash_active_view = sel

    with right:
        view = st.session_state.dash_active_view

        # ---------------- 설교 요약 ----------------
        if view == "설교 요약":
            if not st.session_state.get("sermon_summary_text") or len(st.session_state.sermon_summary_text) < 50:
                st.session_state.sermon_summary_text = build_local_summary(
                    st.session_state.sermon_title, st.session_state.sermon_scripture,
                    st.session_state.full_sermon)

            summary_val = st.session_state.sermon_summary_text
            render_section_top_toolbar(f"{st.session_state.sermon_title}_설교요약",
                                       summary_val, "sermon_sum", ppt_mode="sermon")

            if st.button("✨ AI 강해 요약 생성 (원고 근거 인용 포함)", type="primary", key="btn_regen_ai_summary"):
                with st.spinner("원고를 읽고 대지·근거·적용을 추출하는 중..."):
                    task = """[출력 형식 - 아래 틀 그대로]
🎯 설교 핵심 명제
(이 원고 전체를 관통하는 중심 사상 1~2문장. 원고의 표현을 살려 쓸 것)

🔑 이 설교의 핵심 키워드
(원고에 실제로 반복 등장한 단어 6~8개, 쉼표 구분)

📜 본문 말씀 (인용)
(설교 본문 구절을 개역개정으로 1~3절 인용. 각 줄 맨 앞에 절 번호를 쓸 것)

📌 강해적 3대지
1. (대지 제목 — 원고에 나온 표현으로. 상투어 금지) [근거 성구: 원고에 인용된 구절]
   ▸ 원고 근거: "원고에서 그대로 가져온 문장 1~2개"
   ▸ 주해 설명: (그 대목이 말하는 바를 2~3문장)
2. (동일 형식)
3. (동일 형식)

💡 성도를 위한 실천 적용 3가지
1. (원고의 권면을 구체적 행동으로 — 원고에 없는 일반론 금지)
2.
3.

🙏 결단 및 축복 기도문
(원고의 결론·기도 대목을 반영한 3~4문장)"""
                    res = get_ai_response(build_grounded_prompt(task, ctx_chars=11000), is_json=False,
                                          temperature=0.25, kind="summary")
                    st.session_state.sermon_summary_text = res
                    update_sermon_in_db(st.session_state.get("current_sermon_id", 1), updated_summary=res)
                    st.rerun()

            show_ai_status()
            if editable_section("sermon_sum", "sermon_summary_text", "설교 요약 편집",
                                height=400, persist_summary=True):
                render_body(st.session_state.sermon_summary_text)

            with st.expander("📜 설교문 원고 전문 보기", expanded=False):
                st.text(st.session_state.full_sermon)

        # ---------------- 주보 칼럼 ----------------
        elif view == "✍️ 주보 칼럼":
            txt = st.session_state.get("bulletin_column_text", "")
            render_section_top_toolbar(f"{st.session_state.sermon_title}_주보칼럼", txt, "bulletin")

            st.caption("설교 요약문을 레퍼런스로 삼아, 주보에 그대로 실을 수 있는 목회 칼럼을 씁니다.")
            bc1, bc2, bc3 = st.columns([1.2, 1, 1])
            with bc1:
                col_tone = st.selectbox("칼럼 어조",
                                        ["따뜻한 목회 서신체", "차분한 묵상 에세이", "권면·도전형", "간증·이야기형"],
                                        key="bc_tone")
            with bc2:
                col_len = st.selectbox("분량", ["700자 내외", "900자 내외", "1,200자 내외"], key="bc_len")
            with bc3:
                col_sign = st.text_input("서명", value=st.session_state.get("preacher_name", ""), key="bc_sign")

            if st.button("✍️ 주보 칼럼 생성", type="primary", key="btn_gen_bulletin"):
                with st.spinner("주보 칼럼을 집필하는 중..."):
                    summary_ref = st.session_state.get("sermon_summary_text", "") or \
                        build_local_summary(st.session_state.sermon_title,
                                            st.session_state.sermon_scripture,
                                            st.session_state.full_sermon)
                    task = f"""[참고 자료 — 이번 주 설교 요약문]
{summary_ref[:2500]}

[작업]
위 설교 요약문과 <설교원고>를 레퍼런스로, 교회 주보에 실을 목회 칼럼을 씁니다.

[출력 형식 - 아래 틀 그대로]
✍️ 「(칼럼 제목 — 설교 제목을 그대로 베끼지 말고, 칼럼다운 제목으로 12자 내외)」

(본문 — {col_len}, {col_tone}. 문단 4~5개.
 1문단: 일상의 한 장면이나 짧은 질문으로 시작해 독자를 끌어들입니다.
 2문단: 이번 주 본문({st.session_state.sermon_scripture})이 말하는 바를 쉽게 풀어 씁니다.
 3문단: 설교의 핵심 메시지를 성도의 삶에 붙입니다.
 4문단: 한 주간의 구체적인 권면 한 가지.
 5문단: 짧은 축복의 문장으로 맺습니다.)

📖 이번 주 말씀 — {st.session_state.sermon_scripture}

— {col_sign or '담임목사'} 드림

[반드시 지킬 것]
- 설교 원고와 요약문에 실제로 있는 내용만 씁니다. 새 예화나 통계를 지어내지 마십시오.
- 설교문을 그대로 옮기지 말고, 읽는 글(칼럼)의 호흡으로 다시 쓰십시오.
- 소제목·번호·불릿을 쓰지 말고, 이어지는 문단 산문으로만 쓰십시오.
- 강단 어투('~하시기 바랍니다'의 반복)를 줄이고 편지처럼 담백하게 쓰십시오."""
                    st.session_state.bulletin_column_text = get_ai_response(
                        build_grounded_prompt(task, ctx_chars=7000), is_json=False,
                        kind="bulletin", temperature=0.6)
                    st.rerun()

            if txt:
                show_ai_status()
                st.caption(f"현재 분량: 약 {len(txt.replace(chr(10), '')):,}자")
                if editable_section("bulletin", "bulletin_column_text", "주보 칼럼 편집", height=420):
                    render_body(st.session_state.bulletin_column_text)
            else:
                st.caption("위 버튼을 눌러 주보 칼럼을 생성하세요. (설교 요약문이 자동으로 참고 자료가 됩니다)")

        # ---------------- 소그룹 나눔 ----------------
        elif view == "소그룹 나눔":
            txt = st.session_state.get("small_group_text", "")
            render_section_top_toolbar(f"{st.session_state.sermon_title}_소그룹나눔지", txt, "sm_grp")

            if st.button("✨ 소그룹 나눔지 생성", type="primary", key="btn_gen_sm_grp"):
                with st.spinner("소그룹 나눔지 작성 중..."):
                    task = """[출력 형식]
1. 마음 열기
- [인도자 팁 / 가이드]: (이 설교 주제와 이어지는 도입 멘트)
- 질문 1개 (이 설교의 키워드와 연결될 것)

2. 말씀 속으로
- [인도자 팁 / 가이드]: ...
- 1. (원고의 특정 대목을 직접 인용하며 묻는 질문)
- 2. (원고의 두 번째 대목을 인용하며 묻는 질문)

3. 삶 속으로
- [인도자 팁 / 가이드]: ...
- 1. (원고의 권면을 구체적 상황으로 바꾼 질문)
- 2. (이번 주 결단 질문)

4. 마침 합심 기도문 (원고 내용 반영, 3문장)

※ 모든 질문에는 이 설교 원고에만 나오는 표현이 최소 1개 포함되어야 합니다.
※ 번호는 각 항목(1.마음 열기 / 2.말씀 속으로 / 3.삶 속으로 / 4.기도) 안에서 매번 1번부터 다시 시작합니다.
※ [인도자 팁 / 가이드] 는 각 항목마다 반드시 한 줄 넣으십시오."""
                    st.session_state.small_group_text = get_ai_response(
                        build_grounded_prompt(task), is_json=False, kind="smallgroup")
                    st.rerun()

            if txt:
                show_ai_status()
                if editable_section("sm_grp", "small_group_text", "소그룹 나눔지 편집"):
                    render_body(st.session_state.small_group_text)
            else:
                st.caption("위 버튼을 눌러 소그룹 나눔지를 생성하세요.")

        # ---------------- QT 5일치 ----------------
        elif view == "QT 5일치":
            txt = st.session_state.get("qt5_text", "")
            render_section_top_toolbar(f"{st.session_state.sermon_title}_주간QT5일치", txt, "qt5")

            if st.button("✨ 5일치 QT 묵상지 생성", type="primary", key="btn_gen_qt5"):
                with st.spinner("주간 5일치 QT 작성 중..."):
                    task = """[출력 형식 — 월~금 5일, 각 날짜마다 아래 틀 그대로]
📅 (요일) · (그날의 소제목 — 원고의 서로 다른 대목에서 뽑을 것)

📖 본문 구절
(원고에 인용된 구절 중 하나, 또는 대표 성구 — 장절과 구절 전문)

💡 말씀 묵상
(원고의 해당 대목을 풀어 4~5문장)

❓ 묵상 질문
- 1. (본문 관찰 질문 — 본문이 무엇을 말하는가)
- 2. (내면을 들여다보는 질문 — 나의 삶에서 이 말씀은 어디에 닿는가)
- 3. (결단을 부르는 질문 — 오늘 무엇을 바꿀 것인가)

🎯 삶의 적용
(오늘 실행할 구체적 행동 1가지)

🙏 오늘의 기도
(2~3문장)

※ 5일이 서로 다른 내용이어야 합니다. 같은 말을 다섯 번 반복하지 마십시오.
※ 묵상 질문 3개는 날마다 새로 만들되, 관찰 → 성찰 → 결단 순서를 지키십시오.
※ 각 날짜의 묵상 질문 번호는 항상 1번부터 다시 시작합니다."""
                    st.session_state.qt5_text = get_ai_response(
                        build_grounded_prompt(task), is_json=False, kind="qt")
                    st.rerun()

            if txt:
                show_ai_status()
                if editable_section("qt5", "qt5_text", "QT 5일치 편집"):
                    render_body(st.session_state.qt5_text)
            else:
                st.caption("위 버튼을 눌러 5일치 QT를 생성하세요.")

        # ---------------- 카드뉴스 ----------------
        elif view == "카드뉴스":
            h1, h2 = st.columns([1.2, 1.8])
            with h1:
                st.markdown("<h2 style='margin:0;font-size:24px;font-weight:bold;'>카드뉴스</h2>",
                            unsafe_allow_html=True)
            with h2:
                e1, e2, e3 = st.columns([1, 1.3, 1.4])
                with e1:
                    if st.button("✏️ 편집", key="cn_edit_toggle_btn"):
                        st.session_state.cn_edit_mode = not st.session_state.get("cn_edit_mode", False)
                if st.session_state.get("card_list"):
                    cj = json.dumps(st.session_state.card_list, ensure_ascii=False)
                    with e2:
                        st.download_button("📥 PPT 전체",
                                           data=generate_cardnews_pptx_bytes(cj, st.session_state.cn_church_name,
                                                                             bg_seed(0), current_bg_theme(),
                                                                             st.session_state.sermon_scripture),
                                           file_name=f"{st.session_state.sermon_title}_카드뉴스.pptx",
                                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                           key="cn_dl_ppt")
                    with e3:
                        st.download_button("📦 전체 PNG",
                                           data=generate_cardnews_zip_bytes(cj, st.session_state.sermon_scripture,
                                                                            st.session_state.cn_church_name,
                                                                            bg_seed(0), current_bg_theme()),
                                           file_name=f"{st.session_state.sermon_title}_카드뉴스.zip",
                                           mime="application/zip", key="cn_dl_zip")

            o1, o2 = st.columns([1, 1])
            with o1:
                card_count = st.slider("카드 총 장수", 5, 12, 7, key="cn_count_slider")
            with o2:
                st.session_state.cn_church_name = st.text_input(
                    "교회명", value=st.session_state.cn_church_name, key="cn_church_input")

            if st.button(f"🎨 카드뉴스 {card_count}장 생성 / 다시 생성", type="primary", key="btn_gen_cardnews"):
                with st.spinner(f"원고에서 {card_count}장 카드뉴스 구성 중..."):
                    task = f"""[출력 형식 - JSON만]
{{"cards": [{{"card_number": 1, "headline": "...", "body_text": "..."}}, ...]}}

정확히 {card_count}장을 만드십시오. 구성 규칙:
- 1장: 표제 카드. headline은 설교 제목, body_text는 이 설교의 한 줄 요지(원고 문장 기반).
- 2장~{card_count-2}장: 원고의 서로 다른 대목을 하나씩 담습니다.
  headline은 그 대목의 핵심을 12~20자로 (원고 단어 사용), body_text는 60~110자로 원고 내용을 풀어 씁니다.
- {card_count-1}장: 💡 삶의 적용 — 원고의 권면 3가지를 번호로.
- {card_count}장: 🙏 기도 — 원고의 결론/기도를 반영한 2~3문장.

절대 금지: "말씀의 깊은 은혜", "믿음의 결단", "주님의 신실하신 은혜가 충만하기를" 같이
어느 설교에나 붙는 문구로 카드를 채우는 것. 모든 카드에 이 원고 고유의 단어가 들어가야 합니다."""
                    res = get_ai_response(build_grounded_prompt(task, ctx_chars=9000), is_json=True,
                                          temperature=0.3, kind="cards", card_count=card_count)
                    cards = (res or {}).get("cards", [])
                    if cards:
                        for i, c in enumerate(cards, start=1):
                            c["card_number"] = i
                        st.session_state.card_list = cards[:card_count]
                        st.session_state.cn_card_idx = 0
                    st.rerun()

            show_ai_status()

            if st.session_state.get("cn_edit_mode") and st.session_state.get("card_list"):
                st.markdown("#### ✏️ 카드 텍스트 편집기")
                for i, c in enumerate(st.session_state.card_list):
                    with st.expander(f"CARD {c.get('card_number', i+1)}", expanded=(i == 0)):
                        c["headline"] = st.text_input("헤드라인", value=c.get("headline", ""), key=f"cn_h_{i}")
                        c["body_text"] = st.text_area("본문", value=c.get("body_text", ""), key=f"cn_b_{i}")

            if st.session_state.get("card_list"):
                cards = st.session_state.card_list
                total = len(cards)
                st.session_state.setdefault("cn_card_idx", 0)
                idx = st.session_state.cn_card_idx % total

                n1, n2, n3 = st.columns([1, 4, 1])
                with n1:
                    if st.button("❮ 이전", key="cn_prev"):
                        st.session_state.cn_card_idx = (idx - 1) % total
                        st.rerun()
                with n3:
                    if st.button("다음 ❯", key="cn_next"):
                        st.session_state.cn_card_idx = (idx + 1) % total
                        st.rerun()
                with n2:
                    png = generate_single_card_png_bytes(
                        json.dumps(cards[idx], ensure_ascii=False), idx,
                        st.session_state.sermon_scripture, st.session_state.cn_church_name,
                        f"{bg_seed(0)}|{idx}", current_bg_theme())
                    st_image_full(png, caption=f"{idx+1} / {total} — 실제 다운로드 결과와 동일")
                    st.download_button(f"🖼️ CARD {idx+1} PNG 다운로드", data=png,
                                       file_name=f"{st.session_state.sermon_title}_card_{idx+1}.png",
                                       mime="image/png", key=f"dl_card_{idx}")

                st.write("---")
                st.markdown("#### 인스타그램 캡션")
                _a = analyze_manuscript(st.session_state.full_sermon)
                cap = (f"[{st.session_state.sermon_title}] ({st.session_state.sermon_scripture})\n\n"
                       f"{(_a['key_sentences'][0] if _a['key_sentences'] else '')}\n\n"
                       + " ".join(["#주일설교", "#말씀묵상"] + [f"#{k}" for k in _a['keywords'][:6]]))
                st.code(cap, language="text")
            else:
                st.caption("위 버튼을 눌러 카드뉴스를 생성하세요.")

        # ---------------- 쇼츠 대본 ----------------
        elif view == "쇼츠 대본":
            txt = st.session_state.get("shorts_script_text", "")
            render_section_top_toolbar(f"{st.session_state.sermon_title}_쇼츠대본", txt, "sh_script")

            if st.button("🎬 쇼츠 대본 3종 생성", type="primary", key="btn_gen_shorts_script"):
                with st.spinner("쇼츠 대본 작성 중..."):
                    task = """[출력 형식 — 60초 세로 쇼츠 대본 3종]
각 대본마다:
🎬 (유형명)
- 후킹(0~5초): (원고의 가장 강력한 문장에서 따온 한 마디)
- 본론(5~45초): (원고 내용 기반 3~4문장, 화면 자막 단위로 줄바꿈)
- 결단(45~60초): (한 문장 권면)
- 화면 자막 키워드: (5개)

3종 유형: 1) 감동·위로형  2) 질문·호기심형  3) 결단 선포형
세 대본은 원고의 서로 다른 대목을 사용해야 합니다."""
                    st.session_state.shorts_script_text = get_ai_response(
                        build_grounded_prompt(task), is_json=False, kind="shorts_script")
                    st.rerun()

            if txt:
                show_ai_status()
                if editable_section("sh_script", "shorts_script_text", "쇼츠 대본 편집"):
                    render_body(st.session_state.shorts_script_text)
            else:
                st.caption("위 버튼을 눌러 쇼츠 대본을 생성하세요.")

        # ---------------- 가정예배지 ----------------
        elif view == "🏡 세대별 가정예배지":
            age = st.selectbox("예배 대상 선택", ["👶 영유아용", "🧒 어린이용", "🧑 청소년용", "👨‍👩‍👧 청장년용"],
                               key="sel_age_group")
            fkey = f"family_worship_{age}"
            txt = st.session_state.get(fkey, "")
            render_section_top_toolbar(f"{st.session_state.sermon_title}_가정예배지_{age}", txt, f"fam_{age}")

            if st.button(f"✨ {age} 맞춤 가정예배지 생성", type="primary", key="btn_gen_fam"):
                with st.spinner(f"{age} 가정예배지 작성 중..."):
                    task = f"""[출력 형식 — 대상: {age}]
1. 찬양 및 신앙고백
- [인도자 팁 / 가이드]: (이 설교 주제에 맞는 찬송 제안과 시작 멘트)
2. 함께 읽는 성경 말씀
- [인도자 팁 / 가이드]: ...
3. {age} 눈높이 3분 가족 메시지
- (원고의 핵심 내용을 이 연령대 언어로. 원고의 예화·표현을 반드시 활용)
4. 온 가족 나눔 질문 2가지
- [인도자 팁 / 가이드]: ...
5. 가정을 축복하는 마무리 기도문

※ {age}의 이해 수준에 맞춘 어휘를 쓰되, 내용은 반드시 이 설교 원고에서 나와야 합니다.
※ 다섯 항목 각각에 [인도자 팁 / 가이드] 를 한 줄씩 반드시 넣으십시오.
※ 번호는 각 항목 안에서 1번부터 다시 시작합니다(나눔 질문 1·2 / 기도제목 1·2 형태)."""
                    st.session_state[fkey] = get_ai_response(
                        build_grounded_prompt(task), is_json=False, kind="family")
                    st.rerun()

            if txt:
                show_ai_status()
                if editable_section(f"fam_{age}", fkey, "가정예배지 편집", height=320):
                    render_body(st.session_state[fkey])
            else:
                st.caption(f"위 버튼을 눌러 {age} 맞춤 가정예배지를 생성하세요.")

        # ---------------- 설교 점검 ----------------
        elif view == "🔍 설교 점검 및 제안":
            txt = st.session_state.get("sermon_audit_text", "")
            render_section_top_toolbar(f"{st.session_state.sermon_title}_설교점검및제안", txt, "sermon_audit")

            if st.button("🔍 설교 정밀 점검", type="primary", key="btn_gen_audit"):
                with st.spinner("설교 분석 리포트 작성 중..."):
                    task = """[출력 형식 — 냉정하고 구체적으로. 칭찬만 나열하지 말 것]
1. 📖 본문 주해의 정확성 (점수/100 + 근거)
   - 잘된 점: 원고 문장 인용 + 이유
   - 아쉬운 점: 원고 문장 인용 + 구체적 수정 제안
2. 🏗️ 대지 전개와 구조 (점수/100 + 근거)
   - 실제 대지 구조를 요약한 뒤, 논리 비약이나 중복이 있으면 지적
3. 💡 예화·적용의 적절성 (점수/100)
   - 원고의 예화를 열거하고, 본문과의 연결 강도를 평가
4. 🎙️ 전달력 개선 제안 3가지 (원고의 특정 문장을 지목해서)
5. 📊 종합 총평 + 다음 설교를 위한 권고 3가지

※ 원고에 실제로 없는 내용을 지적하거나 칭찬하지 마십시오."""
                    st.session_state.sermon_audit_text = get_ai_response(
                        build_grounded_prompt(task, ctx_chars=12000), is_json=False,
                        temperature=0.3, kind="audit")
                    st.rerun()

            if txt:
                show_ai_status()
                if editable_section("sermon_audit", "sermon_audit_text", "설교 점검 편집"):
                    render_body(st.session_state.sermon_audit_text)
            else:
                st.caption("위 버튼을 눌러 설교 점검 리포트를 생성하세요.")

        # ---------------- 리더가이드 ----------------
        elif view == "📖 소그룹 리더가이드":
            txt = st.session_state.get("leader_guide_text", "")
            render_section_top_toolbar(f"{st.session_state.sermon_title}_소그룹리더가이드", txt, "ldr_guide")

            if st.button("📖 리더가이드 생성", type="primary", key="btn_gen_ldr_guide"):
                with st.spinner("소그룹 인도자 가이드 작성 중..."):
                    task = """[출력 형식]
1. 🎯 이번 주 모임의 핵심 목표
- [인도자 팁 / 가이드]: ...
2. 📖 본문 배경 및 신학적 핵심 해설 (리더 심화용)
- [인도자 팁 / 가이드]: (성도가 던질 만한 신학적 질문 2개와 답변 방향)
3. 💬 나눔 질문별 예상 답변 & 피드백 요령
- [인도자 팁 / 가이드]: ...
4. ⚠️ 침묵·돌발 상황 대처
- [인도자 팁 / 가이드]: ...
5. 🙏 소그룹 맞춤 중보기도 제목 3가지

※ 2번 해설은 반드시 이 설교의 본문과 원고 내용에 근거해야 합니다.
※ 각 항목의 하위 번호는 항목마다 1번부터 다시 시작합니다.
※ [인도자 팁 / 가이드] 는 항목마다 한 줄씩 반드시 넣으십시오."""
                    st.session_state.leader_guide_text = get_ai_response(
                        build_grounded_prompt(task), is_json=False, kind="leader")
                    st.rerun()

            if txt:
                show_ai_status()
                if editable_section("ldr_guide", "leader_guide_text", "리더가이드 편집"):
                    render_body(st.session_state.leader_guide_text)
            else:
                st.caption("위 버튼을 눌러 소그룹 리더가이드를 생성하세요.")


# ==============================================================================
# 2) 새 설교 등록 / 원고작성
# ==============================================================================
elif app_mode == "📤 새 설교 등록/원고작성":
    st.markdown("<h1 style='font-size:28px;font-weight:800;'>📤 새 설교 등록 및 원고 작성</h1>",
                unsafe_allow_html=True)
    st.caption("직접 타이핑 · 파일 업로드 · AI 강해설교문 자동 생성으로 새 설교를 등록합니다.")

    t1, t2, t3 = st.tabs(["✍️ 직접 타이핑 작성",
                          "📁 파일 업로드 (.docx/.pdf/.txt)",
                          "📖 AI 강해설교문 생성"])

    with t1:
        c1, c2, c3 = st.columns([2, 1.5, 1])
        with c1:
            t_title = st.text_input("설교 제목", placeholder="예: 광야에서 만나는 하나님의 은혜", key="type_title")
        with c2:
            t_scrip = st.text_input("성경 본문", placeholder="예: 출애굽기 16:1-12", key="type_scrip")
        with c3:
            st.text_input("설교자", value=st.session_state.preacher_name, key="type_preach")

        t_tags = st.text_input("태그 (쉼표 구분)", value="주일설교", key="type_tags")
        t_content = st.text_area("설교문 본문 전문", height=380, key="type_text")
        st.caption(f"글자 수: **{len(t_content):,}자**")

        if st.button("💾 새 설교로 등록", type="primary", key="save_type_sermon"):
            if not t_title.strip() or len(t_content.strip()) < 50:
                st.warning("설교 제목과 본문(50자 이상)을 입력해주세요.")
            else:
                testament, book = classify_scripture(t_scrip.strip())
                entry = {"title": t_title.strip(), "scripture": t_scrip.strip(),
                         "testament": testament, "book": book,
                         "topic": (t_tags.split(",")[0].strip() if t_tags else "일반설교"),
                         "theology": "직접작성", "date": datetime.now().strftime("%Y-%m-%d"),
                         "tags": [x.strip() for x in t_tags.split(",") if x.strip()],
                         "summary": "", "text": t_content.strip()}
                saved = add_sermon_to_db(entry)
                load_sermon_to_workspace(saved, idx=len(st.session_state.sermon_library) - 1)
                st.success(f"'{t_title}' 등록 완료! 사이드바에서 [📊 설교 대시보드]로 이동하세요.")
                st.rerun()

    with t2:
        u_file = st.file_uploader("설교 파일 선택", type=["docx", "pdf", "txt"], key="up_sermon_file")
        f_title = st.text_input("설교 제목", value="", placeholder="업로드 설교 제목", key="up_title")
        f_scrip = st.text_input("성경 본문", value="", placeholder="예: 로마서 8:28-39", key="up_scrip")

        if u_file and st.button("📂 파일 읽어와서 등록", type="primary", key="save_up_sermon"):
            text = ""
            fn = u_file.name.lower()
            try:
                if fn.endswith('.txt'):
                    text = u_file.read().decode('utf-8', errors='ignore')
                elif fn.endswith('.docx'):
                    d = Document(u_file)
                    text = "\n".join(p.text for p in d.paragraphs if p.text.strip())
                elif fn.endswith('.pdf'):
                    pdf = PdfReader(u_file)
                    text = "\n".join((pg.extract_text() or "") for pg in pdf.pages)
            except Exception as e:
                st.error(f"파일 읽기 실패: {e}")
                text = ""

            if len(text.strip()) < 50:
                st.warning("본문을 읽지 못했습니다. 파일을 확인해 주세요.")
            else:
                title = f_title.strip() or os.path.splitext(u_file.name)[0]
                scrip = f_scrip.strip() or (extract_scripture_refs(text)[:1] or ["본문 미지정"])[0]
                testament, book = classify_scripture(scrip)
                entry = {"title": title, "scripture": scrip, "testament": testament, "book": book,
                         "topic": "파일등록", "theology": "파일업로드",
                         "date": datetime.now().strftime("%Y-%m-%d"), "tags": ["파일등록"],
                         "summary": "", "text": text.strip()}
                saved = add_sermon_to_db(entry)
                load_sermon_to_workspace(saved, idx=len(st.session_state.sermon_library) - 1)
                st.success(f"등록 완료! (본문 자동 인식: {scrip}) [📊 설교 대시보드]로 이동하세요.")
                st.rerun()

    with t3:
        st.markdown(
            "<div class='hero' style='padding:16px 22px;'><h1 style='font-size:22px;'>"
            "📖 본문 연구실 &amp; 강해설교문 작성</h1>"
            "<div style='margin-top:8px;color:#c3cdf5;font-size:13.5px;'>"
            "본문·주제를 정하면 아래 <b>연구 도구</b>들이 성경 66권 전체·역사·문학·현대 자료에서 "
            "설교 준비 자료를 만들어 옵니다.</div></div>",
            unsafe_allow_html=True)

        a1, a2, a3 = st.columns([1.2, 1.5, 1.3])
        with a1:
            sel_book = st.selectbox("성경 66권", BIBLE_BOOKS,
                                    index=BIBLE_BOOKS.index("시편"), key="sel_ai_book")
        with a2:
            sel_cv = st.text_input("장·절", value="1편 1-6절", key="sel_ai_cv")
        with a3:
            theology = st.selectbox("신학적 관점", list(THEOLOGY_LENSES.keys()), key="sel_ai_theology")
        st.caption(f"🔎 {THEOLOGY_LENSES[theology]['desc']}")

        b1, b2, b3, b4 = st.columns([2, 1, 1, 1])
        with b1:
            topic = st.text_input("설교 주제 / 강조 포인트 (비워두면 본문이 말하는 주제를 AI가 잡습니다)",
                                  value="", placeholder="예: 고난 중에도 끊을 수 없는 하나님의 사랑",
                                  key="sel_ai_topic")
        with b2:
            outline_key = st.selectbox("대지 구조", list(OUTLINE_SHAPES.keys()), index=2, key="sel_ai_outline")
        with b3:
            length_key = st.selectbox("설교 분량", list(SERMON_LENGTHS.keys()), index=5, key="sel_ai_length")
        with b4:
            style = st.selectbox("설교 형태", ["본문중심 강해설교", "절별 주해설교",
                                            "구속사적 복음설교", "원어 주해 중심 강해설교",
                                            "주제(토픽) 설교"],
                                 key="sel_ai_style")

        full_scrip = f"{sel_book} {sel_cv}"
        lens = THEOLOGY_LENSES[theology]
        shape = OUTLINE_SHAPES[outline_key]
        length = SERMON_LENGTHS[length_key]

        # ------------------------------------------------------------------
        # 본문 연구 도구 7종
        # ------------------------------------------------------------------
        render_research_tools(full_scrip, topic, theology)

        st.write("---")
        st.markdown("### ✍️ 강해설교문 전문 작성")

        if st.button(f"🚀 강해설교문 전문 작성 ({length_key})", type="primary", key="btn_gen_ai_sermon"):
            if not get_resolved_api_key():
                st.error("이 기능은 AI 생성 전용입니다. 사이드바 [⚙️ AI 연결 설정]에서 Gemini API 키를 먼저 등록해 주세요.")
            else:
                with st.spinner(f"[{theology}] · {outline_key} · {length_key} 분량으로 집필 중... (1~2분 소요)"):
                    body_struct = "\n".join(
                        f"{i+5}. {name} — (대지 제목) / 본문 주해 / 예화 / 적용"
                        for i, name in enumerate(shape["points"]))
                    prompt = f"""당신은 한국 장로교 강단에서 30년간 설교해 온 목회자입니다.
아래 조건으로 실제 강단에서 그대로 선포할 수 있는 설교 원고 전문을 작성하십시오.

[조건]
- 성경 본문: {full_scrip}
- 설교 주제: {topic}
- 신학적 관점: {theology}
- 설교 형태: {style}
- 대지 구조: {outline_key}
- 분량: 한국어 {length['chars']:,}자 내외 ({length_key} 선포 분량) — 이 분량을 반드시 지키십시오.

[신학적 관점 지침 — 이 설교 전체를 지배해야 합니다]
{lens['guide']}

[대지 구조 지침]
{shape['guide']}

[분량 지침]
{length['guide']}

[반드시 지킬 구조]
1. 제목
2. 본문 봉독 안내 ({full_scrip})
3. 서론 — 청중의 삶에서 출발하는 구체적 도입 (실제 있을 법한 상황 묘사)
4. 본문의 역사적·문학적 배경 설명
{body_struct}
{len(shape['points'])+5}. 결론 — 메시지 요약과 결단 촉구
{len(shape['points'])+6}. 마침 기도

[작성 규칙]
- 각 대지는 반드시 {full_scrip} 본문의 특정 구절을 인용하고 주해할 것.
- 예화는 구체적이고 현실적으로. "어떤 성도가 있었습니다" 식의 막연한 예화 금지.
- '은혜가 충만하기를' 같은 상투적 문구로 분량을 채우지 말 것.
- 100% 한국어. 영어 메모나 머리말 없이 설교 원고 본문만 출력.
- 마크다운 기호(#, **) 없이 평문으로 작성."""
                    res = get_ai_response(prompt, is_json=False, temperature=0.7,
                                          kind="sermon_write", max_tokens=length["tokens"])
                    st.session_state.temp_generated_sermon = res
                    st.session_state.temp_ai_title = f"{sel_book} 강해: {topic[:30]}"
                    st.session_state.temp_ai_scrip = full_scrip
                    st.rerun()

        if st.session_state.get("temp_generated_sermon"):
            st.write("---")
            show_ai_status()
            st.caption(f"생성 분량: {len(st.session_state.temp_generated_sermon):,}자")
            render_section_top_toolbar(st.session_state.temp_ai_title,
                                       st.session_state.temp_generated_sermon, "ai_gen_sermon")
            edited = st.text_area("작성된 강해설교문 검토 및 수정",
                                  value=st.session_state.temp_generated_sermon,
                                  height=450, key="edit_ai_sermon_area")

            if st.button("✅ 서재와 대시보드에 최종 등록", type="primary", key="btn_save_ai_sermon"):
                testament, book = classify_scripture(st.session_state.temp_ai_scrip)
                entry = {"title": st.session_state.temp_ai_title,
                         "scripture": st.session_state.temp_ai_scrip,
                         "testament": testament, "book": book, "topic": sel_book,
                         "theology": theology.split(' (')[0],
                         "date": datetime.now().strftime("%Y-%m-%d"),
                         "tags": [sel_book, theology.split(' (')[0], outline_key, length_key, "강해설교"],
                         "summary": "", "text": edited}
                saved = add_sermon_to_db(entry)
                load_sermon_to_workspace(saved, idx=len(st.session_state.sermon_library) - 1)
                st.success("등록 완료! [📊 설교 대시보드]로 이동하세요.")
                st.rerun()


# ==============================================================================
# 3) AI 보이스오버 스튜디오
# ==============================================================================
elif app_mode == "🎙️ AI 보이스오버 스튜디오":
    st.markdown("<h1 style='font-size:28px;font-weight:800;'>🎙️ AI 보이스오버 스튜디오</h1>",
                unsafe_allow_html=True)

    if not HAS_TTS:
        st.error("edge-tts 패키지가 설치되어 있지 않습니다. `pip install edge-tts` 후 다시 실행해 주세요.")

    c1, c2 = st.columns([1.5, 1])
    with c1:
        src = st.radio("변환할 원문", ["설교 원고 전문", "설교 요약", "직접 입력"], horizontal=True, key="vo_src")
        if src == "설교 원고 전문":
            default_txt = st.session_state.full_sermon
        elif src == "설교 요약":
            default_txt = st.session_state.get("sermon_summary_text", "")
        else:
            default_txt = ""

        vo_text = st.text_area("음성 변환 텍스트", value=default_txt, height=330, key=f"vo_text_{src}")
        vo_voice = st.selectbox("성우 보이스", ["인준 (남성 - 차분하고 신뢰감 있는 톤)", "선희 (여성 - 맑고 또렷한 톤)"],
                                key="vo_voice_sel")
        voice_id = "ko-KR-InJoonNeural" if "인준" in vo_voice else "ko-KR-SunHiNeural"

        st.caption(f"글자 수 {len(vo_text):,}자 · 예상 길이 약 {max(1, len(vo_text)//330)}분")

        if st.button("🎙️ 고음질 음성(TTS) 생성", type="primary", key="btn_gen_tts", disabled=not HAS_TTS):
            if not vo_text.strip():
                st.warning("텍스트를 입력해주세요.")
            else:
                with st.spinner("AI 음성 생성 중..."):
                    try:
                        st.session_state.vo_audio_path = generate_voiceover_audio(vo_text, voice_id)
                        st.success("보이스오버 생성 완료!")
                    except Exception as e:
                        st.error(f"음성 생성 실패: {e}")

    with c2:
        st.markdown("### 🎧 플레이어 & 다운로드")
        p = st.session_state.get("vo_audio_path")
        if p and os.path.exists(p):
            st.audio(p)
            with open(p, "rb") as af:
                st.download_button("📥 MP3 다운로드", data=af.read(),
                                   file_name=f"{st.session_state.sermon_title}_voice.mp3",
                                   mime="audio/mp3", key="dl_vo_mp3")
        else:
            st.info("왼쪽에서 생성하면 이곳에 플레이어가 나타납니다.")


# ==============================================================================
# 4) 쇼츠 만들기
# ==============================================================================
elif app_mode == "🎬 쇼츠 만들기 (스튜디오)":
    st.markdown("<h1 style='font-size:28px;font-weight:800;'>▶️ 쇼츠 만들기 스튜디오</h1>",
                unsafe_allow_html=True)

    if not HAS_VIDEO_ENGINE:
        st.warning("`video_engine.py` 모듈이 없어 영상 렌더링 기능이 비활성화됩니다. 같은 폴더에 파일을 두고 다시 실행해 주세요.")

    tab_yt, tab_ai = st.tabs(["🔗 유튜브 링크에서 숏츠 추출", "🎨 AI 나레이션 & 템플릿 숏츠 제작"])

    with tab_yt:
        st.caption("클라우드 서버 IP는 유튜브가 자주 차단합니다. 실패하면 오른쪽 탭에서 영상을 직접 업로드하세요.")
        yt_url = st.text_input("유튜브 영상 링크", placeholder="https://www.youtube.com/watch?v=...", key="yt_url_input")
        y1, y2, y3 = st.columns(3)
        with y1:
            s_min = st.number_input("시작 분", 0, 300, 12, key="yt_s_min")
        with y2:
            s_sec = st.number_input("시작 초", 0, 59, 30, key="yt_s_sec")
        with y3:
            dur = st.slider("길이(초)", 15, 60, 45, key="yt_dur")

        yt_title = st.text_input("상단 헤드라인", value=st.session_state.sermon_title, key="yt_title")
        yt_sub = st.text_input("강조 자막", value="", key="yt_sub")
        yt_church = st.text_input("교회명 워터마크", value=st.session_state.cn_church_name, key="yt_church")

        if st.button("🚀 9:16 세로 숏츠 추출", type="primary", key="btn_yt_extract"):
            if not yt_url.strip():
                st.warning("유튜브 링크를 입력해주세요.")
            else:
                with st.spinner("영상 추출 중... (20~40초)"):
                    try:
                        st.session_state.yt_extracted_result = extract_youtube_to_shorts(
                            yt_url.strip(), s_min * 60 + s_sec, dur, yt_title, yt_sub, yt_church)
                        st.success("추출 완료!")
                    except Exception as e:
                        st.warning(str(e))

        r = st.session_state.get("yt_extracted_result")
        if r and os.path.exists(r):
            st.write("---")
            v1, v2 = st.columns(2)
            with v1:
                st.video(r)
            with v2:
                with open(r, "rb") as f:
                    st.download_button("📥 MP4 다운로드", data=f.read(),
                                       file_name=f"{yt_title}_shorts.mp4", mime="video/mp4", key="dl_yt")

    with tab_ai:
        s1, s2, s3 = st.columns(3)
        with s1:
            st.link_button("🎥 픽사베이 영상", "https://pixabay.com/ko/videos/")
        with s2:
            st.link_button("📸 펙셀스 비디오", "https://www.pexels.com/ko-kr/videos/")
        with s3:
            st.link_button("🎵 무료 BGM", "https://pixabay.com/ko/music/")

        st.write("---")
        if st.button("✨ 쇼츠 제목 5개 & 해시태그 추천", key="btn_gen_shorts_meta"):
            with st.spinner("제목 분석 중..."):
                task = """[출력 형식 - JSON만]
{"titles": ["1. ...", "2. ...", "3. ...", "4. ...", "5. ..."], "hashtags": ["#...", ...8개]}
제목은 이 설교 원고의 실제 문장·키워드에서 뽑되, 25자 이내로 클릭을 부르게 다듬으십시오."""
                st.session_state.shorts_rec = get_ai_response(
                    build_grounded_prompt(task, ctx_chars=4000), is_json=True, kind="shorts_meta")
                st.rerun()

        selected_title = st.session_state.sermon_title
        if st.session_state.get("shorts_rec"):
            show_ai_status()
            rec = st.session_state.shorts_rec
            tc = st.radio("추천 제목", rec.get("titles", []), key="rad_shorts_title")
            if tc:
                selected_title = re.sub(r"^\d+\.\s*", "", tc)
            st.markdown("**추천 태그:** " + " ".join(rec.get("hashtags", [])))

        st.write("---")
        col_a, col_b = st.columns([1.2, 1])
        with col_a:
            v_title = st.text_input("쇼츠 제목", value=selected_title, key="in_shorts_title")
            _a = analyze_manuscript(st.session_state.full_sermon)
            default_script = "\n".join(_a["key_sentences"][:3]) or "말씀을 붙들고 이번 한 주를 살아갑시다."
            v_script = st.text_area("자막 대본 (줄바꿈 구분)", value=default_script, height=130, key="in_shorts_script")

            r1, r2 = st.columns(2)
            with r1:
                v_ratio = st.radio("비율", ["9:16 (세로)", "16:9 (가로)"], key="rad_ratio")
            with r2:
                v_voice = st.selectbox("보이스", ["인준 (남성)", "선희 (여성)"], key="sel_shorts_voice")

            with st.expander("🎨 폰트 크기 · 자막 위치", expanded=False):
                f1, f2 = st.columns(2)
                with f1:
                    t_fsize = st.slider("제목 폰트(pt)", 32, 72, 48, 2, key="sh_t_fsize")
                    t_ypos = st.slider("제목 Y 위치", 80, 500, 180, 10, key="sh_t_ypos")
                with f2:
                    s_fsize = st.slider("자막 폰트(pt)", 24, 60, 42, 2, key="sh_s_fsize")
                    s_ypos = st.slider("자막 Y 위치", 800, 1700, 1400, 20, key="sh_s_ypos")

            bg_media = st.file_uploader("배경 동영상/사진", type=["mp4", "mov", "jpg", "png"], key="up_shorts_bg")
            bgm_media = st.file_uploader("배경음악", type=["mp3", "wav"], key="up_shorts_bgm")

            if st.button("🚀 비디오 렌더링 시작", type="primary", key="btn_render_video",
                         disabled=not HAS_VIDEO_ENGINE):
                with st.spinner("자막 애니메이션 · BGM 믹싱 중..."):
                    try:
                        os.makedirs("./uploads", exist_ok=True)
                        bg_p = bgm_p = None
                        if bg_media:
                            bg_p = f"./uploads/{bg_media.name}"
                            with open(bg_p, "wb") as f:
                                f.write(bg_media.getbuffer())
                        if bgm_media:
                            bgm_p = f"./uploads/{bgm_media.name}"
                            with open(bgm_p, "wb") as f:
                                f.write(bgm_media.getbuffer())

                        lines = [l.strip() for l in v_script.split("\n") if l.strip()]
                        st.session_state.rendered_shorts_out = create_animated_video(
                            title=v_title, script_paragraphs=lines, bg_media_path=bg_p, bgm_path=bgm_p,
                            aspect_ratio=("9:16" if "9:16" in v_ratio else "16:9"),
                            voice=("ko-KR-InJoonNeural" if "인준" in v_voice else "ko-KR-SunHiNeural"),
                            title_fontsize=t_fsize, sub_fontsize=s_fsize, title_y=t_ypos, sub_y=s_ypos)
                        st.success("렌더링 완료!")
                    except Exception as e:
                        st.error(f"렌더링 실패: {e}")

        with col_b:
            st.markdown("### 🎬 완성된 영상")
            out = st.session_state.get("rendered_shorts_out")
            if out and os.path.exists(out):
                st.video(out)
                with open(out, "rb") as vf:
                    st.download_button("📥 MP4 다운로드", data=vf.read(), file_name="sermon_shorts.mp4",
                                       mime="video/mp4", key="dl_shorts_mp4")
            else:
                st.info("렌더링하면 이곳에 영상이 나타납니다.")


# ==============================================================================
# 5) 말씀카드 이미지
# ==============================================================================
elif app_mode == "📷 말씀카드 이미지":
    st.markdown("<h1 style='font-size:28px;font-weight:800;'>📷 말씀카드 이미지 스튜디오</h1>",
                unsafe_allow_html=True)

    _a = analyze_manuscript(st.session_state.full_sermon)
    candidates = _a["key_sentences"][:8] or [st.session_state.sermon_title]

    c1, c2 = st.columns([1.3, 1.2])
    with c1:
        st.markdown("#### ✏️ 문구 선택 (원고에서 자동 추출)")
        pick = st.selectbox("원고 속 핵심 문장에서 고르기",
                            ["(직접 입력)"] + [s[:70] + ("…" if len(s) > 70 else "") for s in candidates],
                            key="vc_pick")
        if pick != "(직접 입력)":
            idx = ["(직접 입력)"] + candidates
            init_text = candidates[[s[:70] + ("…" if len(s) > 70 else "") for s in candidates].index(pick)]
        else:
            init_text = st.session_state.sermon_title

        v_text = st.text_area("말씀문구 / 메시지", value=init_text, height=120, key=f"vc_text_{pick}")
        v_scrip = st.text_input("성경 구절", value=st.session_state.sermon_scripture, key="vc_scrip_in")
        v_church = st.text_input("교회명 배지", value=st.session_state.cn_church_name, key="vc_church_in")

        st.markdown("#### 🎨 배경 & 스타일 (무한 이미지)")
        b1, b2 = st.columns(2)
        with b1:
            bg_opt = st.radio("배경 방식", ["사진(무한)", "기본", "직접 업로드"], key="vc_bg_radio")
            st.selectbox("이미지 분위기", list(BG_THEMES.keys()), key="vc_theme")
        with b2:
            up_file = st.file_uploader("배경 이미지", type=["jpg", "png", "jpeg"], key="vc_up_file") \
                if bg_opt == "직접 업로드" else None
            # ⚠️ 위젯 key 를 코드에서 직접 수정하면 Streamlit 이 예외를 냅니다.
            #    그래서 카운터는 위젯이 아닌 일반 세션 키로 따로 관리합니다.
            bg_i = int(st.session_state.get("vc_bg_counter", 0))
            st.markdown(f"<div style='padding:6px 0;color:#7dd3fc;font-weight:700;'>"
                        f"현재 이미지 #{bg_i}</div>", unsafe_allow_html=True)
            if st.button("🔀 다른 이미지로 바꾸기", key="vc_shuffle"):
                st.session_state["vc_bg_counter"] = bg_i + 1
                st.rerun()
        bg_opt = "사진" if bg_opt.startswith("사진") else bg_opt

        s1, s2 = st.columns(2)
        with s1:
            fsize = st.slider("폰트 크기", 28, 68, 42, 2, key="vc_fsize")
            lspace = st.slider("줄간격", 10, 40, 18, 2, key="vc_lspace")
        with s2:
            fcolor = st.color_picker("글자 색", "#FDE047", key="vc_fcolor")
            scolor = st.color_picker("테두리 색", "#000000", key="vc_scolor")
            opacity = st.slider("배경 어둡기", 0.2, 0.9, 0.6, 0.05, key="vc_op")

    with c2:
        st.markdown("### 🖼️ 미리보기 & 다운로드")
        png = generate_verse_card_png(v_text, v_scrip, bg_opt, up_file, fsize, lspace,
                                      fcolor, scolor, opacity, v_church, bg_index=bg_i)
        st_image_full(png.getvalue(), caption="1:1 고화질 말씀카드")
        st.download_button("📥 PNG 다운로드", data=png.getvalue(),
                           file_name=f"{st.session_state.sermon_title}_말씀카드.png",
                           mime="image/png", key="dl_verse_card")


# ==============================================================================
# 6) 설교 서재
# ==============================================================================
elif app_mode == "📚 설교 서재 (Sermon Library)":
    st.markdown("<h1 style='font-size:28px;font-weight:800;'>📚 설교 서재 (영구 기록보관소)</h1>",
                unsafe_allow_html=True)

    sermons_db = get_db_sermons()
    st.session_state.sermon_library = sermons_db

    if cloud_store_ready():
        st.success("🗄️ **영구 보관 작동 중** — GitHub 비공개 보관함에 자동 저장됩니다. "
                   "앱을 다시 켜도 설교가 사라지지 않습니다.")
    else:
        st.error("⚠️ **지금은 임시 저장 상태입니다.** 앱이 재시작되면 아래 설교들이 사라집니다. "
                 "왼쪽 사이드바 **[🗄️ 설교 서재 저장소]** 를 열어 영구 보관을 켜 주세요. "
                 "그 전까지는 아래 [💾 전체 백업] 을 꼭 눌러 파일로 보관하세요.")

    t1, t2 = st.columns([1.5, 1.5])
    with t1:
        st.markdown(f"총 **{len(sermons_db):,}편**의 설교문이 보관되어 있습니다.")
    with t2:
        k1, k2 = st.columns(2)
        with k1:
            st.download_button("💾 전체 백업(.json)",
                               data=json.dumps(sermons_db, ensure_ascii=False, indent=2).encode('utf-8'),
                               file_name=f"설교서재_백업_{datetime.now().strftime('%Y%m%d')}.json",
                               mime="application/json", key="dl_backup")
        with k2:
            with st.popover("📂 백업 복원"):
                rf = st.file_uploader("백업 JSON", type=["json"], key="up_restore")
                mode = st.radio("복원 방식", ["기존에 병합(추가)", "완전 교체"], key="restore_mode")
                if rf and st.button("✅ 복원 실행", key="btn_restore"):
                    try:
                        data = json.load(rf)
                        if isinstance(data, list):
                            if mode == "기존에 병합(추가)":
                                exist = {(s.get("title"), s.get("scripture")) for s in sermons_db}
                                merged = list(sermons_db)
                                nid = max([int(s.get("id", 0)) for s in sermons_db] or [0])
                                for s in data:
                                    if (s.get("title"), s.get("scripture")) not in exist:
                                        nid += 1
                                        s["id"] = nid
                                        merged.append(s)
                                data = merged
                            save_db_sermons(data)
                            st.session_state.sermon_library = data
                            st.success(f"복원 완료! 현재 {len(data)}편")
                            st.rerun()
                        else:
                            st.error("올바른 백업 파일이 아닙니다.")
                    except Exception as e:
                        st.error(f"복원 실패: {e}")

    st.write("---")
    f1, f2, f3, f4 = st.columns([1.5, 1.2, 1.5, 1.2])
    with f1:
        kw = st.text_input("검색어 (제목/본문/키워드)", key="lib_kw")
    with f2:
        tf = st.selectbox("구약/신약", ["전체", "구약", "신약"], key="lib_testament")
    with f3:
        bf = st.selectbox("성경 66권", ["전체"] + BIBLE_BOOKS, key="lib_book")
    with f4:
        so = st.selectbox("정렬", ["최신순", "오래된순"], key="lib_sort")

    st.write("---")
    filtered = []
    for s in sermons_db:
        if "testament" not in s or "book" not in s:
            s["testament"], s["book"] = classify_scripture(s.get("scripture", ""))
        if kw:
            hay = " ".join([s.get("title", ""), s.get("scripture", ""), s.get("text", "")]
                           + list(s.get("tags", []))).lower()
            if kw.lower() not in hay:
                continue
        if tf != "전체" and s.get("testament") != tf:
            continue
        if bf != "전체" and s.get("book") != bf:
            continue
        filtered.append(s)

    filtered.sort(key=lambda x: x.get("id", 0), reverse=(so == "최신순"))
    st.markdown(f"**검색 결과:** `{len(filtered):,}편`")

    if not filtered:
        st.info("조건에 맞는 설교문이 없습니다.")
    else:
        for i, s in enumerate(filtered):
            tags_html = ' '.join(
                f'<span style="background:#1e293b;color:#38bdf8;padding:3px 10px;border-radius:6px;'
                f'font-size:12px;margin-right:4px;">#{t}</span>' for t in s.get('tags', []))
            st.markdown(
                f"""<div class="lib-card">
                <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:8px;">
                <h3 style="margin:0;font-size:20px;font-weight:bold;color:#f8fafc;">{s.get('title')}</h3>
                <span style="background:#1e3a8a;color:#fde047;padding:4px 12px;border-radius:12px;font-size:12px;font-weight:bold;">
                {s.get('testament','기타')} · {s.get('book','성경')}</span></div>
                <p style="margin:0 0 10px 0;font-size:14px;color:#94a3b8;">
                📖 <strong>{s.get('scripture')}</strong> · [{s.get('theology','-')}] · 📅 {s.get('date','-')}
                · 📝 {len(s.get('text','')):,}자</p>
                <div style="margin-bottom:6px;">{tags_html}</div></div>""",
                unsafe_allow_html=True
            )
            c1, c2 = st.columns([3, 1])
            with c1:
                if st.button("📖 대시보드로 불러오기", key=f"lib_load_{s.get('id')}_{i}"):
                    try:
                        real_idx = sermons_db.index(s)
                    except ValueError:
                        real_idx = 0
                    load_sermon_to_workspace(s, idx=real_idx)
                    st.success(f"'{s.get('title')}' 로 전체 메뉴가 동기화되었습니다.")
                    st.rerun()
            with c2:
                if st.button("🗑️ 삭제", key=f"lib_del_{s.get('id')}_{i}"):
                    updated = [x for x in get_db_sermons() if x.get('id') != s.get('id')]
                    save_db_sermons(updated)
                    st.success("삭제되었습니다.")
                    st.rerun()
